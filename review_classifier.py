import logging
import os
import json
import numpy as np
import pandas as pd
from functools import wraps
from typing import Dict, List, Any, Optional, Tuple, Union
from fastapi import FastAPI, APIRouter, UploadFile, File, HTTPException, Form, Request, status
from fastapi.responses import JSONResponse, HTMLResponse, Response
from fastapi.middleware.cors import CORSMiddleware
import traceback
import time
import asyncio
from io import StringIO
import tempfile
from openai import AzureOpenAI
from fuzzywuzzy import process
import aiohttp
import tempfile
from typing import Optional, Dict, Any, Union
import mimetypes
from io import BytesIO
from pydantic import BaseModel, ValidationError
from typing import List, Optional
import re
# Core imports
import pandas as pd
import chardet
import pdfplumber
# Try importing Unstructured components
try:
    from unstructured.partition.auto import partition
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False
    partition = None

# Import specific partitioners for fallback
try:
    from unstructured.partition.csv import partition_csv
except ImportError:
    partition_csv = None

try:
    from unstructured.partition.docx import partition_docx
except ImportError:
    partition_docx = None

try:
    from unstructured.partition.pdf import partition_pdf
except ImportError:
    partition_pdf = None

try:
    from unstructured.partition.html import partition_html
except ImportError:
    partition_html = None

try:
    from unstructured.partition.text import partition_text
except ImportError:
    partition_text = None

try:
    from unstructured.partition.email import partition_email
except ImportError:
    partition_email = None

try:
    from unstructured.partition.xlsx import partition_xlsx
except ImportError:
    partition_xlsx = None

try:
    from unstructured.partition.image import partition_image
except ImportError:
    partition_image = None

try:
    from unstructured.partition.json import partition_json
except ImportError:
    partition_json = None

try:
    from unstructured.partition.xml import partition_xml
except ImportError:
    partition_xml = None

try:
    from unstructured.partition.md import partition_md
except ImportError:
    partition_md = None

try:
    from unstructured.partition.pptx import partition_pptx
except ImportError:
    partition_pptx = None

try:
    from unstructured.partition.msg import partition_msg
except ImportError:
    partition_msg = None

try:
    from unstructured.partition.rtf import partition_rtf
except ImportError:
    partition_rtf = None

try:
    from unstructured.partition.epub import partition_epub
except ImportError:
    partition_epub = None

try:
    from unstructured.partition.odt import partition_odt
except ImportError:
    partition_odt = None

try:
    from unstructured.partition.doc import partition_doc
except ImportError:
    partition_doc = None

try:
    from unstructured.partition.ppt import partition_ppt
except ImportError:
    partition_ppt = None

try:
    from unstructured.partition.tsv import partition_tsv
except ImportError:
    partition_tsv = None

try:
    from unstructured.partition.rst import partition_rst
except ImportError:
    partition_rst = None

try:
    from unstructured.partition.org import partition_org
except ImportError:
    partition_org = None

# Fallback imports for when Unstructured isn't available
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import html2text
except ImportError:
    html2text = None

try:
    import markdown
except ImportError:
    markdown = None


    
# Add Pydantic models for validation

# Add these constants at the top of your file
EXTRACT_REVIEWS_API_URL = "https://copilotv2.azurewebsites.net/extract-reviews"
SUPPORTED_DIRECT_FORMATS = {'.csv', '.xlsx', '.xls', '.xlsm'}  # Formats we handle directly
API_TIMEOUT = 60  # Timeout for the extract-reviews API call
_positive_anchor_embedding = None
_negative_anchor_embedding = None
def async_timeout(seconds):
    def decorator(func):
        @wraps(func)
        async def wrapper(*args, **kwargs):
            try:
                return await asyncio.wait_for(func(*args, **kwargs), timeout=seconds)
            except asyncio.TimeoutError:
                logger.error(f"Function {func.__name__} timed out after {seconds} seconds")
                raise HTTPException(
                    status_code=status.HTTP_504_GATEWAY_TIMEOUT,
                    detail=f"Request timed out after {seconds} seconds"
                )
        return wrapper
    return decorator
# Configure detailed logging for development
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - [%(filename)s:%(lineno)d] - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("customer_feedback.log")
    ]
)
logger = logging.getLogger("customer_feedback")

# Azure OpenAI client configuration - reuse from main app
AZURE_ENDPOINT = "https://prodhubfinnew-openai-97de.openai.azure.com/"
AZURE_API_KEY = "97fa8c02f9e64e8ea5434987b11fe6f4" 
AZURE_API_VERSION = "2024-12-01-preview"  # Updated API version

# Embedding model configuration
EMBEDDING_MODEL = "text-embedding-3-small"  # Correct embedding model name
AZURE_CLIENT = AzureOpenAI(
    api_version=AZURE_API_VERSION,
    azure_endpoint=AZURE_ENDPOINT,
    api_key=AZURE_API_KEY
)


# Create standalone app (primary usage mode)
app = FastAPI(
    title="Customer Feedback Analysis API",
    description="API for analyzing customer feedback and identifying key problem areas",
    version="1.0.0"
)

# Add CORS middleware to allow cross-domain requests
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # For development; restrict in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
@app.middleware("http")
async def timeout_middleware(request: Request, call_next):
    try:
        return await asyncio.wait_for(call_next(request), timeout=120)
    except asyncio.TimeoutError:
        return JSONResponse(
            status_code=504,
            content={"detail": "Request timed out after 120 seconds"}
        )

class UnstructuredDocumentExtractor:
    """
    Production-grade universal document text extractor using Unstructured library.
    
    This extractor leverages the powerful Unstructured library for sophisticated
    document parsing, with comprehensive fallbacks for robustness.
    
    Supported formats (via Unstructured):
    - Documents: DOC, DOCX, ODT, PDF, RTF, TXT, LOG
    - Spreadsheets: XLS, XLSX, CSV, TSV
    - Presentations: PPT, PPTX
    - Web: HTML, HTM, XML
    - Email: EML, MSG
    - Images: PNG, JPG, JPEG, TIFF, BMP, HEIC (with OCR)
    - E-books: EPUB
    - Markup: MD, RST, ORG
    - Code: JS, PY, JAVA, CPP, CC, CXX, C, CS, PHP, RB, SWIFT, TS, GO
    - Data: JSON
    """
    
    def __init__(self, logger: Optional[logging.Logger] = None):
        """Initialize the extractor with optional logger."""
        self.logger = logger or logging.getLogger(__name__)
        
        # Log available components
        if UNSTRUCTURED_AVAILABLE:
            self.logger.info("Unstructured library is available")
        else:
            self.logger.warning("Unstructured library not available, using fallback methods")
            
    def extract_text(self, 
                    file_content: Union[bytes, str], 
                    filename: str,
                    encoding: Optional[str] = None,
                    strategy: str = "auto",
                    include_metadata: bool = False,
                    max_partition_length: int = 1500,
                    languages: Optional[List[str]] = None,
                    extract_tables: bool = True) -> str:
        """
        Extract text from any supported document type using Unstructured.
        
        Args:
            file_content: Raw file content as bytes or string
            filename: Original filename for type detection
            encoding: Optional encoding override
            strategy: Partitioning strategy ("auto", "fast", "hi_res", "ocr_only")
            include_metadata: Whether to include element metadata in output
            max_partition_length: Maximum length for text partitions
            languages: OCR languages (e.g., ["eng", "spa"])
            extract_tables: Whether to extract and format tables
            
        Returns:
            Extracted text as string
        """
        try:
            # Convert string to bytes if needed
            if isinstance(file_content, str):
                file_content = file_content.encode('utf-8')
            
            # Get file extension
            file_ext = os.path.splitext(filename)[1].lower()
            self.logger.info(f"Extracting text from {filename} (type: {file_ext})")
            
            # Try Unstructured first if available
            if UNSTRUCTURED_AVAILABLE:
                try:
                    text = self._extract_with_unstructured(
                        file_content, filename, encoding, strategy, 
                        include_metadata, max_partition_length, languages
                    )
                    if text and len(text.strip()) > 10:
                        return text
                except Exception as e:
                    self.logger.warning(f"Unstructured extraction failed: {str(e)}, trying fallbacks")
            
            # Use fallback methods
            return self._extract_with_fallback(file_content, filename, encoding)
            
        except Exception as e:
            self.logger.error(f"Error extracting text from {filename}: {str(e)}")
            self.logger.error(f"Traceback: {traceback.format_exc()}")
            # Emergency fallback
            return self._emergency_text_extraction(file_content)
    
    def _extract_with_unstructured(self, 
                                  file_content: bytes, 
                                  filename: str,
                                  encoding: Optional[str],
                                  strategy: str,
                                  include_metadata: bool,
                                  max_partition_length: int,
                                  languages: Optional[List[str]]) -> str:
        """Extract text using Unstructured library."""
        # Save to temporary file (Unstructured often works better with files)
        with tempfile.NamedTemporaryFile(
            delete=False, 
            suffix=os.path.splitext(filename)[1]
        ) as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        try:
            file_ext = os.path.splitext(filename)[1].lower()
            # Prepare kwargs
            if encoding is None and file_ext in ['.csv', '.txt', '.log', '.md']:
                encoding = self._detect_encoding(file_content)
                self.logger.info(f"Auto-detected encoding: {encoding}")
            
            kwargs = {
                "filename": tmp_path,
                "encoding": encoding or 'utf-8',  # Use detected or default
                "max_partition": max_partition_length,
            }
            
            if file_ext in ['.pdf', '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.heic']:
                kwargs["strategy"] = strategy
                if languages:
                    kwargs["languages"] = languages
            
            # Special handling for specific file types
            if file_ext in ['.html', '.htm']:
                kwargs["include_page_breaks"] = True
            elif file_ext in ['.pdf']:
                kwargs["include_page_breaks"] = True
                kwargs["infer_table_structure"] = True
            elif file_ext in ['.eml', '.msg']:
                kwargs["process_attachments"] = False  # Don't process attachments for text extraction
                kwargs["content_source"] = "text/html"  # Prefer HTML content
                
            # Use the main partition function
            elements = partition(**kwargs)
            
            # Convert elements to text
            text_parts = []
            
            for element in elements:
                # Get element text
                element_text = str(element)
                
                if include_metadata and hasattr(element, 'metadata'):
                    # Add metadata if requested
                    metadata = element.metadata
                    if hasattr(metadata, 'page_number') and metadata.page_number:
                        element_text = f"[Page {metadata.page_number}] {element_text}"
                    if hasattr(metadata, 'section') and metadata.section:
                        element_text = f"[{metadata.section}] {element_text}"
                
                # Handle tables specially
                if hasattr(element, 'metadata') and hasattr(element.metadata, 'text_as_html'):
                    # Convert HTML table to text representation
                    table_text = self._html_table_to_text(element.metadata.text_as_html)
                    if table_text:
                        text_parts.append(table_text)
                    else:
                        text_parts.append(element_text)
                else:
                    text_parts.append(element_text)
            
            return "\n\n".join(text_parts)
            
        finally:
            # Clean up temp file
            try:
                os.unlink(tmp_path)
            except:
                pass
    
    def _extract_with_fallback(self, file_content: bytes, filename: str, encoding: Optional[str]) -> str:
        """Fallback extraction methods when Unstructured is not available or fails."""
        file_ext = os.path.splitext(filename)[1].lower()
        
        # Try specific extractors based on file type
        if file_ext in ['.docx']:
            return self._extract_docx_fallback(file_content)
        elif file_ext in ['.pdf']:
            return self._extract_pdf_fallback(file_content)
        elif file_ext in ['.xlsx', '.xls']:
            return self._extract_excel_fallback(file_content, filename)
        elif file_ext in ['.csv']:
            return self._extract_csv_fallback(file_content, encoding)
        elif file_ext in ['.html', '.htm']:
            return self._extract_html_fallback(file_content, encoding)
        elif file_ext in ['.json']:
            return self._extract_json_fallback(file_content, encoding)
        elif file_ext in ['.xml']:
            return self._extract_xml_fallback(file_content, encoding)
        elif file_ext in ['.md', '.markdown']:
            return self._extract_markdown_fallback(file_content, encoding)
        elif file_ext in ['.pptx']:
            return self._extract_pptx_fallback(file_content)
        elif file_ext in ['.eml']:
            return self._extract_email_fallback(file_content, encoding)
        elif file_ext in ['.msg']:
            return self._extract_msg_fallback(file_content)
        else:
            # Default text extraction
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_docx_fallback(self, file_content: bytes) -> str:
        """Fallback DOCX extraction without Unstructured."""
        if DocxDocument:
            try:
                doc = DocxDocument(BytesIO(file_content))
                paragraphs = []
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                
                # Extract tables
                for table in doc.tables:
                    table_text = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))
                    if table_text:
                        paragraphs.append("\n".join(table_text))
                
                return "\n\n".join(paragraphs)
            except Exception as e:
                self.logger.error(f"DOCX fallback failed: {str(e)}")
        
        # Try XML extraction
        return self._extract_docx_xml_fallback(file_content)
    
    def _extract_docx_xml_fallback(self, file_content: bytes) -> str:
        """Extract text from DOCX by parsing XML."""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            text_parts = []
            
            with zipfile.ZipFile(BytesIO(file_content)) as docx:
                # Extract main document
                if 'word/document.xml' in docx.namelist():
                    xml_content = docx.read('word/document.xml')
                    tree = ET.fromstring(xml_content)
                    
                    namespaces = {
                        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                    }
                    
                    # Extract paragraphs
                    for para in tree.findall('.//w:p', namespaces):
                        para_text = []
                        for t in para.findall('.//w:t', namespaces):
                            if t.text:
                                para_text.append(t.text)
                        if para_text:
                            text_parts.append(''.join(para_text))
            
            return "\n".join(text_parts)
        except:
            return self._extract_text_with_encoding(file_content)
    
    def _extract_pdf_fallback(self, file_content: bytes) -> str:
        """Fallback PDF extraction without Unstructured."""
        try:
            text_parts = []   
            with pdfplumber.open(BytesIO(file_content)) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[Page {i+1}]\n{page_text}")
                        
                        # Extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            table_text = self._format_table(table)
                            text_parts.append(table_text)
                
            return "\n\n".join(text_parts)
        except:
            pass
        
        return self._extract_text_with_encoding(file_content)
    
    def _extract_excel_fallback(self, file_content: bytes, filename: str) -> str:
        """Fallback Excel extraction without Unstructured."""
        try:
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            
            try:
                excel_file = pd.ExcelFile(tmp_path)
                text_parts = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(tmp_path, sheet_name=sheet_name)
                    text_parts.append(f"\n[Sheet: {sheet_name}]\n")
                    
                    if not df.empty:
                        # Convert to readable format
                        text_parts.append(df.to_string())
                
                return "\n".join(text_parts)
            finally:
                os.unlink(tmp_path)
        except:
            return self._extract_csv_fallback(file_content, None)
    
    def _extract_csv_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback CSV extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            
            try:
                df = pd.read_csv(StringIO(text_content))
                return df.to_string()
            except:
                return text_content
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_html_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback HTML extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            
            if html2text:
                h = html2text.HTML2Text()
                h.ignore_links = False
                h.body_width = 0  # Don't wrap lines
                return h.handle(text_content)
            else:
                # Basic HTML stripping
                text = re.sub(r'<script[^>]*>.*?</script>', '', text_content, flags=re.DOTALL)
                text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text)
                return text.strip()
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_json_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback JSON extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            data = json.loads(text_content)
            
            # Pretty print JSON
            return json.dumps(data, indent=2, ensure_ascii=False)
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_xml_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback XML extraction without Unstructured."""
        try:
            import xml.etree.ElementTree as ET
            
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            root = ET.fromstring(text_content)
            
            # Extract all text from XML
            texts = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())
            
            return "\n".join(texts)
        except:
            # Just strip tags
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            text = file_content.decode(encoding, errors='replace')
            text = re.sub(r'<[^>]+>', ' ', text)
            return re.sub(r'\s+', ' ', text).strip()
    
    def _extract_markdown_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback Markdown extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            
            if markdown and html2text:
                # Convert to HTML then to plain text
                html_content = markdown.markdown(text_content)
                h = html2text.HTML2Text()
                h.body_width = 0
                return h.handle(html_content)
            
            return text_content
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_pptx_fallback(self, file_content: bytes) -> str:
        """Fallback PPTX extraction without Unstructured."""
        if Presentation:
            try:
                prs = Presentation(BytesIO(file_content))
                text_parts = []
                
                for i, slide in enumerate(prs.slides):
                    text_parts.append(f"\n[Slide {i + 1}]\n")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_parts.append(shape.text)
                        
                        if shape.has_table:
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    table_text.append(" | ".join(row_text))
                            if table_text:
                                text_parts.append("\n".join(table_text))
                
                return "\n".join(text_parts)
            except:
                pass
        
        return self._extract_text_with_encoding(file_content)
    
    def _extract_email_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback email extraction without Unstructured."""
        try:
            import email
            from email import policy
            
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            # Parse email
            msg = email.message_from_bytes(file_content, policy=policy.default)
            
            text_parts = []
            
            # Add headers
            headers = ['From', 'To', 'Subject', 'Date']
            for header in headers:
                value = msg.get(header)
                if value:
                    text_parts.append(f"{header}: {value}")
            
            text_parts.append("")  # Empty line
            
            # Extract body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            text_parts.append(payload.decode(encoding, errors='replace'))
                    elif part.get_content_type() == "text/html" and not any("text/plain" in p.get_content_type() for p in msg.walk()):
                        payload = part.get_payload(decode=True)
                        if payload:
                            html_text = payload.decode(encoding, errors='replace')
                            # Convert HTML to text
                            if html2text:
                                h = html2text.HTML2Text()
                                text_parts.append(h.handle(html_text))
                            else:
                                # Strip HTML tags
                                text = re.sub(r'<[^>]+>', ' ', html_text)
                                text_parts.append(text)
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    text_parts.append(payload.decode(encoding, errors='replace'))
            
            return "\n".join(text_parts)
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_msg_fallback(self, file_content: bytes) -> str:
        """Fallback MSG extraction without Unstructured."""
        # MSG files are complex binary format, just extract readable text
        return self._extract_text_with_encoding(file_content)
    
    def _html_table_to_text(self, html_table: str) -> str:
        """Convert HTML table to readable text format."""
        try:
            # Simple HTML table parsing
            rows = re.findall(r'<tr[^>]*>(.*?)</tr>', html_table, re.DOTALL)
            table_text = []
            
            for row in rows:
                cells = re.findall(r'<t[hd][^>]*>(.*?)</t[hd]>', row, re.DOTALL)
                if cells:
                    # Clean cell content
                    clean_cells = []
                    for cell in cells:
                        cell_text = re.sub(r'<[^>]+>', '', cell)
                        cell_text = cell_text.strip()
                        clean_cells.append(cell_text)
                    
                    table_text.append(" | ".join(clean_cells))
            
            return "\n".join(table_text)
        except:
            return ""
    
    def _format_table(self, table_data: List[List[Any]]) -> str:
        """Format table data as readable text."""
        if not table_data:
            return ""
        
        formatted_rows = []
        for row in table_data:
            if row:
                formatted_row = " | ".join(str(cell) if cell is not None else "" for cell in row)
                formatted_rows.append(formatted_row)
        
        return "\n".join(formatted_rows)
    
    def _extract_text_with_encoding(self, file_content: bytes, encoding: Optional[str] = None) -> str:
        """Extract text with automatic encoding detection."""
        if encoding is None:
            encoding = self._detect_encoding(file_content)
        
        try:
            return file_content.decode(encoding, errors='replace')
        except:
            # Try common encodings
            for enc in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1', 'utf-16']:
                try:
                    return file_content.decode(enc, errors='replace')
                except:
                    continue
            
            # Force UTF-8
            return file_content.decode('utf-8', errors='replace')
    
    def _detect_encoding(self, file_content: bytes) -> str:
        """Detect file encoding using chardet."""
        try:
            # Sample first 100KB for performance
            sample = file_content[:100000] if len(file_content) > 100000 else file_content
            result = chardet.detect(sample)
            
            encoding = result.get('encoding')
            confidence = result.get('confidence', 0)
            
            if encoding and confidence > 0.7:
                self.logger.info(f"Detected encoding: {encoding} (confidence: {confidence:.2f})")
                return encoding
        except:
            pass
        
        return 'utf-8'
    
    def _emergency_text_extraction(self, file_content: bytes) -> str:
        """Emergency fallback to extract any readable text."""
        try:
            # First, check if it's a known binary format that we shouldn't try to decode
            file_header = file_content[:10] if len(file_content) >= 10 else file_content
            
            # Check for PDF header
            if file_header.startswith(b'%PDF'):
                self.logger.error("Emergency extraction called on PDF file - cannot extract without proper tools")
                return "PDF file detected but text extraction failed. Please ensure PDF extraction libraries are installed."
                
            # Check for other binary formats
            binary_headers = [
                b'\x50\x4b\x03\x04',  # ZIP/DOCX/XLSX
                b'\xd0\xcf\x11\xe0',  # DOC/XLS
                b'\x89\x50\x4e\x47',  # PNG
                b'\xff\xd8\xff',      # JPEG
            ]
            
            for header in binary_headers:
                if file_header.startswith(header):
                    self.logger.error("Emergency extraction called on binary file - cannot extract")
                    return "Binary file detected but text extraction failed. Please ensure document processing libraries are installed."
            
            # Try to extract printable ASCII and common Unicode
            text_parts = []
            i = 0
            
            while i < len(file_content):
                # Try to decode as UTF-8
                for length in range(4, 0, -1):
                    if i + length <= len(file_content):
                        try:
                            char = file_content[i:i+length].decode('utf-8')
                            if char.isprintable() or char in '\n\r\t':
                                text_parts.append(char)
                                i += length
                                break
                        except:
                            pass
                else:
                    # Skip non-decodable byte
                    i += 1
            
            text = ''.join(text_parts)
            
            # Clean up
            text = re.sub(r'\s+', ' ', text)
            text = re.sub(r' +', ' ', text)
            
            result = text.strip()
            
            # If we got very little text from a large file, it's probably binary
            if len(result) < 100 and len(file_content) > 1000:
                return "Unable to extract meaningful text from this document. The file may be corrupted or in an unsupported format."
                
            return result if result else "Unable to extract text from this document."
            
        except Exception as e:
            self.logger.error(f"Emergency extraction failed: {str(e)}")
            return "Unable to extract text from this document."
async def extract_text_internal(
    file_content: bytes,
    filename: str,
    strategy: str = "auto",
    languages: Optional[List[str]] = None,
    encoding: Optional[str] = None,
    logger: Optional[logging.Logger] = None
) -> str:
    max_chars = 50000
    """
    Internal function to extract text from document content.
    This can be called from other API endpoints.
    
    Args:
        file_content: Raw file content as bytes
        filename: Original filename for type detection
        strategy: Extraction strategy
        languages: OCR languages list
        encoding: Optional encoding override
        logger: Logger instance
        
    Returns:
        Extracted text as string
        
    Raises:
        Exception: If extraction fails
    """
    if logger is None:
        logger = logging.getLogger(__name__)
    
    try:
        extractor = UnstructuredDocumentExtractor(logger=logger)
        extracted_text = extractor.extract_text(
            file_content=file_content,
            filename=filename,
            encoding=encoding,
            strategy=strategy,
            languages=languages
        )
        if max_chars and len(extracted_text) > max_chars:
            logger.info(f"Truncating extracted text from {len(extracted_text)} to {max_chars} characters")
            # Try to truncate at a sentence boundary
            truncated = extracted_text[:max_chars]
            last_period = truncated.rfind('.')
            last_newline = truncated.rfind('\n')
            cutoff = max(last_period, last_newline)
            if cutoff > max_chars * 0.8:  # Only use if it's not too far back
                extracted_text = truncated[:cutoff + 1]
            else:
                extracted_text = truncated + "..."
                
        return extracted_text
    except Exception as e:
        logger.error(f"Internal text extraction failed: {str(e)}")
        raise


class FileConversionError(Exception):
    '''Custom exception for file conversion errors'''
    pass
async def extract_raw_content_from_file(file_content: bytes, filename: str) -> str:
    '''
    Extracts raw text content from file bytes for fallback processing.
    Handles multiple file types and encodings robustly.
    Uses the same convert_file_to_csv logic as analyze_feedback endpoint.
    
    Args:
        file_content: Raw file bytes
        filename: Original filename for type detection
        
    Returns:
        Raw text content as string
    '''
    try:
        file_ext = os.path.splitext(filename)[1].lower()
        logger.info(f"Extracting raw content from {filename} (type: {file_ext})")
        
        # For Excel files, try to extract with pandas
        if file_ext in ['.xlsx', '.xls', '.xlsm']:
            try:
                import io
                excel_file = pd.ExcelFile(io.BytesIO(file_content))
                all_text = []
                for sheet in excel_file.sheet_names:
                    df = pd.read_excel(io.BytesIO(file_content), sheet_name=sheet)
                    all_text.append(f"Sheet: {sheet}\n")
                    all_text.append(df.to_string())
                return '\n'.join(all_text)
            except Exception as excel_error:
                logger.warning(f"Failed to parse Excel file {filename}: {str(excel_error)}, trying conversion API")
        
        # For CSV, ensure proper parsing
        if file_ext == '.csv':
            try:
                # Try multiple encodings for CSV
                for encoding in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                    try:
                        text = file_content.decode(encoding)
                        # Validate it's actually CSV-like
                        if ',' in text or '\t' in text or '|' in text:
                            return text
                    except:
                        continue
            except:
                logger.warning(f"Failed to parse CSV file {filename}, trying conversion API")
        
        # For PDF, DOCX, TXT and other files, use convert_file_to_csv (same as analyze_feedback)
        if file_ext in ['.pdf', '.docx', '.doc', '.txt', '.text', '.rtf', '.odt'] or file_ext not in SUPPORTED_DIRECT_FORMATS:
            logger.info(f"Using convert_file_to_csv API for {filename}")
            
            # Build extraction prompt for getting ALL text (not just feedback)
            extraction_prompt = '''
Extract ALL text content from this document into CSV format.

CRITICAL INSTRUCTIONS:
1. Extract EVERY piece of text in the document - not just feedback
2. Each paragraph or logical section should be a separate row
3. Include ALL content: headers, footers, body text, tables, lists, etc.
4. Extract text EXACTLY as written - do not modify or summarize
5. Preserve the original text structure as much as possible

OUTPUT COLUMNS (use EXACTLY these names):
- text_content: The complete text content (can be multiline)
- section: Section identifier or "Main" if not clear
- page_number: Page number if available, else "N/A"

IMPORTANT: This is for complete text extraction. Extract ALL text, not just customer feedback.'''
            
            # Try conversion with retries (same pattern as analyze_feedback)
            conversion_attempts = 3 if file_ext in ['.txt', '.text'] else 1
            
            for attempt in range(conversion_attempts):
                try:
                    logger.info(f"Conversion attempt {attempt + 1}/{conversion_attempts}")
                    
                    # Adjust prompt for retries on text files
                    if attempt > 0 and file_ext in ['.txt', '.text']:
                        extraction_prompt = "Extract to CSV format. Each line/paragraph is a separate row. Columns: text_content,section,page_number"
                    
                    # Call convert_file_to_csv - EXACT same function used by analyze_feedback
                    csv_content = await convert_file_to_csv(
                        file_content=file_content,
                        filename=filename,
                        prompt=extraction_prompt,
                        mode="extract"
                    )
                    
                    if csv_content and csv_content.strip():
                        # Extract text from the CSV
                        try:
                            df = pd.read_csv(StringIO(csv_content))
                            
                            # Look for text content in various possible columns
                            text_columns = ['text_content', 'review_text', 'text', 'content', 'feedback']
                            text_parts = []
                            
                            for col in text_columns:
                                if col in df.columns:
                                    parts = df[col].dropna().astype(str).tolist()
                                    if parts:
                                        text_parts.extend(parts)
                                        break
                            
                            # If no specific text columns found, use all text from all columns
                            if not text_parts:
                                for col in df.columns:
                                    if df[col].dtype == 'object':
                                        parts = df[col].dropna().astype(str).tolist()
                                        text_parts.extend(parts)
                            
                            if text_parts:
                                extracted_text = '\n'.join(text_parts)
                                logger.info(f"Successfully extracted {len(extracted_text)} characters from {filename}")
                                return extracted_text
                            else:
                                # Return the whole dataframe as string
                                return df.to_string()
                                
                        except Exception as parse_error:
                            logger.warning(f"Failed to parse CSV, returning raw content: {str(parse_error)}")
                            return csv_content
                            
                except FileConversionError as e:
                    logger.error(f"Conversion attempt {attempt + 1} failed: {str(e)}")
                    if attempt < conversion_attempts - 1:
                        await asyncio.sleep(1)
                    else:
                        logger.warning("All conversion attempts failed, trying direct text extraction")
                        break
                        
                except Exception as e:
                    logger.error(f"Unexpected error in conversion: {str(e)}")
                    if attempt < conversion_attempts - 1:
                        await asyncio.sleep(1)
                    else:
                        break
        
        # Fallback: Try multiple encodings for direct text extraction
        logger.info(f"Attempting direct text extraction for {filename}")
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1', 'utf-16', 'utf-16-le', 'utf-16-be', 'ascii']
        
        for encoding in encodings:
            try:
                text = file_content.decode(encoding)
                # Basic validation - check if we got reasonable text
                if len(text) > 0 and (text.isprintable() or '\n' in text or '\r' in text):
                    logger.info(f"Successfully decoded {filename} using {encoding} encoding")
                    return text
            except:
                continue
        
        # If all encodings fail, try with error handling
        for encoding in ['utf-8', 'latin-1']:
            try:
                text = file_content.decode(encoding, errors='ignore')
                if text.strip():
                    logger.warning(f"Decoded {filename} using {encoding} with errors ignored")
                    return text
            except:
                continue
        
        # Last resort - force UTF-8 with replacement
        logger.warning(f"Using UTF-8 with replacement characters for {filename}")
        return file_content.decode('utf-8', errors='replace')
        
    except Exception as e:
        logger.error(f"Error extracting content from {filename}: {str(e)}")
        # Return whatever we can extract
        return str(file_content)[:150000]  # Limit to prevent memory issues
async def generate_insight_summary_direct(
    client: AzureOpenAI,
    content: str,
    max_chars: int = 150000  # Approximately 50k tokens
) -> Dict[str, str]:
    '''
    Directly generates insight summary from raw content using Azure OpenAI.
    This bypasses all preprocessing and generates insights from any content.
    
    Args:
        client: Azure OpenAI client
        content: Raw content from any file type
        max_chars: Maximum characters to process
        
    Returns:
        Dictionary with user_loves, feature_request, pain_point, and overall_summary
    '''
    try:
        logger.info(f"Generating direct insight summary from {len(content)} characters")
        
        # Truncate if too long
        if len(content) > max_chars:
            logger.info(f"Truncating content from {len(content)} to {max_chars} characters")
            content = content[:max_chars]
        
        prompt = f'''
You are an elite customer feedback analyst at a top-tier consulting firm. Your expertise in behavioral psychology, data analysis, and business strategy allows you to extract game-changing insights from customer feedback that others miss.

🎯 CRITICAL MISSION: Perform an EXHAUSTIVE, FORENSIC-LEVEL analysis of the customer feedback below. Your insights will directly influence product strategy and could impact millions in revenue.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📊 RAW CUSTOMER FEEDBACK DATA FOR ANALYSIS:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{content}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

🔬 DEEP ANALYSIS FRAMEWORK:

1️⃣ **WHAT USERS GENUINELY LOVE** (user_loves)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
🔍 INVESTIGATION TECHNIQUES:
• Sentiment Analysis: Look for superlatives ("amazing", "love", "fantastic", "game-changer")
• Repeat Mentions: Features mentioned positively by multiple users
• Emotional Indicators: Exclamation marks, caps, emojis indicating delight
• Comparison Praise: "Better than [competitor]", "Finally a solution that..."
• Stickiness Factors: "Can't imagine working without", "Use it every day"
• Word-of-Mouth Indicators: "Recommended to my team", "Told my friends"

📈 WHAT TO EXTRACT:
• The #1 most beloved feature/aspect (with usage %)
• Unexpected delighters that create loyalty
• Features that solve critical pain points
• Elements that differentiate from competitors

✅ EXCELLENT EXAMPLES:
• "Real-time collaboration with instant sync loved by 73% of users, especially the conflict resolution"
• "The AI suggestions save 4.2 hours/week on average, called 'career-changing' by power users"
• "Mobile app's offline mode lets field workers stay productive, mentioned in 89% of 5-star reviews"

❌ AVOID THESE:
• "Users like the interface" (too vague, no specifics)
• "Good product" (meaningless without context)
• "Some positive feedback" (lazy analysis)

2️⃣ **MOST CRITICAL FEATURE REQUESTS** (feature_request)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
🔍 INVESTIGATION TECHNIQUES:
• Request Frequency: Count exact mentions of each request
• Language Patterns: "wish", "hope", "would love", "missing", "needs", "should have"
• Workaround Indicators: Users describing manual processes they do
• Competitive Mentions: "Like [competitor] has", "The only thing missing"
• Business Impact: Requests tied to revenue/efficiency gains
• Urgency Indicators: "desperately need", "dealbreaker", "considering alternatives"

📊 PRIORITIZATION MATRIX:
• HIGH IMPACT + HIGH FREQUENCY = Critical
• HIGH IMPACT + LOW FREQUENCY = Important  
• LOW IMPACT + HIGH FREQUENCY = Nice-to-have
• LOW IMPACT + LOW FREQUENCY = Backlog

✅ EXCELLENT EXAMPLES:
• "API access (62% of enterprise users) and Zapier integration (41%) to connect with existing workflows"
• "Dark mode requested by 156 users, especially those working night shifts in healthcare"
• "Bulk operations for processing 100+ items, blocking 3 enterprise deals worth $180k ARR"

❌ AVOID THESE:
• "Users want more features" (which features specifically?)
• "Some improvements requested" (quantify and specify)
• "Various enhancements" (too vague to action)

3️⃣ **CRITICAL PAIN POINTS & DEAL BREAKERS** (pain_point)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
🔍 INVESTIGATION TECHNIQUES:
• Emotion Mining: Frustration, anger, disappointment language
• Churn Indicators: "Switching to", "Looking for alternatives", "Canceled because"
• Severity Markers: "Critical", "Blocker", "Unusable", "Broken"
• Frequency + Recency: How often + how recently reported
• Cascade Analysis: One issue causing multiple problems
• Support Ticket Themes: Repeated issues in customer complaints
• Performance Metrics: Load times, crash rates, error frequencies

🚨 SEVERITY CLASSIFICATION:
• CRITICAL: Causes data loss, security issues, or complete workflow blockage
• HIGH: Significantly slows productivity or causes regular frustration  
• MEDIUM: Annoying but has workarounds
• LOW: Minor inconveniences

✅ EXCELLENT EXAMPLES:
• "Data sync fails 31% of time for files >10MB, causing 3-hour recovery process and data loss"
• "Mobile app crashes on Android 12+ devices (affecting 42% of Android users) when switching accounts"
• "Search function takes 8-15 seconds, making real-time customer service impossible for support teams"

❌ AVOID THESE:
• "Some bugs exist" (which bugs? how severe?)
• "Performance could be better" (specific metrics needed)
• "Users face difficulties" (too general to fix)

4️⃣ **STRATEGIC SYNTHESIS & RECOMMENDATIONS** (overall_summary)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📝 MUST INCLUDE (50-80 words):
• Overall sentiment ratio (positive:negative)
• Top 2 value drivers keeping users
• Top 2 issues risking churn
• Most requested feature category
• Strategic recommendation
• Urgency level (Critical/High/Medium)
• Potential revenue impact if addressed

🎯 ANALYSIS TECHNIQUES TO APPLY:

1. **FREQUENCY ANALYSIS**
   - Count exact mentions of each issue/feature/praise
   - Note clustering of feedback around specific topics
   - Identify outliers that might indicate emerging trends

2. **SENTIMENT INTENSITY MAPPING**
   - Score emotional language (1-10 scale)
   - Track sentiment progression through time
   - Identify emotional tipping points

3. **USER SEGMENT ANALYSIS**
   - Group feedback by user type (enterprise, SMB, individual)
   - Identify segment-specific patterns
   - Note different priorities across segments

4. **COMPETITIVE INTELLIGENCE**
   - Extract all competitor mentions
   - Identify feature gaps vs. competitors
   - Note switching triggers

5. **BUSINESS IMPACT CALCULATION**
   - Estimate revenue at risk from issues
   - Calculate potential gains from improvements
   - Prioritize by ROI potential

6. **PATTERN RECOGNITION**
   - Look for seasonal trends
   - Identify correlation between features
   - Spot unexpected connections

📏 QUALITY METRICS FOR YOUR ANALYSIS:
• Specificity: Use exact numbers, percentages, and examples
• Actionability: Product team can immediately act on insights
• Evidence-Based: Every claim backed by feedback data
• Prioritized: Clear what to fix/build first
• Quantified: Impact measured where possible

🎨 EXCEPTIONAL OUTPUT EXAMPLE:
{{
    "user_loves": "The AI-powered auto-tagging feature (praised by 67% of users) saves average 3.5 hours/week, with the predictive analytics dashboard helping teams spot trends 2 weeks earlier than competitors, driving strong word-of-mouth growth",
    "feature_request": "Highest demand for bulk import via CSV (127 requests, blocking 8 enterprise deals) and real-time collaborative editing (89 requests from teams 5+), followed by mobile offline sync for field sales teams (72 requests)",
    "pain_point": "Critical: Dashboard loading exceeds 12 seconds for 35% of users with 1000+ records, causing daily workflow interruptions; login failures spike to 23% during peak hours (9-11 AM EST), resulting in 5 enterprise accounts threatening to churn",
    "overall_summary": "Strong product-market fit evidenced by 72% positive sentiment, primarily driven by AI features and time savings. However, performance issues at scale risk 23% of MRR ($127K). Fixing load times and authentication would retain at-risk accounts. Bulk import feature could unlock $430K in stalled deals. Priority: Fix critical issues within 30 days, then build bulk features."
}}

⚠️ FINAL CRITICAL INSTRUCTIONS:
1. If feedback volume is low (<10 items), explicitly state this
2. Always prefer specific numbers over vague terms
3. Connect every insight to business impact
4. Never invent data - say "insufficient data" if needed
5. Focus on patterns, not individual complaints
6. Highlight both quick wins and strategic opportunities

🚫 ABSOLUTELY AVOID:
- Generic statements without evidence
- Insights that aren't actionable
- Mixing multiple issues into one
- Ignoring positive feedback when analyzing problems
- Missing critical business context

Return ONLY a JSON object with exactly these 4 fields. Make every word count toward driving business decisions.
'''
        
        # Retry logic
        max_retries = 3
        for attempt in range(max_retries):
            try:
                response = client.chat.completions.create(
                    model="gpt-4.1-mini",
                    messages=[
                        {
                            "role": "system", 
                            "content": "You are a senior customer feedback analyst. Extract insights based ONLY on provided content. Never make up information."
                        },
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.2,
                    max_tokens=1500,
                    response_format={"type": "json_object"}
                )
                
                result = json.loads(response.choices[0].message.content)
                
                # Validate required fields
                required_fields = ["user_loves", "feature_request", "pain_point", "overall_summary"]
                for field in required_fields:
                    if field not in result:
                        raise ValueError(f"Missing required field: {field}")
                    if not isinstance(result[field], str) or len(result[field].strip()) < 10:
                        raise ValueError(f"Invalid or empty field: {field}")
                
                logger.info("Successfully generated insight summary")
                return result
                
            except Exception as e:
                logger.error(f"Attempt {attempt + 1} failed: {str(e)}")
                if attempt < max_retries - 1:
                    await asyncio.sleep(1)
                else:
                    # Return default insights on final failure
                    return {
                        "user_loves": "Unable to determine what users love from the provided content",
                        "feature_request": "Unable to identify feature requests from the provided content",
                        "pain_point": "Unable to identify pain points from the provided content",
                        "overall_summary": "Analysis could not be completed due to processing errors. Please ensure the file contains valid customer feedback."
                    }
                    
    except Exception as e:
        logger.error(f"Critical error in generate_insight_summary_direct: {str(e)}")
        return {
            "user_loves": "Unable to analyze positive feedback aspects",
            "feature_request": "Unable to analyze feature requests",
            "pain_point": "Unable to analyze pain points",
            "overall_summary": "Critical error prevented feedback analysis. Please try again."
        }
def sanity_check_analysis_results(analysis_results: Dict[str, Any]) -> bool:
    '''
    Checks if analysis results from CSV flow are valid and meaningful.
    Returns True if results are good, False if we should use fallback.
    
    Args:
        analysis_results: Results from process_csv_data or process_excel_data
        
    Returns:
        Boolean indicating if results are valid
    '''
    try:
        # Check if we have the basic structure
        if not analysis_results:
            logger.warning("Sanity check: No analysis results")
            return False
            
        if "key_areas" not in analysis_results:
            logger.warning("Sanity check: Missing key_areas")
            return False
            
        key_areas = analysis_results.get("key_areas", [])
        
        # Must have at least 2 key areas
        if len(key_areas) < 2:
            logger.warning(f"Sanity check: Only {len(key_areas)} key areas found")
            return False
            
        # Check for diversity in area types
        area_types = set()
        for area in key_areas:
            if isinstance(area, dict):
                area_type = area.get("area_type", "issue")
                area_types.add(area_type)
        
        # Should have both features and issues
        if len(area_types) < 2:
            logger.warning(f"Sanity check: Limited diversity in area types: {area_types}")
            return False
            
        # Check classified feedback
        classified_feedback = analysis_results.get("classified_feedback", {})
        if not classified_feedback:
            logger.warning("Sanity check: No classified feedback")
            return False
            
        # Count total feedback items
        total_feedback = sum(len(feedbacks) for feedbacks in classified_feedback.values())
        if total_feedback < 5:  # Minimum threshold
            logger.warning(f"Sanity check: Only {total_feedback} total feedback items")
            return False
            
        # Check insight summary quality
        insights = analysis_results.get("insight_summary", {})
        if not insights:
            logger.warning("Sanity check: No insight summary")
            return False
            
        # Check if insights are meaningful (not just error messages)
        for key in ["user_loves", "feature_request", "pain_point"]:
            value = insights.get(key, "")
            if "unable to" in value.lower() or "error" in value.lower() or len(value) < 30:
                logger.warning(f"Sanity check: Poor quality insight for {key}: {value[:50]}...")
                return False
        
        logger.info("Sanity check: Analysis results are valid")
        return True
        
    except Exception as e:
        logger.error(f"Sanity check error: {str(e)}")
        return False
async def generate_complete_analysis_via_completions(
    client: AzureOpenAI,
    raw_content: str,
    source: str = None,
    return_raw_feedback: bool = False,
    max_retries: int = 5  # Increased retries
) -> Dict[str, Any]:
    '''
    Ultimate fallback: Generates complete analysis structure from raw content.
    Returns exact same format as normal processing flow.
    
    Args:
        client: Azure OpenAI client
        raw_content: Raw content from any file
        source: Source identifier
        return_raw_feedback: Whether to include raw feedback samples
        max_retries: Number of retry attempts
        
    Returns:
        Complete analysis matching expected format
    '''
    logger.info(f"Generating complete analysis via completions ({len(raw_content)} chars)")
    
    # Truncate if needed
    max_chars = 100000
    if len(raw_content) > max_chars:
        logger.info(f"Truncating content from {len(raw_content)} to {max_chars} chars")
        raw_content = raw_content[:max_chars]
    
    prompt = f'''
You are an expert customer feedback analyst. Your task is to analyze content and return EXACTLY the JSON structure shown below.

⚠️ CRITICAL: The UI will BREAK if you don't follow this EXACT structure. Every key must be spelled correctly. Every field must be present.

INPUT CONTENT:
{raw_content}

YOU MUST RETURN THIS EXACT STRUCTURE (copy the keys exactly):

{{
    "analysis_results": [
        {{
            "key_area": "string - 2-4 words describing the area",
            "customer_problem": "string - one clear sentence describing the issue/request",
            "number_of_users": integer - realistic count based on feedback,
            "type": "string - MUST be exactly 'feature' or 'issue' (lowercase)",
            "raw_feedbacks": [
                {{
                    "Received": "string - date in YYYY-MM-DD format or exactly 'N/A'",
                    "Source": "string - source or exactly '{source or 'N/A'}'",
                    "Customer Feedback": "string - the actual feedback text",
                    "Name": "string - customer name or exactly 'N/A'",
                    "sentiment_score": float - number between -1.0 and 1.0
                }}
            ]
        }}
    ],
    "summary": {{
        "total_feedback_items": integer - sum of all number_of_users,
        "total_key_areas": integer - count of items in analysis_results
    }},
    "insight_summary": {{
        "user_loves": "string - what users genuinely appreciate",
        "feature_request": "string - most requested features",
        "pain_point": "string - biggest problems",
        "overall_summary": "string - 50-80 word balanced summary"
    }}
}}

🔴 CRITICAL RULES:
1. EVERY key must be EXACTLY as shown above (case-sensitive)
2. "type" MUST be exactly "feature" or "issue" (lowercase, no variations)
3. Every raw_feedback item MUST have ALL 5 fields: Received, Source, Customer Feedback, Name, sentiment_score
4. Use EXACTLY "N/A" (not "n/a", "NA", "Unknown", etc.) for missing values
5. sentiment_score MUST be a float between -1.0 and 1.0
6. number_of_users MUST be an integer (not string, not float)
7. total_feedback_items MUST equal the sum of all number_of_users
8. total_key_areas MUST equal the count of analysis_results items

📋 CLASSIFICATION RULES:
- "feature" = requests for new functionality, enhancements, "I wish", "want", "need", "add"
- "issue" = problems, bugs, errors, complaints, "broken", "slow", "crash", "doesn't work"

📊 SENTIMENT SCORING:
- -1.0 to -0.7: Extremely negative (hate, terrible, worst)
- -0.6 to -0.3: Negative (frustrated, disappointed)
- -0.2 to 0.2: Neutral (okay, fine, average)
- 0.3 to 0.6: Positive (good, like, helpful)
- 0.7 to 1.0: Very positive (love, excellent, amazing)

📌 REAL EXAMPLE 1 (Reference Management System):
{{
    "analysis_results": [
        {{
            "key_area": "Reference Management",
            "customer_problem": "Users struggle to renumber references after adding a new one and want automatic placement in correct order.",
            "number_of_users": 12,
            "type": "feature",
            "raw_feedbacks": [
                {{
                    "Received": "N/A",
                    "Source": "N/A",
                    "Customer Feedback": "Users struggled to renumber references after adding a new one at the end of the list; they wanted it to be placed in its correct order.",
                    "Name": "N/A",
                    "sentiment_score": -0.038966266735042165
                }},
                {{
                    "Received": "N/A",
                    "Source": "N/A", 
                    "Customer Feedback": "PubMed/CrossRef integration was appreciated, as it ensured only valid references.",
                    "Name": "N/A",
                    "sentiment_score": 0.14952857496364988
                }}
            ]
        }},
        {{
            "key_area": "Affiliation Editing",
            "customer_problem": "Adding a new affiliation is not intuitive and users mistakenly click 'cite affiliation' instead.",
            "number_of_users": 8,
            "type": "issue",
            "raw_feedbacks": [
                {{
                    "Received": "N/A",
                    "Source": "N/A",
                    "Customer Feedback": "Adding a new affiliation was not intuitive, as they mistakenly clicked 'cite affiliation' instead.",
                    "Name": "N/A",
                    "sentiment_score": -0.040920173348273664
                }}
            ]
        }}
    ],
    "summary": {{
        "total_feedback_items": 20,
        "total_key_areas": 2
    }},
    "insight_summary": {{
        "user_loves": "Users appreciate the PubMed/CrossRef integration for ensuring valid references and the clear task list interface.",
        "feature_request": "Users most commonly request improved reference management with automatic renumbering and correct ordering when adding new references.",
        "pain_point": "The biggest frustration is the unintuitive affiliation editing process where users confuse 'cite affiliation' with adding new affiliations.",
        "overall_summary": "While users value the reference validation features, they struggle with reference ordering and affiliation management. The interface needs clearer labeling and automatic reference renumbering to improve the user experience."
    }}
}}

📌 REAL EXAMPLE 2 (Mobile App):
{{
    "analysis_results": [
        {{
            "key_area": "App Performance",
            "customer_problem": "App crashes frequently during photo uploads larger than 2MB causing data loss.",
            "number_of_users": 23,
            "type": "issue",
            "raw_feedbacks": [
                {{
                    "Received": "2024-01-15",
                    "Source": "App Store",
                    "Customer Feedback": "The app keeps crashing when I try to upload photos from my vacation. Lost all my edits!",
                    "Name": "John Smith",
                    "sentiment_score": -0.8
                }}
            ]
        }},
        {{
            "key_area": "Offline Mode",
            "customer_problem": "Users need ability to access and edit data without internet connection.",
            "number_of_users": 18,
            "type": "feature",
            "raw_feedbacks": [
                {{
                    "Received": "2024-01-20",
                    "Source": "Feature Request Form",
                    "Customer Feedback": "Please add offline mode so I can work during flights without losing progress.",
                    "Name": "Sarah Johnson",
                    "sentiment_score": -0.3
                }}
            ]
        }}
    ],
    "summary": {{
        "total_feedback_items": 41,
        "total_key_areas": 2
    }},
    "insight_summary": {{
        "user_loves": "Users appreciate the intuitive interface design and fast cloud sync when network is available.",
        "feature_request": "Offline mode is the most requested feature, with 44% of users needing to work without internet connectivity.",
        "pain_point": "App crashes during photo uploads affect 56% of users, causing significant data loss and workflow disruption.",
        "overall_summary": "While the interface is valued, critical stability issues and missing offline capability are driving users to competitors. Photo upload crashes need immediate attention to prevent user churn."
    }}
}}

🔄 PROCESSING STEPS:
1. Identify individual feedback items in the content
2. Group similar feedback into 3-8 key areas
3. Determine if each area is "feature" or "issue"
4. Count realistic number of users per area
5. Extract 1-5 sample feedbacks per area
6. Calculate sentiment scores for each feedback
7. Generate honest insights based on the feedback

⚠️ FINAL REMINDERS:
- If limited feedback exists, create fewer areas (minimum 1)
- Always base everything on actual content - no hallucination
- Every analysis_results item MUST have ALL fields
- Every raw_feedback MUST have EXACTLY these 5 fields: Received, Source, Customer Feedback, Name, sentiment_score
- Use EXACTLY "N/A" for missing values
- Double-check that total_feedback_items equals sum of number_of_users
- Ensure all keys are spelled EXACTLY as shown

7. Ensure all keys are spelled EXACTLY as shown

RESPOND WITH VALID JSON ONLY. The UI depends on this exact structure.'''
    
    for attempt in range(max_retries):
        try:
            logger.info(f"Attempt {attempt + 1}/{max_retries} for complete analysis")
            
            response = client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[
                    {
                        "role": "system",
                        "content": "Generate structured feedback analysis. Never hallucinate. Base everything on provided content."
                    },
                    {"role": "user", "content": prompt}
                ],
                temperature=0.2,
                max_tokens=4000,
                response_format={"type": "json_object"}
            )
            
            result = json.loads(response.choices[0].message.content)
            
            # STRICT validation of structure
            if not isinstance(result, dict):
                raise ValueError("Response is not a dictionary")
            
            # Check top-level keys
            required_top_keys = ["analysis_results", "summary", "insight_summary"]
            missing_keys = [k for k in required_top_keys if k not in result]
            if missing_keys:
                raise ValueError(f"Missing required top-level keys: {missing_keys}")
            
            # Validate analysis_results
            if not isinstance(result["analysis_results"], list):
                raise ValueError("analysis_results must be a list")
            
            if len(result["analysis_results"]) == 0:
                result["analysis_results"] = [{
                    "key_area": "General Feedback",
                    "customer_problem": "Limited feedback data available",
                    "number_of_users": 1,
                    "type": "issue",
                    "raw_feedbacks": [{
                        "Received": "N/A",
                        "Source": source or "N/A",
                        "Customer Feedback": "No specific feedback extracted",
                        "Name": "N/A",
                        "sentiment_score": 0.0
                    }]
                }]
            
            # Validate each analysis result
            for i, area in enumerate(result["analysis_results"]):
                # Check required fields
                required_area_keys = ["key_area", "customer_problem", "number_of_users", "type"]
                area_missing = [k for k in required_area_keys if k not in area]
                if area_missing:
                    raise ValueError(f"Analysis result {i} missing keys: {area_missing}")
                
                # Validate types
                if not isinstance(area["key_area"], str) or len(area["key_area"]) < 2:
                    area["key_area"] = f"Area {i+1}"
                
                if not isinstance(area["customer_problem"], str) or len(area["customer_problem"]) < 5:
                    area["customer_problem"] = "Unspecified customer concern"
                
                if not isinstance(area["number_of_users"], int) or area["number_of_users"] < 0:
                    area["number_of_users"] = 1
                
                if area["type"] not in ["feature", "issue"]:
                    # Try to infer from problem description
                    problem_lower = area["customer_problem"].lower()
                    if any(word in problem_lower for word in ["want", "need", "request", "add", "missing"]):
                        area["type"] = "feature"
                    else:
                        area["type"] = "issue"
                
                # Validate raw_feedbacks
                if "raw_feedbacks" not in area:
                    area["raw_feedbacks"] = []
                
                if not isinstance(area["raw_feedbacks"], list):
                    area["raw_feedbacks"] = []
                
                # Validate each feedback
                for j, feedback in enumerate(area["raw_feedbacks"]):
                    if not isinstance(feedback, dict):
                        area["raw_feedbacks"][j] = {
                            "Received": "N/A",
                            "Source": source or "N/A",
                            "Customer Feedback": str(feedback),
                            "Name": "N/A",
                            "sentiment_score": 0.0
                        }
                        continue
                    
                    # Check all required feedback fields
                    required_feedback_keys = ["Received", "Source", "Customer Feedback", "Name", "sentiment_score"]
                    for key in required_feedback_keys:
                        if key not in feedback:
                            if key == "Received":
                                feedback[key] = "N/A"
                            elif key == "Source":
                                feedback[key] = source or "N/A"
                            elif key == "Customer Feedback":
                                feedback[key] = "No feedback text"
                            elif key == "Name":
                                feedback[key] = "N/A"
                            elif key == "sentiment_score":
                                feedback[key] = 0.0
                    
                    # Validate sentiment score
                    try:
                        score = float(feedback["sentiment_score"])
                        feedback["sentiment_score"] = max(-1.0, min(1.0, score))
                    except:
                        feedback["sentiment_score"] = 0.0
                
                # Remove raw_feedbacks if not requested
                if not return_raw_feedback and "raw_feedbacks" in area:
                    del area["raw_feedbacks"]
            
            # Validate summary
            if not isinstance(result["summary"], dict):
                result["summary"] = {}
            
            total_users = sum(area.get("number_of_users", 0) for area in result["analysis_results"])
            result["summary"]["total_feedback_items"] = total_users
            result["summary"]["total_key_areas"] = len(result["analysis_results"])
            
            # Validate insight_summary
            if not isinstance(result["insight_summary"], dict):
                result["insight_summary"] = {}
            
            required_insights = ["user_loves", "feature_request", "pain_point", "overall_summary"]
            for field in required_insights:
                if field not in result["insight_summary"] or not isinstance(result["insight_summary"][field], str) or len(result["insight_summary"][field]) < 10:
                    result["insight_summary"][field] = f"Unable to determine {field.replace('_', ' ')} from provided content"
            
            logger.info("Successfully generated complete analysis")
            logger.info(f"Structure validation passed - {len(result['analysis_results'])} areas, {result['summary']['total_feedback_items']} total items")
            return result
            
        except json.JSONDecodeError as e:
            logger.error(f"Attempt {attempt + 1} failed - JSON parsing error at position {e.pos}: {str(e)}")
            logger.error(f"Raw response: {response.choices[0].message.content[:500]}...")
            if attempt < max_retries - 1:
                await asyncio.sleep(2 ** attempt)
                
        except ValueError as e:
            logger.error(f"Attempt {attempt + 1} failed - Validation error: {str(e)}")
            logger.error(f"Response structure: {json.dumps(result, indent=2)[:1000]}...")
            if attempt < max_retries - 1:
                await asyncio.sleep(2 ** attempt)
                
        except Exception as e:
            logger.error(f"Attempt {attempt + 1} failed - Unexpected error: {type(e).__name__}: {str(e)}")
            if attempt < max_retries - 1:
                await asyncio.sleep(2 ** attempt)
    
    # All retries failed
    logger.error("All attempts failed, returning minimal valid structure")
    
    # Return minimal but VALID structure with all required fields
    minimal_feedback = {
        "Received": "N/A",
        "Source": source or "N/A", 
        "Customer Feedback": "Unable to process feedback due to technical error",
        "Name": "N/A",
        "sentiment_score": 0.0
    }
    
    minimal_result = {
        "analysis_results": [{
            "key_area": "Processing Error",
            "customer_problem": "Unable to analyze feedback due to technical issues",
            "number_of_users": 0,
            "type": "issue"
        }],
        "summary": {
            "total_feedback_items": 0,
            "total_key_areas": 1
        },
        "insight_summary": {
            "user_loves": "Unable to process feedback to determine what users love",
            "feature_request": "Unable to process feedback to identify feature requests",
            "pain_point": "System encountered errors while processing the feedback data",
            "overall_summary": "Analysis could not be completed due to technical errors. Please ensure the file contains valid customer feedback and try again."
        }
    }
    
    # Add raw_feedbacks if requested
    if return_raw_feedback:
        minimal_result["analysis_results"][0]["raw_feedbacks"] = [minimal_feedback]
    
    return minimal_result
async def convert_file_to_csv(
    file_content: bytes,
    filename: str,
    prompt: Optional[str] = None,
    mode: str = "extract"
) -> str:
    '''
    Convert any file to CSV using the extract-reviews API.
    
    Args:
        file_content: Raw file content
        filename: Original filename
        prompt: Optional prompt for extraction
        mode: Mode for extraction (extract/generate/auto)
        
    Returns:
        CSV content as a string
        
    Raises:
        FileConversionError: If conversion fails at any point
    '''
    try:
        timeout = aiohttp.ClientTimeout(total=API_TIMEOUT)

        async with aiohttp.ClientSession(timeout=timeout) as session:
            for attempt in range(3):  # Retry up to 3 times
                # Rebuild FormData each try (cannot reuse a consumed FormData)
                form_data = aiohttp.FormData()
                form_data.add_field(
                    "file",
                    file_content,
                    filename=filename,
                    content_type=mimetypes.guess_type(filename)[0] or "application/octet-stream"
                )
                form_data.add_field("mode", mode)
                form_data.add_field("output_format", "csv")
                if prompt:
                    form_data.add_field("prompt", prompt)
                else:
                    default_prompt = '''
                    Extract customer feedback/reviews from this document.
                    
                    IMPORTANT RULES:
                    1. ONLY extract content that is explicitly customer feedback, reviews, or direct customer comments
                    2. DO NOT include: product descriptions, company statements, marketing copy, or general information
                    3. Each review MUST contain actual customer opinion or experience
                    4. Extract ALL customer reviews - do not skip any, even if they are positive/5-star
                    5. Do NOT create or infer reviews that don't exist in the document
                    6. Do NOT modify or summarize the review content - extract it exactly as written
                    
                    What qualifies as a review:
                    - Direct customer testimonials or feedback
                    - User comments with opinions about products/services
                    - Ratings accompanied by customer text
                    - Customer complaints or suggestions
                    - Review text from any platform (app stores, websites, surveys, etc.)
                    
                    Structure the data with these specific columns:
                    - user: Customer name or reviewer identifier (use 'Anonymous' if not provided)
                    - review_text: The EXACT feedback or review content as written by the customer
                    - date: When the review was submitted (use 'N/A' if not available)
                    - source: Where the review came from (platform, channel, etc.) (use 'N/A' if not clear)
                    - star_rating: Numerical rating (1-5 stars) or sentiment score (use 'N/A' if no rating)
                    
                    CRITICAL: 
                    - Return ONLY actual customer reviews/feedback
                    - Keep the EXACT wording - do not paraphrase or summarize
                    - Include EVERY review found in the document
                    - If unsure whether something is a customer review, check if it contains personal opinion or experience from a customer perspective
                    '''
                    form_data.add_field("prompt", default_prompt.strip())

                form_data.add_field("columns", "user,review_text,date,source,star_rating")

                logger.info(f"Calling extract-reviews API (attempt {attempt + 1}/3) for file: {filename}")

                try:
                    async with session.post(EXTRACT_REVIEWS_API_URL, data=form_data) as resp:
                        # If the extract-reviews endpoint always returns JSON:
                        if resp.status != 200:
                            error_text = await resp.text()
                            logger.error(f"API error (status {resp.status}): {error_text}")
                            if attempt == 2:
                                raise FileConversionError(
                                    f"Failed to convert file after 3 attempts. "
                                    f"Status: {resp.status}, Error: {error_text}"
                                )
                            await asyncio.sleep(2 ** attempt)
                            continue

                        # Parse JSON envelope
                        try:
                            payload = await resp.json()
                        except Exception:
                            raise FileConversionError("Failed to parse JSON from API response")

                        if payload.get("status") != "success":
                            msg = payload.get("message", "Unknown API error")
                            raise FileConversionError(f"API reported failure: {msg}")

                        download_url = payload.get("download_url")
                        if not download_url:
                            raise FileConversionError("No download_url in API response")

                    # ── Second request: fetch the actual CSV ──
                    async with session.get(download_url) as csv_resp:
                        if csv_resp.status != 200:
                            error_text = await csv_resp.text()
                            raise FileConversionError(
                                f"CSV download failed (status {csv_resp.status}): {error_text}"
                            )
                        csv_content = await csv_resp.text()

                    # ── Validate the CSV content ──
                    if not csv_content.strip():
                        raise FileConversionError("Downloaded CSV content is empty")

                    lines = csv_content.strip().splitlines()
                    if len(lines) < 2:  # Expect at least header + 1 row
                        raise FileConversionError("Downloaded CSV appears invalid (too few lines)")

                    logger.info(
                        f"Successfully converted {filename} to CSV ({len(csv_content)} chars, {len(lines)} lines)"
                    )
                    return csv_content

                except asyncio.TimeoutError:
                    logger.error(f"Timeout on attempt {attempt + 1}")
                    if attempt == 2:
                        raise FileConversionError(f"API timeout after {API_TIMEOUT} seconds")
                    await asyncio.sleep(2 ** attempt)

                except aiohttp.ClientError as e:
                    logger.error(f"Client error on attempt {attempt + 1}: {e}")
                    if attempt == 2:
                        raise FileConversionError(f"Network error: {str(e)}")
                    await asyncio.sleep(2 ** attempt)

        # If we exit the loop without returning, raise a generic error
        raise FileConversionError("Failed to convert file after all retries")

    except FileConversionError:
        # Re-raise our custom exception without modification
        raise
    except Exception as e:
        # Catch anything unexpected
        logger.error(f"Error converting file {filename}: {e}")
        raise FileConversionError(f"Failed to convert file: {e}")
async def generate_summary(
    client: AzureOpenAI,
    review_embeddings: List[List[float]],
    review_texts: List[str],
    classified_feedback: Dict[str, List] = None,
    key_areas: List[Dict] = None,
    slow_mode: bool = False,
    skip_insights: bool = False 
) -> Dict[str, str]:
    '''
    Analyzes customer feedback to identify what users love, what features they request, their pain points,
    and generates an overall summary of the feedback.
    
    Args:
        client: Azure OpenAI client
        review_embeddings: List of review embedding vectors
        review_texts: List of original review text
        classified_feedback: Dictionary mapping areas to feedback (optional, for area_type filtering)
        key_areas: List of key areas with area_type information
        slow_mode: If True, use slower but more methodical processing
        
    Returns:
        Dictionary with user_loves, feature_request, pain_point, and overall_summary
    '''
    
    try:
        logger.info(f"Generating summary insights from {len(review_texts)} reviews")
        if skip_insights:
            logger.info("Skipping insights generation as skip_insights=True")
            return {
                "user_loves": "Insights will be generated via completions",
                "feature_request": "Insights will be generated via completions",
                "pain_point": "Insights will be generated via completions",
                "overall_summary": "Insights will be generated via completions"
            }
        # More robust check for insufficient data
        if len(review_texts) < 2:
            logger.warning(f"Insufficient reviews ({len(review_texts)}) for reliable summary generation")
            return {
                "user_loves": "Not enough reviews to determine what users love",
                "feature_request": "Not enough reviews to identify feature requests",
                "pain_point": "Not enough reviews to identify pain points",
                "overall_summary": "Not enough reviews to generate an overall summary"
            }
            
        # Check for meaningful content
        meaningful_reviews = [r for r in review_texts if r and len(r.strip()) > 10]
        if len(meaningful_reviews) < 2:
            logger.warning(f"Only {len(meaningful_reviews)} meaningful reviews found - not enough for summary")
            return {
                "user_loves": "Not enough meaningful reviews to determine what users love",
                "feature_request": "Not enough meaningful reviews to identify feature requests",
                "pain_point": "Not enough meaningful reviews to identify pain points",
                "overall_summary": "Not enough meaningful reviews to generate an overall summary"
            }
        
        # If we have area_type information, filter reviews accordingly
        feature_reviews = []
        issue_reviews = []
        
        if classified_feedback and key_areas:
            # Create area_type mapping
            area_type_map = {area["area"]: area.get("area_type", "issue") for area in key_areas}
            
            # Collect reviews by type
            for area, feedbacks in classified_feedback.items():
                area_type = area_type_map.get(area, "issue")
                for feedback in feedbacks:
                    # Extract the review text from the feedback dict
                    if isinstance(feedback, dict):
                        review_text = feedback.get("Customer Feedback", feedback.get("text", ""))
                    else:
                        review_text = str(feedback)
                    
                    if area_type == "feature":
                        feature_reviews.append(review_text)
                    else:
                        issue_reviews.append(review_text)
        
        # If no classified feedback provided, use all reviews for both
        if not feature_reviews and not issue_reviews:
            feature_reviews = meaningful_reviews
            issue_reviews = meaningful_reviews
            
        logger.info(f"Processing summaries - Feature reviews: {len(feature_reviews)}, Issue reviews: {len(issue_reviews)}")
        
        # Use the meaningful reviews only
        review_texts = meaningful_reviews
        
        # Define anchor texts for each category (same as before)
        user_loves_anchors = [
            "I really love this feature",
            "This is my favorite part about the product",
            "The best thing about this service is",
            "What I appreciate most is",
            "This works exceptionally well",
            "I'm impressed by how well it handles",
            "The most useful feature for me is",
            "This makes my life so much easier"
        ]
        
        feature_request_anchors = [
            "I wish this product had",
            "It would be great if you could add",
            "Please consider adding a feature for",
            "The one thing this product needs is",
            "I would love to see this improved with",
            "Could you please implement",
            "This would be perfect if it included",
            "What's missing is"
        ]
        
        pain_point_anchors = [
            "The most frustrating thing is",
            "This doesn't work properly when",
            "I'm having trouble with",
            "The biggest problem I have is",
            "This is consistently failing to",
            "It's really annoying when",
            "What needs to be fixed urgently is",
            "I can't stand how it"
        ]
        
        # Generate embeddings for anchor texts
        all_anchors = user_loves_anchors + feature_request_anchors + pain_point_anchors
        anchor_embeddings = await get_embeddings(client, all_anchors)
        
        # Split anchor embeddings by category
        user_loves_embeddings = anchor_embeddings[:len(user_loves_anchors)]
        feature_request_embeddings = anchor_embeddings[len(user_loves_anchors):len(user_loves_anchors) + len(feature_request_anchors)]
        pain_point_embeddings = anchor_embeddings[len(user_loves_anchors) + len(feature_request_anchors):]
        
        # Calculate average embeddings for each category
        user_loves_embedding = np.mean(user_loves_embeddings, axis=0)
        feature_request_embedding = np.mean(feature_request_embeddings, axis=0)
        pain_point_embedding = np.mean(pain_point_embeddings, axis=0)
        
        # Convert to numpy arrays for vectorized operations
        review_matrix = np.array(review_embeddings[:len(meaningful_reviews)])  # Use only meaningful reviews
        user_loves_vec = np.array(user_loves_embedding)
        feature_request_vec = np.array(feature_request_embedding)
        pain_point_vec = np.array(pain_point_embedding)
        
        # Normalize vectors
        user_loves_vec = user_loves_vec / np.linalg.norm(user_loves_vec)
        feature_request_vec = feature_request_vec / np.linalg.norm(feature_request_vec)
        pain_point_vec = pain_point_vec / np.linalg.norm(pain_point_vec)
        
        review_norms = np.linalg.norm(review_matrix, axis=1, keepdims=True)
        # Avoid division by zero
        review_norms = np.where(review_norms == 0, 1e-10, review_norms)
        review_matrix_norm = review_matrix / review_norms
        
        # Calculate similarity scores for each category
        user_loves_scores = np.dot(review_matrix_norm, user_loves_vec)
        feature_request_scores = np.dot(review_matrix_norm, feature_request_vec)
        pain_point_scores = np.dot(review_matrix_norm, pain_point_vec)
        
        # Function to select top reviews with minimum similarity threshold
        def get_top_reviews(scores, texts, n=15, min_threshold=0.2):
            if len(scores) == 0:
                return []
            
            # For small datasets, adjust threshold dynamically
            if len(texts) <= 5:
                min_threshold = 0.2  # Lower threshold for small datasets
            
            # Get indices of reviews with scores above threshold
            valid_indices = np.where(scores > min_threshold)[0]
            
            # If no reviews meet the threshold and we have few reviews, use top scoring ones
            if len(valid_indices) == 0 and len(texts) <= 5:
                # Take the top 50% of reviews by score
                num_to_take = max(1, len(texts) // 2)
                sorted_indices = np.argsort(scores)[-num_to_take:]
                return [texts[i] for i in sorted_indices]
            
            # If no reviews meet the threshold, return empty list
            if len(valid_indices) == 0:
                return []
                
            # Sort the valid indices by score
            sorted_indices = valid_indices[np.argsort(scores[valid_indices])[-min(n, len(valid_indices)):]]
            
            # Return the corresponding reviews
            return [texts[i] for i in sorted_indices]
        
        # Check similarity scores before collecting reviews
        top_user_loves = get_top_reviews(user_loves_scores, review_texts)
        
        # For feature requests, prefer reviews from feature areas
        if feature_reviews:
            top_feature_requests = feature_reviews[:15]  # Take top feature reviews
        else:
            top_feature_requests = get_top_reviews(feature_request_scores, review_texts)
        
        # For pain points, prefer reviews from issue areas
        if issue_reviews:
            top_pain_points = issue_reviews[:15]  # Take top issue reviews
        else:
            top_pain_points = get_top_reviews(pain_point_scores, review_texts)
        
        # Log the number of matching reviews found for each category
        logger.info(f"Found {len(top_user_loves)} user loves, {len(top_feature_requests)} feature requests, {len(top_pain_points)} pain points")
        
        # Function to generate summary using OpenAI (same as before, just updated prompts)
        async def generate_category_summary(reviews, category_name):
            # If no matching reviews found, return appropriate message
            if not reviews:
                return f"No clear {category_name} identified in the reviews"
            
            # Take a maximum of 15 reviews to avoid token limits
            sample_reviews = reviews[:15]
            
            # Combine reviews into a single text
            combined_reviews = "\n".join([f"- {review}" for review in sample_reviews])
            
            prompt_templates = {
                "user_loves": f'''
                Below are {len(sample_reviews)} customer reviews that indicate features or aspects users love about the product:
                
                {combined_reviews}
                
                Based ONLY on these specific reviews, summarize in ONE concise sentence what users love most about the product.
                Focus on concrete features or aspects mentioned in the reviews, not general satisfaction.
                
                IMPORTANT: If the reviews don't clearly indicate what users love or contain too little information, 
                respond with EXACTLY this phrase: "No clear indication of what users love in the provided reviews."
                DO NOT make up or hallucinate information not present in the reviews.
                ''',
                
                "feature_request": f'''
                Below are {len(sample_reviews)} customer reviews that suggest features users would like to see added:
                
                {combined_reviews}
                
                Based ONLY on these specific reviews, summarize in ONE concise sentence what feature or improvement users most commonly request.
                Focus on clear feature requests explicitly mentioned in the reviews, not general complaints.
                
                IMPORTANT: If the reviews don't clearly indicate feature requests or contain too little information, 
                respond with EXACTLY this phrase: "No clear feature requests identified in the provided reviews."
                DO NOT make up or hallucinate information not present in the reviews.
                ''',
                
                "pain_point": f'''
                Below are {len(sample_reviews)} customer reviews that highlight pain points or issues:
                
                {combined_reviews}
                
                Based ONLY on these specific reviews, summarize in ONE concise sentence what is the biggest pain point or issue users are experiencing.
                Focus on specific problems explicitly mentioned in the reviews, not general dissatisfaction.
                
                IMPORTANT: If the reviews don't clearly indicate pain points or contain too little information, 
                respond with EXACTLY this phrase: "No clear pain points identified in the provided reviews."
                DO NOT make up or hallucinate information not present in the reviews.
                ''',
                
                "overall_summary": f'''
                Below are {len(sample_reviews)} selected customer reviews:
                
                {combined_reviews}
                
                ONLY if there is sufficient meaningful feedback in these reviews, provide a comprehensive yet concise summary 
                (approximately 50 words) of the overall customer sentiment and key themes mentioned in these specific reviews.
                
                Your summary should:
                1. Reflect only what's actually mentioned in these specific reviews
                2. Capture the main themes and sentiments expressed
                3. Be a fresh analysis, not a combination of previous categorizations
                4. Stay neutral in your analysis terminology and avoid using category terms like "pain points"
                {f"5. Note that approximately {len(feature_reviews)} reviews are about feature requests and {len(issue_reviews)} are about issues/problems" if feature_reviews or issue_reviews else ""}
                
                IMPORTANT: If the reviews contain insufficient meaningful feedback or are too sparse, 
                respond with EXACTLY this phrase: "Insufficient meaningful feedback to generate an overall summary."
                DO NOT make up or hallucinate information not present in the reviews.
                '''
            }
            
            try:
                response = client.chat.completions.create(
                    model="gpt-4.1-mini",
                    messages=[
                        {"role": "system", "content": "You are a customer insight specialist who extracts clear, actionable insights from reviews. You NEVER make up information that isn't explicitly present in the reviews. If there isn't sufficient information, you clearly state that."},
                        {"role": "user", "content": prompt_templates[category_name]}
                    ],
                    temperature=0.2,  # Lower temperature to reduce creativity
                    max_tokens=100
                )
                
                summary = response.choices[0].message.content.strip()
                # Remove any quotes that might be around the summary
                summary = summary.strip('"\'')
                return summary
                
            except Exception as e:
                logger.error(f"Error generating summary for {category_name}: {str(e)}")
                return f"Unable to generate summary for {category_name} due to an error"
        
        # Generate summaries in parallel for individual categories
        user_loves_summary_task = asyncio.create_task(generate_category_summary(top_user_loves, "user_loves"))
        feature_request_summary_task = asyncio.create_task(generate_category_summary(top_feature_requests, "feature_request"))
        pain_point_summary_task = asyncio.create_task(generate_category_summary(top_pain_points, "pain_point"))
        
        # Wait for all category summaries to complete
        user_loves_summary = await user_loves_summary_task
        feature_request_summary = await feature_request_summary_task
        pain_point_summary = await pain_point_summary_task
        
        # Create a fresh combined set of reviews for overall summary (up to 10 from each category)
        # Select reviews but avoid duplicates to ensure we're getting a diverse set
        overall_reviews = []
        
        # Function to add reviews while avoiding duplicates
        def add_unique_reviews(reviews_list, source_reviews, count=10):
            # Create a set of existing review texts for easy comparison
            existing_texts = set(r.lower() for r in reviews_list)
            added = 0
            
            for review in source_reviews:
                # Only add if not already in the list (case-insensitive comparison)
                if review.lower() not in existing_texts:
                    reviews_list.append(review)
                    existing_texts.add(review.lower())
                    added += 1
                    if added >= count:
                        break
        
        # Add unique reviews from each category
        if top_user_loves:
            add_unique_reviews(overall_reviews, top_user_loves, 10)
        if top_feature_requests:
            add_unique_reviews(overall_reviews, top_feature_requests, 10)
        if top_pain_points:
            add_unique_reviews(overall_reviews, top_pain_points, 10)
        
        # Only generate overall summary if we have enough meaningful data
        if len(overall_reviews) >= 1:
            overall_summary = await generate_category_summary(overall_reviews, "overall_summary")
        else:
            overall_summary = "Insufficient meaningful feedback to generate an overall summary"
        
        # Return the summaries
        return {
            "user_loves": user_loves_summary,
            "feature_request": feature_request_summary,
            "pain_point": pain_point_summary,
            "overall_summary": overall_summary
        }
        
    except Exception as e:
        logger.error(f"Error generating summary insights: {str(e)}\n{traceback.format_exc()}")
        # Return default values on error
        return {
            "user_loves": "Unable to determine what users love due to an error",
            "feature_request": "Unable to identify feature requests due to an error",
            "pain_point": "Unable to identify pain points due to an error",
            "overall_summary": "Unable to generate an overall summary due to an error"
        }
async def classify_feedback_by_type(
    client: AzureOpenAI,
    feedback_embeddings: List[List[float]], 
    feedback_texts: List[str]
) -> List[str]:
    '''
    Classifies each feedback item as either a feature request or an issue.
    
    Args:
        client: Azure OpenAI client
        feedback_embeddings: List of feedback embedding vectors
        feedback_texts: List of feedback text
        
    Returns:
        List of classifications: ["feature", "issue", ...] for each feedback
    '''
    try:
        logger.info(f"Classifying {len(feedback_texts)} feedback items as feature/issue")
        
        # Define anchor texts for features and issues
        feature_anchors = [
            "I wish this product had",
            "It would be great if you could add", 
            "Please implement",
            "Would love to see",
            "Feature request",
            "Missing functionality",
            "Should include",
            "Needs to have",
            "Consider adding",
            "Enhancement suggestion",
            "Can you add",
            "Want this feature",
            "Requesting feature",
            "Please add support for"
        ]
        
        issue_anchors = [
            "This is broken",
            "Doesn't work properly",
            "Failed to",
            "Error occurs when", 
            "Problem with",
            "Bug in the system",
            "Crashes when",
            "Not functioning",
            "Glitch in",
            "Malfunction",
            "Cannot access",
            "Stopped working",
            "Issues with",
            "Not responding"
        ]
        
        # Generate embeddings for anchor texts
        all_anchors = feature_anchors + issue_anchors
        anchor_embeddings = await get_embeddings(client, all_anchors)
        
        # Split anchor embeddings
        feature_embeddings = anchor_embeddings[:len(feature_anchors)]
        issue_embeddings = anchor_embeddings[len(feature_anchors):]
        
        # Calculate average embeddings for each type
        feature_embedding = np.mean(feature_embeddings, axis=0)
        issue_embedding = np.mean(issue_embeddings, axis=0)
        
        # Convert to numpy arrays
        feedback_matrix = np.array(feedback_embeddings)
        feature_vec = np.array(feature_embedding)
        issue_vec = np.array(issue_embedding)
        
        # Normalize vectors
        feature_vec = feature_vec / np.linalg.norm(feature_vec)
        issue_vec = issue_vec / np.linalg.norm(issue_vec)
        
        feedback_norms = np.linalg.norm(feedback_matrix, axis=1, keepdims=True)
        feedback_norms = np.where(feedback_norms == 0, 1e-10, feedback_norms)
        feedback_matrix_norm = feedback_matrix / feedback_norms
        
        # Calculate similarity scores
        feature_scores = np.dot(feedback_matrix_norm, feature_vec)
        issue_scores = np.dot(feedback_matrix_norm, issue_vec)
        
        # Classify based on higher similarity
        classifications = []
        for i in range(len(feedback_texts)):
            if feature_scores[i] > issue_scores[i]:
                classifications.append("feature")
            else:
                classifications.append("issue")
        
        # Log classification statistics
        feature_count = classifications.count("feature")
        issue_count = classifications.count("issue")
        logger.info(f"Classified feedback: {feature_count} features, {issue_count} issues")
        
        return classifications
        
    except Exception as e:
        logger.error(f"Error classifying feedback by type: {str(e)}")
        # Return all as issues as fallback
        return ["issue"] * len(feedback_texts)
def standardize_dataframe(df: pd.DataFrame, columns: Dict[str, str], source: str = None) -> pd.DataFrame:
    '''
    Standardizes a dataframe to have consistent column names and structure.
    Handles cases where column detection failed or columns don't exist.
    
    Args:
        df: Original DataFrame
        columns: Dictionary mapping column types to actual column names
        source: Optional source value to override the Source column
        
    Returns:
        Standardized DataFrame with consistent columns
    '''
    try:
        logger.info("Standardizing DataFrame to consistent column structure")
        
        # Handle case where df has no columns
        if len(df.columns) == 0:
            logger.warning("DataFrame has no columns, creating empty standardized DataFrame")
            return pd.DataFrame({
                "Received": ["N/A"],
                "Name": ["N/A"],
                "Customer Feedback": ["No data available"],
                "Source": [source if source else "N/A"]
            })
        
        # Create a copy to avoid modifying the original
        result_df = df.copy()
        
        # Clean column names - remove trailing/leading whitespace
        result_df.columns = [col.strip() if isinstance(col, str) else str(col) for col in result_df.columns]
        
        # Update column references in the columns dict if they had trailing spaces
        for col_type in columns:
            if columns[col_type] and isinstance(columns[col_type], str):
                columns[col_type] = columns[col_type].strip()
        
        # Validate that referenced columns actually exist
        for col_type, col_name in list(columns.items()):
            if col_name and col_name not in result_df.columns:
                logger.warning(f"Column '{col_name}' referenced for {col_type} but not found in DataFrame")
                columns[col_type] = None
        
        # Handle specific column name mappings for common patterns
        specific_mappings = {
            "Reviwer": "Name",
            "Reviewer": "Name", 
            "Customer Name": "Name",
            "User": "Name",
            "Username": "Name",
            "Author": "Name",
            "Submitted By": "Name",
            "Customer": "Name",
            "Reviewed By": "Name",
            
            "Review": "Customer Feedback",
            "Reviews": "Customer Feedback",
            "Comment": "Customer Feedback",
            "Comments": "Customer Feedback",
            "Feedback": "Customer Feedback",
            "Review Text": "Customer Feedback",
            "Customer Review": "Customer Feedback",
            "User Feedback": "Customer Feedback",
            "Message": "Customer Feedback",
            "Description": "Customer Feedback",
            "review_text": "Customer Feedback",
            
            "Date": "Received",
            "Timestamp": "Received",
            "Time": "Received",
            "Created At": "Received",
            "Submitted": "Received",
            "Review Date": "Received",
            "Posted": "Received",
            "Date Posted": "Received",
            "date": "Received",
            
            "Platform": "Source",
            "Channel": "Source",
            "Website": "Source",
            "Origin": "Source",
            "Review Source": "Source",
            "From": "Source",
            "Via": "Source",
            "source": "Source"
        }
        
        # Apply specific mappings first
        for old_col, new_col in specific_mappings.items():
            if old_col in result_df.columns:
                result_df = result_df.rename(columns={old_col: new_col})
                # Update the columns dict accordingly
                for col_type, col_name in columns.items():
                    if col_name == old_col:
                        columns[col_type] = new_col
        
        # Create a mapping from original column names to standard names
        column_mapping = {}
        
        # Only map columns that were actually identified and exist
        if columns.get("feedback_col") and columns["feedback_col"] in result_df.columns:
            column_mapping[columns["feedback_col"]] = "Customer Feedback"
        if columns.get("received_col") and columns["received_col"] in result_df.columns:
            column_mapping[columns["received_col"]] = "Received"
        if columns.get("name_col") and columns["name_col"] in result_df.columns:
            column_mapping[columns["name_col"]] = "Name"
        if columns.get("source_col") and columns["source_col"] in result_df.columns:
            column_mapping[columns["source_col"]] = "Source"
            
        # Rename the columns that were identified
        if column_mapping:
            result_df = result_df.rename(columns=column_mapping)
        
        # Get current columns and desired columns
        current_columns = set(result_df.columns)
        desired_columns = {"Received", "Name", "Customer Feedback", "Source"}
        
        # Ensure we have a Customer Feedback column
        if "Customer Feedback" not in current_columns:
            # Try to find any text column to use as feedback
            feedback_found = False
            for col in result_df.columns:
                if result_df[col].dtype == 'object' and col not in desired_columns:
                    # Check if it has meaningful content
                    non_null = result_df[col].dropna()
                    if len(non_null) > 0:
                        avg_length = non_null.astype(str).str.len().mean()
                        if avg_length > 10:  # Likely contains feedback
                            result_df = result_df.rename(columns={col: "Customer Feedback"})
                            feedback_found = True
                            logger.info(f"Using column '{col}' as Customer Feedback")
                            break
            
            if not feedback_found:
                # Use any text column as feedback
                for col in result_df.columns:
                    if result_df[col].dtype == 'object' and col not in desired_columns:
                        result_df = result_df.rename(columns={col: "Customer Feedback"})
                        feedback_found = True
                        logger.info(f"Using column '{col}' as Customer Feedback (fallback)")
                        break
            
            if not feedback_found:
                # Use the first column as feedback as last resort
                if len(result_df.columns) > 0:
                    first_col = result_df.columns[0]
                    result_df = result_df.rename(columns={first_col: "Customer Feedback"})
                    logger.warning(f"Using first column '{first_col}' as Customer Feedback")
                else:
                    result_df["Customer Feedback"] = "No feedback data"
        
        # Update current columns after potential renaming
        current_columns = set(result_df.columns)
        
        # Add any missing desired columns with N/A values
        for col in desired_columns:
            if col not in current_columns:
                logger.info(f"Adding missing column '{col}' with N/A values")
                result_df[col] = "N/A"
        
        # Override Source column if source parameter is provided
        if source:
            logger.info(f"Overriding Source column with value: {source}")
            result_df["Source"] = source
                
        # Keep only the desired columns plus any rating column that exists
        columns_to_keep = list(desired_columns)
        
        # Add the rating column if it exists and is different from other columns
        if columns.get("rating_col"):
            # Check if the rating column still exists after renaming
            if columns["rating_col"] in result_df.columns:
                columns_to_keep.append(columns["rating_col"])
            else:
                # Check if it was renamed to a standard name
                for old_name, new_name in column_mapping.items():
                    if old_name == columns["rating_col"] and new_name in result_df.columns:
                        # It was already renamed to a standard column
                        break
                else:
                    # Look for any remaining numeric column that could be rating
                    for col in result_df.columns:
                        if col not in desired_columns and pd.api.types.is_numeric_dtype(result_df[col]):
                            columns_to_keep.append(col)
                            logger.info(f"Including numeric column '{col}' as potential rating")
                            break
        
        # Filter to keep only desired columns
        available_columns = [col for col in columns_to_keep if col in result_df.columns]
        result_df = result_df[available_columns]
        
        # Ensure we have at least one row
        if len(result_df) == 0:
            logger.warning("DataFrame is empty after standardization, adding placeholder row")
            result_df = pd.DataFrame({
                "Received": ["N/A"],
                "Name": ["N/A"],
                "Customer Feedback": ["No feedback data available"],
                "Source": [source if source else "N/A"]
            })
        
        # Validate the Customer Feedback column has actual content
        if "Customer Feedback" in result_df.columns:
            # Replace empty strings and whitespace with NaN
            result_df["Customer Feedback"] = result_df["Customer Feedback"].replace(r'^\s*$', np.nan, regex=True)
            # If all feedback is empty, add a placeholder
            if result_df["Customer Feedback"].isna().all():
                logger.warning("All Customer Feedback values are empty")
                result_df["Customer Feedback"] = result_df["Customer Feedback"].fillna("No feedback content")
        
        logger.info(f"Standardized DataFrame structure: {list(result_df.columns)}")
        logger.info(f"Standardized DataFrame shape: {result_df.shape}")
        
        return result_df
        
    except Exception as e:
        logger.error(f"Error standardizing DataFrame: {str(e)}")
        # Return a basic dataframe with the required columns if anything goes wrong
        try:
            fallback_df = pd.DataFrame({
                "Received": ["N/A"] * len(df) if len(df) > 0 else ["N/A"],
                "Name": ["N/A"] * len(df) if len(df) > 0 else ["N/A"],
                "Customer Feedback": df.iloc[:, 0].astype(str) if len(df.columns) > 0 and len(df) > 0 else ["Error processing data"],
                "Source": [source if source else "N/A"] * (len(df) if len(df) > 0 else 1)
            })
            return fallback_df
        except:
            # Ultimate fallback
            return pd.DataFrame({
                "Received": ["N/A"],
                "Name": ["N/A"],
                "Customer Feedback": ["Error processing data"],
                "Source": [source if source else "N/A"]
            })
async def get_sentiment_anchor_embeddings(client: AzureOpenAI) -> Tuple[List[float], List[float]]:
    '''
    Generates and caches embeddings for positive and negative sentiment anchors.
    
    Args:
        client: Azure OpenAI client
        
    Returns:
        Tuple of (positive_anchor_embedding, negative_anchor_embedding)
    '''
    global _positive_anchor_embedding, _negative_anchor_embedding
    
    # Return cached embeddings if already generated
    if _positive_anchor_embedding is not None and _negative_anchor_embedding is not None:
        return _positive_anchor_embedding, _negative_anchor_embedding
    
    # Define positive and negative sentiment anchor texts
    positive_anchors = [
    "excellent",
    "amazing",
    "fantastic",
    "great",
    "perfect",
    "love it",
    "impressive",
    "satisfied",
    "wonderful",
    "delightful"
    ]
    
    # Better negative anchors (sentiment-pure)
    negative_anchors = [
        "terrible",
        "awful",
        "horrible",
        "disappointing",
        "frustrating",
        "useless",
        "hate it",
        "poor",
        "annoying",
        "dreadful"
    ]
        
    try:
        # Generate embeddings for all anchor texts
        all_anchors = positive_anchors + negative_anchors
        anchor_embeddings = await get_embeddings(client, all_anchors)
        
        # Calculate average embeddings for positive and negative anchors
        positive_embeddings = anchor_embeddings[:len(positive_anchors)]
        negative_embeddings = anchor_embeddings[len(positive_anchors):]
        
        # Calculate average embeddings
        _positive_anchor_embedding = np.mean(positive_embeddings, axis=0).tolist()
        _negative_anchor_embedding = np.mean(negative_embeddings, axis=0).tolist()
        
        logger.info("Generated sentiment anchor embeddings successfully")
        
        return _positive_anchor_embedding, _negative_anchor_embedding
    
    except Exception as e:
        logger.error(f"Error generating sentiment anchor embeddings: {str(e)}")
        
        # Generate fallback embeddings with the same dimension as expected
        dim = 1536  # text-embedding-3-large dimension
        _positive_anchor_embedding = list(np.random.normal(0.5, 0.1, dim))
        _negative_anchor_embedding = list(np.random.normal(-0.5, 0.1, dim))
        
        return _positive_anchor_embedding, _negative_anchor_embedding

async def calculate_sentiment_scores(review_embeddings: List[List[float]], client: AzureOpenAI) -> List[float]:
    '''
    Calculates sentiment scores for reviews by comparing their embeddings to positive/negative anchors.
    
    Args:
        review_embeddings: List of review embedding vectors
        client: Azure OpenAI client for generating anchor embeddings if needed
        
    Returns:
        List of sentiment scores between -1 and 1 for each review
    '''
    try:
        # Get positive and negative anchor embeddings
        positive_embedding, negative_embedding = await get_sentiment_anchor_embeddings(client)
        
        # Convert to numpy arrays for vectorized operations
        review_matrix = np.array(review_embeddings)
        pos_embedding = np.array(positive_embedding)
        neg_embedding = np.array(negative_embedding)
        
        # Normalize anchor embeddings
        pos_embedding = pos_embedding / np.linalg.norm(pos_embedding)
        neg_embedding = neg_embedding / np.linalg.norm(neg_embedding)
        
        # Normalize review embeddings
        review_norms = np.linalg.norm(review_matrix, axis=1, keepdims=True)
        review_norm = review_matrix / review_norms
        
        # Calculate similarity to positive and negative anchors using vectorized operations
        positive_sim = np.dot(review_norm, pos_embedding)
        negative_sim = np.dot(review_norm, neg_embedding)
        
        # Calculate sentiment scores: scale from -1 to 1
        # Higher positive similarity and lower negative similarity = higher score
        sentiment_scores = positive_sim - negative_sim
        
        # Clip values to ensure they stay in [-1, 1] range (should be already, but just to be safe)
        sentiment_scores = np.clip(sentiment_scores, -1, 1)
        
        logger.info(f"Calculated sentiment scores for {len(review_embeddings)} reviews")
        return sentiment_scores.tolist()
        
    except Exception as e:
        logger.error(f"Error calculating sentiment scores: {str(e)}")
        # Return neutral sentiment scores as fallback
        return [0.0] * len(review_embeddings)
async def process_excel_data(
    client: AzureOpenAI, 
    file_content: bytes,
    filename: str,
    source: str = None,
    skip_insights: bool = False
) -> Dict[str, Any]:
    '''
    Process Excel file data to extract and classify feedback from all sheets.
    
    Args:
        client: Azure OpenAI client
        file_content: Raw Excel file content
        filename: Original filename for logging
        source: Optional source value to override the Source column
        
    Returns:
        Dictionary with key areas, classified feedback, and insight summaries
    '''
    try:
        # Save the content to a temporary file for pandas to read
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp:
            temp.write(file_content)
            temp_path = temp.name
            logger.info(f"Saved Excel file to temporary path: {temp_path}")
        
        # Read all sheets from the Excel file
        logger.info(f"Reading Excel file with multiple sheets: {filename}")
        excel_file = pd.ExcelFile(temp_path)
        sheet_names = excel_file.sheet_names
        logger.info(f"Found {len(sheet_names)} sheets: {sheet_names}")
        
        # Process each sheet and collect results
        all_feedback_items = []
        sheet_dfs = {}
        
        for sheet_name in sheet_names:
            logger.info(f"Processing sheet: {sheet_name}")
            try:
                df = pd.read_excel(temp_path, sheet_name=sheet_name)
                if not df.empty:
                    sheet_dfs[sheet_name] = df
                    logger.info(f"Sheet {sheet_name}: {len(df)} rows, {len(df.columns)} columns")
                    logger.info(f"Original columns: {list(df.columns)}")
                    
                    # Use identify_relevant_columns to find feedback and rating columns
                    columns = await identify_relevant_columns(AZURE_CLIENT, df)
                    feedback_col = columns.get("feedback_col")
                    rating_col = columns.get("rating_col")
                    
                    logger.info(f"Sheet {sheet_name} columns: feedback={feedback_col}, rating={rating_col}")
                    
                    # Skip processing if no feedback column identified
                    if not feedback_col or feedback_col not in df.columns:
                        logger.warning(f"No feedback column identified in sheet {sheet_name}, skipping")
                        continue
                    
                    # Apply standardization FIRST - this is critical
                    df = standardize_dataframe(df, columns, source)
                    logger.info(f"Standardized columns: {list(df.columns)}")
                    
                    # Update column references to use standardized names
                    feedback_col = "Customer Feedback"  # Always use standardized name
                    
                    # Skip 5-star reviews if rating column exists and was kept
                    rating_col = None
                    for col in df.columns:
                        if col not in ["Received", "Name", "Customer Feedback", "Source"] and "rating" in col.lower():
                            rating_col = col
                            break
                    
                    # if rating_col:
                    #     try:
                    #         df[rating_col] = pd.to_numeric(df[rating_col], errors='coerce')
                    #         high_rating_mask = df[rating_col] >= 5
                    #         five_star_count = high_rating_mask.sum()
                    #         logger.info(f"Skipping {five_star_count} high-rated reviews (5-star)")
                    #         df = df[~high_rating_mask]
                    #     except Exception as e:
                    #         logger.warning(f"Could not filter out 5-star reviews, processing all: {str(e)}")
                    
                    # Extract feedback using the standardized dataframe
                    for idx, row in df.iterrows():
                        if pd.notna(row[feedback_col]) and str(row[feedback_col]).strip():
                            # Construct feedback text
                            feedback_text = str(row[feedback_col])
                            
                            # Add rating info if available
                            if rating_col and pd.notna(row[rating_col]):
                                feedback_text = f"[Rating: {row[rating_col]}] {feedback_text}"
                            
                            # Create row dictionary from standardized dataframe
                            row_dict = row.replace([float('nan'), float('inf'), float('-inf')], -1).to_dict()
                            
                            # Add to feedback items
                            all_feedback_items.append({
                                "text": feedback_text,
                                "original_index": idx,
                                "original_row": row_dict,  # Now contains standardized column names
                                "sheet_name": sheet_name
                            })
            except Exception as sheet_error:
                logger.error(f"Error processing sheet {sheet_name}: {str(sheet_error)}")
                # Continue with other sheets even if one fails
        
        # Clean up temporary file
        try:
            os.unlink(temp_path)
            logger.info(f"Removed temporary file: {temp_path}")
        except Exception as e:
            logger.warning(f"Failed to remove temporary file {temp_path}: {str(e)}")
        
        if not all_feedback_items:
            raise ValueError("No valid feedback found in any sheet of the Excel file")
            
        logger.info(f"Extracted {len(all_feedback_items)} feedback items from {len(sheet_dfs)} sheets")
            
        # Sample for key area identification (limit to 100 items to avoid token limits)
        # Sample for key area identification - ensure we use all reviews for small datasets
        if len(all_feedback_items) <= 10:
            sample_feedbacks = [item["text"] for item in all_feedback_items]
        else:
            sample_feedbacks = [item["text"] for item in all_feedback_items[:min(1000, len(all_feedback_items) // 2)]]
        #sample_feedbacks = [item["text"] for item in all_feedback_items[:min(1000,len(all_feedback_items) // 2)]]
        sample_feedback_text = '\n'.join(sample_feedbacks)
        
        # Identify key areas
        key_areas = await identify_key_areas(AZURE_CLIENT, sample_feedback_text)
            
        logger.info(f"Identified {len(key_areas)} key areas")
        
        # Initialize result structure with empty lists - use the area name as key
        classified_feedback = {area.get('area', f"Area {i+1}"): [] for i, area in enumerate(key_areas)}
        
        # Use embeddings for classification
        try:
            logger.info("Using embeddings for feedback classification")
            
            # Create a list of just the area names for embeddings
            problem_texts = [area.get('problem', f"Problem {i+1}") for i, area in enumerate(key_areas)]
            key_area_embeddings = await get_embeddings(AZURE_CLIENT, problem_texts)
            area_names = [area.get('area', f"Area {i+1}") for i, area in enumerate(key_areas)]
            
            # Process feedback in batches for embedding
            feedback_texts = [item["text"] for item in all_feedback_items]
            batch_size = 300
            all_feedback_embeddings = []
            
            for i in range(0, len(feedback_texts), batch_size):
                batch = feedback_texts[i:i + batch_size]
                logger.info(f"Generating embeddings for feedback batch {i//batch_size + 1}/{(len(feedback_texts) + batch_size - 1)//batch_size}")
                batch_embeddings = await get_embeddings(AZURE_CLIENT, batch)
                all_feedback_embeddings.extend(batch_embeddings)
            
            # Classify feedback by type (feature vs issue)
            feedback_types = await classify_feedback_by_type(AZURE_CLIENT, all_feedback_embeddings, feedback_texts)
            
            # Create area type mapping
            area_type_map = {area["area"]: area.get("area_type", "issue") for area in key_areas}
            
            # Quick preliminary classification for summary generation
            key_area_matrix = np.array(key_area_embeddings)
            feedback_matrix = np.array(all_feedback_embeddings)
            
            # Verify dimensions match
            key_area_dim = key_area_matrix.shape[1]
            feedback_dim = feedback_matrix.shape[1]
            logger.info(f"Embedding dimensions - Key areas: {key_area_matrix.shape}, Feedback: {feedback_matrix.shape}")
            
            if key_area_dim != feedback_dim:
                logger.error(f"Embedding dimensions mismatch: key_areas({key_area_dim}) != feedback({feedback_dim})")
                raise ValueError(f"Embedding dimensions don't match: {key_area_dim} vs {feedback_dim}")
            
            # Calculate similarity
            def cosine_similarity_matrix(A, B):
                A_norm = A / np.linalg.norm(A, axis=1, keepdims=True)
                B_norm = B / np.linalg.norm(B, axis=1, keepdims=True)
                return np.dot(A_norm, B_norm.T)
            
            similarity_matrix = cosine_similarity_matrix(feedback_matrix, key_area_matrix)
            
            # Quick preliminary classification for summary generation
            preliminary_classified = {area.get('area', f"Area {i+1}"): [] for i, area in enumerate(key_areas)}
            for i, similarities in enumerate(similarity_matrix):
                best_match = np.argmax(similarities)
                if best_match < len(area_names):
                    area = area_names[best_match]
                    preliminary_classified[area].append(feedback_texts[i])
            
            # Generate insight summaries using the same embeddings
            logger.info("Generating insight summaries from feedback embeddings")
            insight_summary = await generate_summary(
                AZURE_CLIENT, 
                all_feedback_embeddings, 
                feedback_texts,
                preliminary_classified,  # Pass preliminary classification
                key_areas,  # Pass key areas with area_type
                skip_insights=skip_insights
            )
            
            # Calculate sentiment scores using the same embeddings we already generated
            logger.info("Calculating sentiment scores for all feedback items")
            sentiment_scores = await calculate_sentiment_scores(all_feedback_embeddings, AZURE_CLIENT)
            
            # Create a mapping from text to sentiment score
            text_to_sentiment = {feedback_texts[i]: sentiment_scores[i] for i in range(len(feedback_texts))}
            
            # Apply type-based weighting for final classification
            for i in range(similarity_matrix.shape[0]):
                feedback_type = feedback_types[i]
                for j in range(similarity_matrix.shape[1]):
                    area_name = area_names[j]
                    area_type = area_type_map.get(area_name, "issue")
                    
                    # Apply weighting based on type match
                    if feedback_type == area_type:
                        similarity_matrix[i, j] *= 1.2  # Boost similarity for matching types
                    else:
                        similarity_matrix[i, j] *= 0.8  # Reduce similarity for mismatched types
            
            # Classify feedback based on similarity
            # Adjust similarity threshold based on dataset size
            if len(all_feedback_items) <= 5:
                similarity_threshold = 0.65  # Lower threshold for small datasets
            else:
                similarity_threshold = 0.8
            
            for i, similarities in enumerate(similarity_matrix):
                # Get indices where similarity exceeds threshold
                matches = np.where(similarities > similarity_threshold)[0]
                
                # If no matches, use the best match
                if len(matches) == 0:
                    best_match = np.argmax(similarities)
                    matches = [best_match]
                
                # Add feedback to all matching areas
                for match_idx in matches:
                    if match_idx < len(area_names):
                        area = area_names[match_idx]
                        classified_feedback[area].append(feedback_texts[i])
            
            logger.info(f"Successfully classified {len(all_feedback_items)} feedback items using embeddings")
            
        except Exception as e:
            logger.error(f"Error with embedding classification: {str(e)}")
            logger.info("Falling back to direct OpenAI classification")
            
            # Fallback for insight summary if embeddings failed
            insight_summary = {
                "user_loves": "Insights will be generated via completions" if skip_insights else "Unable to determine what users love due to an error in embeddings processing",
                "feature_request": "Insights will be generated via completions" if skip_insights else "Unable to identify feature requests due to an error in embeddings processing",
                "pain_point": "Insights will be generated via completions" if skip_insights else "Unable to identify pain points due to an error in embeddings processing",
                "overall_summary": "Insights will be generated via completions" if skip_insights else "Unable to generate an overall summary due to an error"
            }
            
            # Fall back to chunked classification via OpenAI
            chunk_size = 10
            for i in range(0, len(all_feedback_items), chunk_size):
                logger.info(f"Processing feedback chunk {i//chunk_size + 1}/{(len(all_feedback_items) + chunk_size - 1)//chunk_size}")
                
                chunk = all_feedback_items[i:i + min(chunk_size, len(all_feedback_items) - i)]
                chunk_texts = [item["text"] for item in chunk]
                
                # Create prompt and classify
                areas_text = "\n".join([f"{j+1}. {area.get('area', f'Area {j+1}')} ({area.get('area_type', 'issue')})" for j, area in enumerate(key_areas)])
                feedback_text = "\n".join([f"Feedback {j+1}: {text}" for j, text in enumerate(chunk_texts)])
                
                prompt = f'''
                I need to classify these customer feedback items into the most relevant categories.
                
                CATEGORIES (with their type):
                {areas_text}
                
                FEEDBACK ITEMS:
                {feedback_text}
                
                For each feedback item, tell me which categories it belongs to (multiple categories allowed).
                Prefer assigning feedback to categories of matching type (feature requests to feature categories, issues to issue categories).
                
                Format as JSON:
                {{
                    "classifications": [
                        {{
                            "feedback_index": 0,
                            "category_indices": [0, 2]
                        }},
                        ...
                    ]
                }}
                
                Category indices are 0-based, corresponding to the numbered list above.
                '''
                
                try:
                    response = client.chat.completions.create(
                        model="gpt-4.1-mini",
                        messages=[
                            {"role": "system", "content": "You are a customer feedback classification system."},
                            {"role": "user", "content": prompt}
                        ],
                        temperature=0.1,
                        response_format={"type": "json_object"}
                    )
                    
                    result = json.loads(response.choices[0].message.content)
                    
                    # Process classifications
                    if "classifications" in result:
                        for classification in result["classifications"]:
                            feedback_idx = classification.get("feedback_index", 0)
                            category_indices = classification.get("category_indices", [])
                            
                            if feedback_idx < len(chunk_texts):
                                feedback_text = chunk_texts[feedback_idx]
                                
                                # If no categories matched, assign to first category
                                if not category_indices and key_areas:
                                    first_area = key_areas[0].get('area', 'Area 1')
                                    classified_feedback[first_area].append(feedback_text)
                                else:
                                    # Add to all matching categories
                                    for cat_idx in category_indices:
                                        if 0 <= cat_idx < len(key_areas):
                                            area = key_areas[cat_idx].get('area', f'Area {cat_idx+1}')
                                            classified_feedback[area].append(feedback_text)
                
                except Exception as chunk_error:
                    logger.error(f"Error processing chunk {i//chunk_size + 1}: {str(chunk_error)}")
                    # Assign all feedbacks in this chunk to the first category as fallback
                    if key_areas:
                        first_area = key_areas[0].get('area', 'Area 1')
                        classified_feedback[first_area].extend(chunk_texts)
                
                # Avoid rate limits
                await asyncio.sleep(0.5)
        
        # Log classification stats
        logger.info("Classification results:")
        for area, feedbacks in classified_feedback.items():
            area_type = area_type_map.get(area, "unknown") if 'area_type_map' in locals() else "unknown"
            logger.info(f"  - {area} [{area_type}]: {len(feedbacks)} items")
        
        # Create a mapping from text to item and add sentiment scores
        text_to_item = {item["text"]: item for item in all_feedback_items}
        enriched_feedback = {}
        
        for area, feedbacks in classified_feedback.items():
            enriched_feedback[area] = []
            for feedback_text in feedbacks:
                if feedback_text in text_to_item:
                    # Include original row data (which now has standardized column names)
                    original_item = text_to_item[feedback_text]
                    original_row = original_item.get("original_row", {"text": feedback_text})
                    
                    # Add sentiment score if available
                    if feedback_text in text_to_sentiment:
                        original_row["sentiment_score"] = text_to_sentiment[feedback_text]
                    else:
                        # Use a neutral sentiment as fallback
                        original_row["sentiment_score"] = 0.0
                        
                    enriched_feedback[area].append(original_row)
                else:
                    # Fallback if text not found in mapping - ensure standardized columns
                    enriched_feedback[area].append({
                        "Customer Feedback": feedback_text, 
                        "Received": "N/A",
                        "Name": "N/A",
                        "Source": source if source else "N/A",
                        "sentiment_score": 0.0
                    })
        
        # Replace classified_feedback with enriched version that includes sentiment
        classified_feedback = enriched_feedback

        return {
            "key_areas": key_areas,  # Now includes area_type
            "classified_feedback": classified_feedback,
            "insight_summary": insight_summary
        }
        
    except Exception as e:
        logger.error(f"Error processing Excel data: {str(e)}\n{traceback.format_exc()}")
        # Fall back to raw data processing if Excel processing fails
        try:
            # Convert Excel to CSV format as fallback
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp:
                temp.write(file_content)
                temp_path = temp.name
            
            # Read the first sheet as fallback
            df = pd.read_excel(temp_path, sheet_name=0)
            csv_data = df.to_csv(index=False)
            
            # Clean up temp file
            os.unlink(temp_path)
            
            logger.info("Falling back to CSV processing for Excel data")
            return await process_csv_data(AZURE_CLIENT, csv_data, source, skip_insights=skip_insights)
        except Exception as fallback_error:
            logger.error(f"Excel fallback also failed: {str(fallback_error)}")
            # Final fallback - raw data
            logger.info("Falling back to raw data processing")
            excel_text = str(file_content)  # Very crude fallback
            return await analyze_raw_data_chunks(AZURE_CLIENT, excel_text)
async def identify_key_areas(client: AzureOpenAI, data_sample: str, max_areas: int = 15) -> List[Dict[str, str]]:
    '''
    Identifies key problem areas from customer feedback using OpenAI.
    Now also classifies each area as either 'feature' or 'issue'.
    
    Args:
        client: Azure OpenAI client
        data_sample: Sample of feedback data to analyze
        max_areas: Maximum number of key areas to identify
        
    Returns:
        List of dictionaries with 'area', 'problem', and 'area_type' keys
    '''
    try:
        logger.info(f"Beginning key problem area identification using {len(data_sample)} characters of sample data")
        # For very small datasets, use a simpler approach
        if len(data_sample) < 500 or data_sample.count('\n') < 3:
            logger.info("Small dataset detected, using simplified key area identification")
            max_areas = min(max_areas, 3)  # Limit areas for small datasets
        # Prepare the prompt for OpenAI
        prompt = f'''
        # Customer Feedback Analysis Task

        You are a senior product insights analyst specializing in customer feedback pattern recognition. 
        Your expertise lies in identifying underlying themes and problems from diverse customer feedback.
        You are also expert in extracting app/ui specific problems which ruin the customer experience.

        ## Your Objective
        Analyze the provided customer feedback data and identify the most significant problem areas that customers are experiencing, with a focus on user experience, app-related issues while still capturing important physical/in-store issues when relevant. These insights will directly inform product development priorities.

        IMPORTANT: You must identify BOTH:
        1. Feature Request Areas - What customers want added/improved (mark as "feature")
        2. Issue/Problem Areas - What's currently broken or problematic (mark as "issue")

        Aim for a balanced mix of both types (at least 40% feature areas and 40% issue areas).

        ## Customer Feedback Data Sample
        ```
        {data_sample}
        ```
        ## Analysis Instructions:
        1. Carefully read and understand all customer feedback in the sample

        2. Identify the top {max_areas} most significant problem areas based on:
        - Frequency (how many customers mention this issue)
        - Severity (how impactful the problem seems to be)
        - Specificity (clear, actionable problem statements)
        - Business impact (issues affecting core product value)

        3. For each identified problem area:
        - Create a concise, high-level area (2-3 words) reflecting the functional area
        - Write a specific problem statement from the customer's perspective (1 sentence)
        - Classify as either "feature" (for feature requests/enhancements) or "issue" (for bugs/problems)
        - Ensure the problem is concrete enough to be actionable
        - Capture the essence of what customers are struggling with or requesting
        - Note that the same high-level area (e.g., "App Performance") can have multiple specific problems

        4. You may:
        - Select directly from the reference categories if they match well
        - Adapt a reference category with more specific wording
        - Create entirely new categories if the existing ones don't capture the feedback
        - Combine similar issues into a single coherent problem area
        - Group related issues under the same area with different problem statements

        ## Response Format Requirements
        You must format your response as a JSON array with each problem area having 'area', 'problem', and 'area_type' keys:

        ```json
        [
            {{"area": "Performance", "problem": "App frequently crashes when switching between multiple workout screens", "area_type": "issue"}},
            {{"area": "Mobile Features", "problem": "Users want dark mode and offline capabilities in the mobile app", "area_type": "feature"}},
            {{"area": "User Interface", "problem": "App navigation requires too many taps to access core tracking features", "area_type": "issue"}},
            {{"area": "Advanced Analytics", "problem": "Users need more detailed reporting and data visualization options", "area_type": "feature"}},
            {{"area": "Data Management", "problem": "App loses workout history when syncing with cloud services", "area_type": "issue"}},
            {{"area": "Integration Options", "problem": "Customers want integration with third-party fitness and health apps", "area_type": "feature"}},
            {{"area": "Connectivity", "problem": "App fails to maintain Bluetooth connection with heart rate monitors", "area_type": "issue"}},
            {{"area": "Customization", "problem": "Users request ability to create custom workout plans and routines", "area_type": "feature"}},
            {{"area": "Tracking Accuracy", "problem": "App calorie calculations show inconsistent results compared to similar services", "area_type": "issue"}},
            {{"area": "Social Features", "problem": "Users want to share achievements and compete with friends", "area_type": "feature"}},
            {{"area": "Battery Life", "problem": "App continues to drain battery significantly when running in background", "area_type": "issue"}},
            {{"area": "Customer Support", "problem": "App troubleshooting guides don't address common synchronization problems", "area_type": "issue"}},
            {{"area": "Personalization", "problem": "Users want AI-powered workout recommendations based on their progress", "area_type": "feature"}},
            {{"area": "App Stability", "problem": "Application crashes during data synchronization with wearables", "area_type": "issue"}},
            {{"area": "Export Options", "problem": "Users need ability to export their data in multiple formats", "area_type": "feature"}}
        ]
        ```

        Ensure you have a good balance of "feature" and "issue" types. Think carefully about naming each area and classifying correctly.
        '''
        
        # Call OpenAI API
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            messages=[
                {"role": "system", "content": "You are a product analytics specialist who identifies key problem areas from customer feedback and classifies them as features or issues."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
            response_format={"type": "json_object"}
        )
        
        # Extract and parse the response
        result_text = response.choices[0].message.content
        try:
            result = json.loads(result_text)
            
            # Handle different possible response formats
            if isinstance(result, dict):
                # Check for common response structures
                if "response" in result:
                    key_areas = result["response"]
                    logger.info("Found key areas in 'response' field")
                elif "areas" in result:
                    key_areas = result["areas"]
                    logger.info("Found key areas in 'areas' field")
                elif "key_areas" in result:
                    key_areas = result["key_areas"]
                    logger.info("Found key areas in 'key_areas' field")
                # If no recognized fields but has area/problem structure, use the whole object
                elif result.get("area") and result.get("problem"):
                    key_areas = [result]
                    logger.info("Found single key area object")
                else:
                    # Try to extract any list with area/problem structure
                    for key, value in result.items():
                        if isinstance(value, list) and len(value) > 0 and isinstance(value[0], dict):
                            if "area" in value[0] and "problem" in value[0]:
                                key_areas = value
                                logger.info(f"Found key areas in '{key}' field")
                                break
                    else:
                        logger.warning(f"Unrecognized dictionary structure in OpenAI response, using entire object")
                        key_areas = [result]
            elif isinstance(result, list):
                # Check if list items have area/problem structure
                if result and isinstance(result[0], dict) and "area" in result[0] and "problem" in result[0]:
                    key_areas = result
                    logger.info("Found key areas in list format")
                else:
                    logger.warning("List items in response don't have expected structure")
                    key_areas = []
            else:
                logger.warning(f"Response is neither dict nor list: {type(result)}")
                key_areas = []
                
            if not key_areas:
                logger.warning(f"Could not extract key areas from response: {result_text[:500]}...")
                key_areas = []
                
        except json.JSONDecodeError:
            logger.error(f"Failed to parse response as JSON: {result_text[:500]}...")
            key_areas = []
            
        # Log the raw response if we had issues
        if not key_areas:
            logger.warning(f"Raw response from OpenAI: {result_text[:1000]}...")
            # Fallback to common key areas with area_type
            fallback_areas = [
                {"area": "Feature Requests", "problem": "Users want new features and enhancements", "area_type": "feature"},
                {"area": "Performance Issues", "problem": "App performance and stability problems", "area_type": "issue"},
                {"area": "User Experience", "problem": "Interface and usability improvements needed", "area_type": "feature"},
                {"area": "Technical Problems", "problem": "Various technical issues and bugs", "area_type": "issue"},
                {"area": "Mobile Features", "problem": "Mobile app feature enhancements requested", "area_type": "feature"},
                {"area": "Integration Issues", "problem": "Problems with third-party integrations", "area_type": "issue"}
            ]
            key_areas = fallback_areas[:max_areas]
            logger.info(f"Using {len(key_areas)} fallback key areas")
            
        # Limit to max_areas and ensure consistent structure
        key_areas = key_areas[:max_areas]
        
        # Ensure each area has 'area', 'problem', and 'area_type' fields
        for area in key_areas:
            if not isinstance(area, dict):
                logger.warning(f"Key area is not a dictionary: {area}")
                continue
            if "area" not in area:
                area["area"] = "Unknown Area"
            if "problem" not in area:
                area["problem"] = "Unspecified problem"
            if "area_type" not in area:
                # Try to infer area_type from the problem description
                problem_lower = area.get("problem", "").lower()
                if any(word in problem_lower for word in ["want", "need", "request", "add", "feature", "enhance", "improve"]):
                    area["area_type"] = "feature"
                else:
                    area["area_type"] = "issue"
        
        # Ensure we have a balanced mix
        feature_count = sum(1 for area in key_areas if area.get("area_type") == "feature")
        issue_count = sum(1 for area in key_areas if area.get("area_type") == "issue")
        
        # Log the identified key areas in detail
        logger.info(f"IDENTIFIED {len(key_areas)} KEY PROBLEM AREAS (Features: {feature_count}, Issues: {issue_count}):")
        for i, area in enumerate(key_areas):
            logger.info(f"  {i+1}. [{area.get('area_type', 'unknown').upper()}] {area['area']}: {area['problem']}")
        
        return key_areas
        
    except Exception as e:
        logger.error(f"Error identifying key areas: {str(e)}")
        # Return a minimal structure to continue processing with area_type
        return [
            {"area": "General Issues", "problem": "Various customer problems and feedback", "area_type": "issue"},
            {"area": "Feature Requests", "problem": "Customer feature requests and enhancements", "area_type": "feature"}
        ]
async def get_embeddings(client: AzureOpenAI, texts: List[str], slow_mode: bool = False) -> List[List[float]]:
    '''
    Generate embeddings for a list of texts.
    
    Args:
        client: Azure OpenAI client
        texts: List of text strings to generate embeddings for
        slow_mode: If True, use the original sequential processing
        
    Returns:
        List of embedding vectors
    '''
    try:
        # Process in batches to avoid request size limits
        batch_size = 100 if slow_mode else 1000
        all_embeddings = []
        
        logger.info(f"Generating embeddings for {len(texts)} texts using model {EMBEDDING_MODEL}")
        
        if slow_mode:
            # Original sequential implementation
            for i in range(0, len(texts), batch_size):
                batch_texts = texts[i:i + batch_size]
                
                try:
                    response = AZURE_CLIENT.embeddings.create(
                        model=EMBEDDING_MODEL,  # Use specific model name
                        input=batch_texts
                    )
                    
                    # Extract embeddings from response
                    batch_embeddings = [item.embedding for item in response.data]
                    all_embeddings.extend(batch_embeddings)
                    
                    logger.info(f"Successfully generated {len(batch_embeddings)} embeddings (batch {i//batch_size + 1})")
                    
                    # Check embedding dimension for debugging
                    if batch_embeddings:
                        dim = len(batch_embeddings[0])
                        logger.info(f"Embedding dimension: {dim}")
                    
                    # Avoid rate limits
                    if i + batch_size < len(texts):
                        await asyncio.sleep(0.5)
                        
                except Exception as batch_error:
                    logger.error(f"Error generating embeddings for batch {i//batch_size + 1}: {str(batch_error)}")
                    # For this batch, generate random embeddings as fallback
                    # Use fixed dimension of 3072 to match the model
                    dim = 1536     # text-embedding-3-large dimension
                    logger.warning(f"Using fallback random embeddings for batch {i//batch_size + 1} with dimension {dim}")
                    for _ in range(len(batch_texts)):
                        random_embedding = list(np.random.normal(0, 0.1, dim))
                        all_embeddings.append(random_embedding)
        else:
            # Optimized parallel implementation
            async def process_batch(batch_idx, batch_texts):
                try:
                    response = AZURE_CLIENT.embeddings.create(
                        model=EMBEDDING_MODEL,
                        input=batch_texts
                    )
                    batch_embeddings = [item.embedding for item in response.data]
                    logger.info(f"Successfully generated {len(batch_embeddings)} embeddings (batch {batch_idx+1})")
                    return batch_embeddings
                except Exception as batch_error:
                    logger.error(f"Error generating embeddings for batch {batch_idx+1}: {str(batch_error)}")
                    dim = 1536   # text-embedding-3-large dimension
                    return [list(np.random.normal(0, 0.1, dim)) for _ in range(len(batch_texts))]
            
            # Split texts into batches
            batches = [(i//batch_size, texts[i:i + batch_size]) 
                       for i in range(0, len(texts), batch_size)]
            
            # Process batches in parallel with rate limiting
            semaphore = asyncio.Semaphore(5)  # Limit to 5 concurrent API calls
            
            async def process_with_semaphore(batch_idx, batch_texts):
                async with semaphore:
                    return await process_batch(batch_idx, batch_texts)
            
            # Gather results from all batches
            batch_results = await asyncio.gather(
                *[process_with_semaphore(idx, batch) for idx, batch in batches]
            )
            
            # Flatten results
            for batch_emb in batch_results:
                all_embeddings.extend(batch_emb)
        
        return all_embeddings
    
    except Exception as e:
        logger.error(f"Error generating embeddings: {str(e)}")
        # Generate random embeddings as fallback
        logger.warning("Using fallback random embeddings for all texts")
        dim = 1536   # text-embedding-3-large dimension
        return [list(np.random.normal(0, 0.1, dim)) for _ in range(len(texts))]
def cosine_similarity_batch(A: np.ndarray, B: np.ndarray, slow_mode: bool = False) -> np.ndarray:
    '''
    Calculate cosine similarity between two sets of vectors efficiently.
    
    Args:
        A: Matrix of shape (n_samples_A, n_features)
        B: Matrix of shape (n_samples_B, n_features)
        slow_mode: If True, use a slower for-loop implementation
        
    Returns:
        Similarity matrix of shape (n_samples_A, n_samples_B)
    '''
    if slow_mode:
        # Slower iterative implementation (for compatibility)
        n_samples_A, _ = A.shape
        n_samples_B, _ = B.shape
        similarity_matrix = np.zeros((n_samples_A, n_samples_B))
        
        for i in range(n_samples_A):
            for j in range(n_samples_B):
                dot_product = np.dot(A[i], B[j])
                norm_A = np.linalg.norm(A[i])
                norm_B = np.linalg.norm(B[j])
                
                if norm_A == 0 or norm_B == 0:
                    similarity_matrix[i, j] = 0
                else:
                    similarity_matrix[i, j] = dot_product / (norm_A * norm_B)
        
        return similarity_matrix
    else:
        # Fast vectorized implementation
        # Normalize the matrices
        A_norm = A / np.linalg.norm(A, axis=1, keepdims=True)
        B_norm = B / np.linalg.norm(B, axis=1, keepdims=True)
        # Calculate similarity matrix in one operation
        return np.dot(A_norm, B_norm.T)
async def cosine_similarity(vec1: List[float], vec2: List[float]) -> float:
    '''Calculate cosine similarity between two vectors.'''
    # Convert to numpy arrays for vectorized operations
    vec1_np = np.array(vec1)
    vec2_np = np.array(vec2)
    
    # Calculate dot product and norms
    dot_product = np.dot(vec1_np, vec2_np)
    norm_a = np.linalg.norm(vec1_np)
    norm_b = np.linalg.norm(vec2_np)
    
    if norm_a == 0 or norm_b == 0:
        return 0
    
    return dot_product / (norm_a * norm_b)

async def identify_relevant_columns(client: AzureOpenAI, df: pd.DataFrame) -> Dict[str, str]:
    '''
    Uses OpenAI to identify which columns contain ratings, feedback text, received date, name, and source.
    Now with extensive fallback logic and validation to ensure proper column identification.
    
    Args:
        client: Azure OpenAI client
        df: Pandas DataFrame containing the data
        
    Returns:
        Dictionary with column mappings
    '''
    try:
        logger.info(f"Beginning column detection on DataFrame with shape: {df.shape}")
        logger.info(f"Available columns: {list(df.columns)}")
        
        # Initialize result
        columns = {
            "feedback_col": None, 
            "rating_col": None,
            "received_col": None,
            "name_col": None,
            "source_col": None
        }
        
        # Clean column names first
        df_columns = [col.strip() if isinstance(col, str) else str(col) for col in df.columns]
        
        # Get a sample of the data
        sample_size = min(5, len(df))
        if sample_size == 0:
            logger.error("Empty DataFrame provided")
            return columns
            
        sample = df.sample(sample_size).to_csv(index=False)
        columns_list = list(df.columns)
        columns_with_indices = [f"{i}: {col}" for i, col in enumerate(columns_list)]
        columns_text = "\n".join(columns_with_indices)
        
        # Try OpenAI first
        try:
            prompt = f'''
            You are a specialized data analyst focusing on customer feedback analysis. Your task is to precisely identify the most relevant columns in this customer feedback dataset.

            DATASET SAMPLE:
            ```
            {sample}
            ```

            AVAILABLE COLUMNS (with indices):
            ```
            {columns_text}
            ```

            TASK:
            Carefully examine the dataset and identify these types of columns:

            1. CUSTOMER FEEDBACK COLUMN:
               - Contains textual customer opinions, comments, reviews, or feedback
               - Common names: "Review", "review_text", "Customer Feedback", "feedback", "comment", "text", "description", "comments", "Reviews"
               - Usually the column with the LONGEST text content
               - Usually includes sentences, paragraphs, or detailed comments
               - May contain terms like "liked", "disliked", "issue with", "problem", etc.
               - IMPORTANT: If you see a column named "Review" or "Reviews", it's almost certainly the feedback column
            
            2. RATING COLUMN:
               - Contains numerical scores (1-5, 1-10) or textual ratings ("Excellent", "Poor")
               - Common names: "Rating", "star_rating", "rating", "score", "stars", "satisfaction"
               - May be presented as numbers or categorical values
               - IMPORTANT: If you see a column named "Rating" with numeric values, it's almost certainly the rating column

            3. RECEIVED COLUMN:
               - Contains date or timestamp information when the feedback was received
               - Common names: "date", "received", "submitted", "timestamp", "Received", "created_at"
               - May be formatted as date, datetime, or timestamp

            4. NAME COLUMN:
               - Contains the name of the reviewer or customer
               - Common names: "Name", "user", "customer", "reviewer", "username", "customer_name"
               - Usually contains full names, first names, or usernames
               - IMPORTANT: If you see a column named "Name" with person names, it's almost certainly the name column
               - This should have SHORTER text than the feedback column

            5. SOURCE COLUMN:
               - Contains information about where the feedback came from
               - Common names: "source", "platform", "channel", "website", "Source", "origin"
               - May contain values like "Google", "Amazon", "Website", "App", etc.

            CRITICAL RULES:
            - The feedback column should have the LONGEST average text length
            - If a column is named "Review" or "Reviews", it's very likely the feedback column
            - If a column is named "Rating" with numbers 1-5, it's very likely the rating column
            - If a column is named "Name" with short text entries, it's very likely the name column
            - Do NOT confuse "Review" (feedback) with "Reviewer" (name)
            - Do NOT select a name/ID column as the feedback column

            For each column you identify, consider:
            - Column content and data type
            - Column name relevance (exact matches like "Review", "Rating", "Name" are highly likely)
            - Amount of information in the column (feedback should be longest)
            - Uniqueness of values

            IMPORTANT INSTRUCTIONS:
            - You MUST select from the provided column indices (0 to {len(columns_list)-1})
            - You must specify both the index and exact column name
            - If a certain type of column doesn't exist, set the value to null
            - Pay special attention to simple column names like "Review", "Rating", "Name"

            RESPONSE FORMAT:
            You must respond with column indices and names in this exact JSON format:
            {{
                "feedback_col": "3: [Exact Column Name]",
                "rating_col": "1: [Exact Column Name]",
                "received_col": "0: [Exact Column Name]",
                "name_col": "2: [Exact Column Name]",
                "source_col": "4: [Exact Column Name]"
            }}
            '''
            
            response = client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[
                    {"role": "system", "content": "You are a data analysis assistant that identifies column types in customer feedback datasets. You recognize both traditional column names and API-generated column names."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                response_format={"type": "json_object"}
            )
            
            result_text = response.choices[0].message.content
            result = json.loads(result_text)
            
            # Extract column information from OpenAI response
            for col_type in columns.keys():
                if result.get(col_type):
                    # Parse the response format "index: Column Name"
                    try:
                        col_parts = result[col_type].split(":", 1)
                        if len(col_parts) == 2:
                            idx = int(col_parts[0].strip())
                            col_name = col_parts[1].strip()
                            
                            # Verify the column exists
                            if 0 <= idx < len(columns_list) and columns_list[idx] == col_name:
                                columns[col_type] = col_name
                            else:
                                # Try fuzzy matching
                                best_match, score = process.extractOne(col_name, columns_list)
                                if score > 80:  # Good match threshold
                                    columns[col_type] = best_match
                                else:
                                    logger.warning(f"Column '{col_name}' not found and no good match")
                    except (ValueError, IndexError) as e:
                        logger.warning(f"Error parsing column specification '{result[col_type]}': {e}")
                        
        except Exception as e:
            logger.error(f"OpenAI column identification failed: {str(e)}")
        
        # Apply comprehensive fallback strategies
        
        # Feedback column fallback - most critical
        if not columns["feedback_col"] or columns["feedback_col"] not in df.columns:
            logger.info("Applying fallback logic for feedback column")
            
            # Extended keywords list with priority order
            exact_keywords = ["review", "reviews", "customer feedback", "feedback", "comment", "comments", "review_text", "review text"]
            partial_keywords = ["review", "feedback", "comment", "text", "description", "message", "opinion", "testimonial", "remarks"]
            
            # First try exact match (case-insensitive)
            for col in df.columns:
                col_lower = str(col).lower().strip()
                for keyword in exact_keywords:
                    if col_lower == keyword:
                        logger.info(f"Selected '{col}' as feedback column based on exact keyword match")
                        columns["feedback_col"] = col
                        break
                if columns["feedback_col"]:
                    break
            
            # If no exact match, try partial match with length validation
            if not columns["feedback_col"]:
                potential_columns = []
                for col in df.columns:
                    col_lower = str(col).lower()
                    for keyword in partial_keywords:
                        if keyword in col_lower and df[col].dtype == 'object':
                            # Calculate average text length
                            non_null_values = df[col].dropna().astype(str)
                            if len(non_null_values) > 0:
                                avg_length = non_null_values.str.len().mean()
                                if avg_length >= 10:  # Lower threshold to catch shorter reviews
                                    potential_columns.append((col, avg_length))
                                    break
                
                # Sort by average length and select the longest
                if potential_columns:
                    potential_columns.sort(key=lambda x: x[1], reverse=True)
                    columns["feedback_col"] = potential_columns[0][0]
                    logger.info(f"Selected '{columns['feedback_col']}' as feedback column based on keyword match and text length ({potential_columns[0][1]:.1f})")
            
            # If still not found, use column with longest average text length
            if not columns["feedback_col"]:
                logger.warning("No feedback column identified via keywords - using text length heuristics")
                
                text_columns = []
                for col in df.columns:
                    if df[col].dtype == 'object':
                        non_null_values = df[col].dropna().astype(str)
                        if len(non_null_values) > len(df) * 0.1:  # Lowered from 0.3 to 0.1 (10% non-null)
                            avg_length = non_null_values.str.len().mean()
                            
                            # Skip if average length is too short (likely names or IDs)
                            if avg_length < 10:  # Lowered from 20 to accommodate shorter reviews
                                continue
                                
                            # Skip if all values are very similar length (likely IDs)
                            length_std = non_null_values.str.len().std()
                            if length_std < 5 and avg_length < 50:
                                continue
                            
                            # Check for diversity in text (not all same values)
                            unique_ratio = non_null_values.nunique() / len(non_null_values)
                            if unique_ratio < 0.1:  # Less than 10% unique values
                                continue
                            
                            text_columns.append((col, avg_length))
                
                # Select column with longest average text
                if text_columns:
                    text_columns.sort(key=lambda x: x[1], reverse=True)
                    columns["feedback_col"] = text_columns[0][0]
                    logger.info(f"Selected '{columns['feedback_col']}' as feedback column based on longest text content (avg: {text_columns[0][1]:.1f})")
                else:
                    # Absolute last resort - use first non-numeric column
                    for col in df.columns:
                        if df[col].dtype == 'object':
                            columns["feedback_col"] = col
                            logger.warning(f"Using '{col}' as feedback column as last resort")
                            break
        
        # Rating column fallback
        if not columns["rating_col"] or columns["rating_col"] not in df.columns:
            logger.info("Applying fallback logic for rating column")
            
            # Extended keywords with priority order
            exact_keywords = ["rating", "star_rating", "stars", "score", "star rating"]
            partial_keywords = ["rating", "score", "stars", "grade", "rank", "satisfaction", "rate"]
            
            # First try exact match (case-insensitive)
            for col in df.columns:
                if col == columns["feedback_col"]:  # Skip if already assigned
                    continue
                    
                col_lower = str(col).lower().strip()
                for keyword in exact_keywords:
                    if col_lower == keyword:
                        logger.info(f"Selected '{col}' as rating column based on exact keyword match")
                        columns["rating_col"] = col
                        break
                if columns["rating_col"]:
                    break
            
            # Try partial match with numeric validation
            if not columns["rating_col"]:
                for col in df.columns:
                    if col == columns["feedback_col"]:
                        continue
                        
                    col_lower = str(col).lower()
                    # Skip if it contains "review" (unless it's exactly "review rating")
                    if 'review' in col_lower and 'rating' not in col_lower:
                        continue
                        
                    for keyword in partial_keywords:
                        if keyword in col_lower:
                            # Check if it's numeric or can be converted
                            if pd.api.types.is_numeric_dtype(df[col]):
                                non_null_values = df[col].dropna()
                                if len(non_null_values) > 0:
                                    min_val = non_null_values.min()
                                    max_val = non_null_values.max()
                                    unique_count = non_null_values.nunique()
                                    
                                    # Check for typical rating patterns
                                    if ((min_val >= 0 and max_val <= 5 and unique_count <= 6) or
                                        (min_val >= 1 and max_val <= 10 and unique_count <= 10) or
                                        (min_val >= 0 and max_val <= 100 and unique_count <= 20)):
                                        columns["rating_col"] = col
                                        logger.info(f"Selected '{col}' as rating column based on keyword and numeric range")
                                        break
                    if columns["rating_col"]:
                        break
            
            # Look for any numeric column with rating-like values
            if not columns["rating_col"]:
                for col in df.columns:
                    if col == columns["feedback_col"]:
                        continue
                        
                    if pd.api.types.is_numeric_dtype(df[col]):
                        non_null_values = df[col].dropna()
                        if len(non_null_values) > len(df) * 0.3:  # At least 30% non-null
                            min_val = non_null_values.min()
                            max_val = non_null_values.max()
                            
                            # Check for common rating ranges
                            if (0 <= min_val <= 1 and 4 <= max_val <= 5) or \
                               (1 <= min_val <= 2 and 9 <= max_val <= 10):
                                columns["rating_col"] = col
                                logger.info(f"Selected '{col}' as rating column based on numeric range [{min_val}, {max_val}]")
                                break
        
        # Name column fallback
        if not columns["name_col"] or columns["name_col"] not in df.columns:
            logger.info("Applying fallback logic for name column")
            
            exact_keywords = ["name", "customer", "reviewer", "user", "customer name", "username"]
            partial_keywords = ["name", "customer", "reviewer", "user", "author", "by", "client", "person"]
            
            # First try exact match
            for col in df.columns:
                if col in [columns["feedback_col"], columns["rating_col"]]:
                    continue
                    
                col_lower = str(col).lower().strip()
                for keyword in exact_keywords:
                    if col_lower == keyword:
                        logger.info(f"Selected '{col}' as name column based on exact keyword match")
                        columns["name_col"] = col
                        break
                if columns["name_col"]:
                    break
            
            # Try partial match with length validation
            if not columns["name_col"]:
                for col in df.columns:
                    if col in [columns["feedback_col"], columns["rating_col"]]:
                        continue
                        
                    col_lower = str(col).lower()
                    # Skip if it's a review column (unless it's reviewer)
                    if 'review' in col_lower and 'reviewer' not in col_lower:
                        continue
                        
                    for keyword in partial_keywords:
                        if keyword in col_lower and df[col].dtype == 'object':
                            # Check average length (names should be short)
                            avg_length = df[col].dropna().astype(str).str.len().mean()
                            if avg_length < 50:  # Names are typically short
                                columns["name_col"] = col
                                logger.info(f"Selected '{col}' as name column based on keyword and text length")
                                break
                    if columns["name_col"]:
                        break
        
        # Received column fallback
        if not columns["received_col"] or columns["received_col"] not in df.columns:
            keywords = ["date", "time", "received", "created", "timestamp", "submitted", "posted", "when"]
            
            for col in df.columns:
                if col in [columns["feedback_col"], columns["rating_col"], columns["name_col"]]:
                    continue
                    
                col_lower = str(col).lower()
                
                # Check column name
                for keyword in keywords:
                    if keyword in col_lower:
                        columns["received_col"] = col
                        logger.info(f"Selected '{col}' as received column based on keyword match")
                        break
                
                if columns["received_col"]:
                    break
                
                # Check if it's a datetime column
                if df[col].dtype == 'datetime64[ns]':
                    columns["received_col"] = col
                    logger.info(f"Selected '{col}' as received column based on datetime type")
                    break
        
        # Source column fallback
        if not columns["source_col"] or columns["source_col"] not in df.columns:
            keywords = ["source", "platform", "channel", "origin", "website", "site", "from", "via"]
            
            for col in df.columns:
                if col in [columns["feedback_col"], columns["rating_col"], columns["name_col"], columns["received_col"]]:
                    continue
                    
                col_lower = str(col).lower()
                for keyword in keywords:
                    if keyword in col_lower:
                        columns["source_col"] = col
                        logger.info(f"Selected '{col}' as source column based on keyword match")
                        break
                if columns["source_col"]:
                    break
        
        # Final validation before returning
        logger.info("Performing final validation of column assignments")
        
        # Ensure feedback column exists and has meaningful content
        if columns["feedback_col"]:
            if columns["feedback_col"] not in df.columns:
                logger.error(f"Feedback column '{columns['feedback_col']}' not in DataFrame columns")
                columns["feedback_col"] = None
            else:
                # Validate it has actual feedback content
                sample_values = df[columns["feedback_col"]].dropna().astype(str).head(3).tolist()
                avg_length = df[columns["feedback_col"]].dropna().astype(str).str.len().mean()
                
                if avg_length < 5:  # Lowered from 10 to be more permissive
                    logger.warning(f"Feedback column has very short content (avg: {avg_length:.1f}), looking for better option")
                    # Try to find a better column
                    better_found = False
                    for col in df.columns:
                        if df[col].dtype == 'object' and col != columns["feedback_col"]:
                            other_avg = df[col].dropna().astype(str).str.len().mean()
                            if other_avg > avg_length * 2:
                                logger.info(f"Found better feedback column: '{col}' (avg length: {other_avg:.1f})")
                                columns["feedback_col"] = col
                                better_found = True
                                break
        
        # Ensure no column is assigned to multiple types
        assigned_cols = {}
        priority = {"feedback_col": 0, "rating_col": 1, "name_col": 2, "received_col": 3, "source_col": 4}
        
        for col_type, col_name in list(columns.items()):
            if col_name and col_name in df.columns:
                if col_name in assigned_cols:
                    # Keep the higher priority assignment
                    if priority.get(col_type, 5) < priority.get(assigned_cols[col_name], 5):
                        columns[assigned_cols[col_name]] = None
                        assigned_cols[col_name] = col_type
                    else:
                        columns[col_type] = None
                else:
                    assigned_cols[col_name] = col_type
            elif col_name and col_name not in df.columns:
                logger.warning(f"Column '{col_name}' assigned to {col_type} but not found in DataFrame")
                columns[col_type] = None
        
        # Log final results
        logger.info(f"FINAL COLUMN ASSIGNMENT:")
        for col_type, col_name in columns.items():
            if col_name:
                logger.info(f"  {col_type}: '{col_name}'")
            else:
                logger.info(f"  {col_type}: Not identified")
        
        # Ensure at least feedback column is identified
        if not columns["feedback_col"]:
            # Last resort: use the first text column
            for col in df.columns:
                if df[col].dtype == 'object':
                    columns["feedback_col"] = col
                    logger.warning(f"Using '{col}' as feedback column as absolute last resort")
                    break
        
        return columns
        
    except Exception as e:
        logger.error(f"Critical error in identify_relevant_columns: {str(e)}")
        # Return safe defaults
        columns = {
            "feedback_col": None, 
            "rating_col": None,
            "received_col": None,
            "name_col": None,
            "source_col": None
        }
        
        # Try to at least identify a text column as feedback
        try:
            for col in df.columns:
                if df[col].dtype == 'object':
                    columns["feedback_col"] = col
                    break
        except:
            pass
            
        return columns
async def analyze_raw_data_chunks(
    client: AzureOpenAI, 
    raw_data: str,
    chunk_size: int = 8000,
    overlap: int = 500,
    slow_mode: bool = False
) -> Dict[str, Any]:
    '''
    Analyzes raw unstructured data by chunking and processing with OpenAI.
    
    Args:
        client: Azure OpenAI client
        raw_data: Raw text data
        chunk_size: Size of chunks to process
        overlap: Overlap between chunks
        slow_mode: If True, use the original sequential processing
        
    Returns:
        Dictionary with key areas, classified feedback, and insight summaries
    '''
    try:
        # Split data into chunks
        chunks = []
        for i in range(0, len(raw_data), chunk_size - overlap):
            chunk = raw_data[i:i + chunk_size]
            chunks.append(chunk)
        
        logger.info(f"Split raw data into {len(chunks)} chunks for processing")
        
        # STAGE 1: Extract feedback items and preliminary areas from chunks
        all_feedbacks = []
        chunk_areas = []
        
        for i, chunk in enumerate(chunks):
            logger.info(f"Processing chunk {i+1}/{len(chunks)}")
            
            # Extract feedback from chunk
            prompt = f'''
            Analyze this raw customer feedback data:
            
            {chunk}
            
            1. Identify distinct customer feedback items/comments.
            2. For each feedback item, extract the exact text.
            3. Ignore any 5-star or extremely positive reviews.
            
            Format as a JSON array of feedback items.
            Example: ["I can't generate the reports I need", "The mobile app lacks key features"]
            '''
            
            try:
                response = client.chat.completions.create(
                    model="gpt-4.1-mini",
                    messages=[
                        {"role": "system", "content": "You are a customer feedback analyzer."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.2,
                    response_format={"type": "json_object"}
                )
                
                result_text = response.choices[0].message.content
                result = json.loads(result_text)
                
                # Extract feedback items
                if isinstance(result, dict) and "feedback" in result:
                    feedbacks = result["feedback"]
                elif isinstance(result, list):
                    feedbacks = result
                else:
                    feedbacks = []
                
                # Store feedbacks with original index and chunk info
                for j, feedback in enumerate(feedbacks):
                    all_feedbacks.append({
                        "text": feedback,
                        "original_index": len(all_feedbacks),
                        "chunk_index": i
                    })
                
                # Identify preliminary areas in this chunk
                if feedbacks:
                    sample_feedback = '\n'.join(feedbacks[:50])
                    areas = await identify_key_areas(AZURE_CLIENT, sample_feedback, max_areas=15)
                    chunk_areas.extend(areas)
                
            except Exception as chunk_error:
                logger.error(f"Error processing chunk {i+1}: {str(chunk_error)}")
                continue
            
            # Avoid rate limits
            await asyncio.sleep(1)
        
        # STAGE 2: Consolidate problem areas across all chunks
        if chunk_areas:
            # Use OpenAI to consolidate
            areas_text = json.dumps(chunk_areas, indent=2)
            prompt = f'''
            I have identified multiple potential problem areas from customer feedback chunks. 
            Please consolidate these into 5-8 main categories, ensuring a balanced mix of feature requests and issues.
            
            {areas_text}
            
            For each category:
            1. Provide a short, consistent name (2-3 words)
            2. Provide a specific problem statement from the customer's perspective
            3. Classify as either "feature" (for feature requests) or "issue" (for problems/bugs)
            
            Ensure at least 40% are feature categories and 40% are issue categories.
            
            Format your response as a JSON array of objects with 'area', 'problem', and 'area_type' keys:
            [
                {{"area": "User Interface", "problem": "The app navigation requires too many taps to access core features", "area_type": "issue"}},
                {{"area": "Advanced Features", "problem": "Users want AI-powered recommendations and analytics", "area_type": "feature"}}
            ]
            '''
            
            response = client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[
                    {"role": "system", "content": "You are a product analytics specialist who identifies key problem areas from customer feedback."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.2,
                response_format={"type": "json_object"}
            )
            
            result_text = response.choices[0].message.content
            result = json.loads(result_text)
            
            # Extract final areas
            final_areas = []
            if isinstance(result, dict) and any(isinstance(v, list) for v in result.values()):
                for k, v in result.items():
                    if isinstance(v, list) and all(isinstance(x, dict) for x in v):
                        if all('area' in x and 'area_type' in x for x in v):
                            final_areas = v
                            break
            elif isinstance(result, list) and all(isinstance(x, dict) for x in result):
                if all('area' in x and 'area_type' in x for x in result):
                    final_areas = result
            
            # If we couldn't parse the consolidated areas properly, use the raw chunk areas
            if not final_areas:
                logger.warning("Could not parse consolidated areas, using raw chunk areas")
                # Try to handle both dictionary and string formats
                if chunk_areas and isinstance(chunk_areas[0], dict) and 'area' in chunk_areas[0]:
                    final_areas = chunk_areas
                else:
                    # Convert strings to dictionaries
                    final_areas = [{"area": area, "problem": f"Issues related to {area.lower()}", "area_type": "issue"} 
                                  for area in chunk_areas if isinstance(area, str)]
        else:
            # Create default areas if none were identified
            final_areas = [
                {"area": "General Issues", "problem": "Various general problems with the service or product", "area_type": "issue"},
                {"area": "Feature Requests", "problem": "Customer requests for new features and enhancements", "area_type": "feature"},
                {"area": "Equipment Problems", "problem": "Issues with physical equipment or hardware", "area_type": "issue"},
                {"area": "Enhanced Capabilities", "problem": "Users want advanced features and integrations", "area_type": "feature"}
            ]
        
        logger.info(f"Consolidated into {len(final_areas)} final areas")
        for i, area in enumerate(final_areas):
            area_type = area.get('area_type', 'unknown')
            logger.info(f"  {i+1}. [{area_type.upper()}] {area.get('area')}: {area.get('problem')}")
            
        # STAGE 3: Use embeddings for fast classification
        # Generate embeddings for key areas - use only the area names for embedding
        problem_texts = [area.get('problem', f"Problem related to {area.get('area', 'unknown')}") for area in final_areas]
        key_area_embeddings = await get_embeddings(AZURE_CLIENT, problem_texts)
        area_names = [area.get('area') for area in final_areas]
        
        # Generate embeddings for all feedback
        feedback_texts = [item["text"] for item in all_feedbacks]
        
        # Process in batches to handle potential size limitations
        batch_size = 300
        all_feedback_embeddings = []
        
        for i in range(0, len(feedback_texts), batch_size):
            batch = feedback_texts[i:i + batch_size]
            logger.info(f"Generating embeddings for feedback batch {i//batch_size + 1}/{(len(feedback_texts) + batch_size - 1)//batch_size}")
            batch_embeddings = await get_embeddings(AZURE_CLIENT, batch)
            all_feedback_embeddings.extend(batch_embeddings)
        
        # Classify feedback by type (feature vs issue)
        feedback_types = await classify_feedback_by_type(AZURE_CLIENT, all_feedback_embeddings, feedback_texts)
        
        # Create area type mapping
        area_type_map = {area["area"]: area.get("area_type", "issue") for area in final_areas}
        
        # Quick preliminary classification for summary generation
        key_area_matrix = np.array(key_area_embeddings)
        feedback_matrix = np.array(all_feedback_embeddings)
        
        # Verify dimensions match
        key_area_dim = key_area_matrix.shape[1]
        feedback_dim = feedback_matrix.shape[1]
        logger.info(f"Embedding dimensions - Key areas: {key_area_matrix.shape}, Feedback: {feedback_matrix.shape}")
        
        if key_area_dim != feedback_dim:
            logger.error(f"Embedding dimensions mismatch: key_areas({key_area_dim}) != feedback({feedback_dim})")
            raise ValueError(f"Embedding dimensions don't match: {key_area_dim} vs {feedback_dim}")
            
        similarity_matrix = cosine_similarity_matrix(feedback_matrix, key_area_matrix)
        
        # Quick preliminary classification for summary generation
        preliminary_classified = {area.get('area', f"Area {i+1}"): [] for i, area in enumerate(final_areas)}
        for i, similarities in enumerate(similarity_matrix):
            best_match = np.argmax(similarities)
            if best_match < len(area_names):
                area = area_names[best_match]
                preliminary_classified[area].append(feedback_texts[i])
        
        # Generate insight summaries based on embeddings
        logger.info("Generating insight summaries from feedback embeddings")
        insight_summary = await generate_summary(
            client, 
            all_feedback_embeddings, 
            feedback_texts,
            preliminary_classified,  # Pass preliminary classification
            final_areas  # Pass key areas with area_type
        )
        
        # Calculate sentiment scores using the same embeddings we already generated
        logger.info("Calculating sentiment scores for all feedback items")
        sentiment_scores = await calculate_sentiment_scores(all_feedback_embeddings, client)
        
        # Create a mapping from text to sentiment score
        text_to_sentiment = {feedback_texts[i]: sentiment_scores[i] for i in range(len(feedback_texts))}
        
        # Function to calculate cosine similarity matrix
        def cosine_similarity_matrix(A, B):
            # Normalize the matrices
            A_norm = A / np.linalg.norm(A, axis=1, keepdims=True)
            B_norm = B / np.linalg.norm(B, axis=1, keepdims=True)
            # Calculate similarity matrix
            similarity = np.dot(A_norm, B_norm.T)
            return similarity
        
        # Apply type-based weighting for final classification
        for i in range(similarity_matrix.shape[0]):
            feedback_type = feedback_types[i]
            for j in range(similarity_matrix.shape[1]):
                area_name = area_names[j]
                area_type = area_type_map.get(area_name, "issue")
                
                # Apply weighting based on type match
                if feedback_type == area_type:
                    similarity_matrix[i, j] *= 1.2  # Boost similarity for matching types
                else:
                    similarity_matrix[i, j] *= 0.8  # Reduce similarity for mismatched types
        
        # Classify feedback based on similarity
        similarity_threshold = 0.5  # Lower threshold to match more feedback
        classified_feedback = {area.get('area'): [] for area in final_areas}
        
        for i, similarities in enumerate(similarity_matrix):
            # Get indices where similarity exceeds threshold
            matches = np.where(similarities > similarity_threshold)[0]
            
            # If no matches, use the best match
            if len(matches) == 0:
                best_match = np.argmax(similarities)
                matches = [best_match]
            
            # Add feedback to all matching areas
            for match_idx in matches:
                if match_idx < len(final_areas):
                    area = final_areas[match_idx].get('area')
                    classified_feedback[area].append(feedback_texts[i])
        
        logger.info(f"Classified {len(all_feedbacks)} feedback items into {len(final_areas)} areas")
        
        # Create enriched feedback with sentiment scores
        text_to_item = {item["text"]: item for item in all_feedbacks}
        def clean_json_values(obj):
            if isinstance(obj, dict):
                return {k: clean_json_values(v) for k, v in obj.items()}
            elif isinstance(obj, list):
                return [clean_json_values(i) for i in obj]
            elif isinstance(obj, float) and (np.isnan(obj) or np.isinf(obj)):
                return -1
            else:
                return obj
                
        enriched_feedback = {}
        
        for area, feedbacks in classified_feedback.items():
            enriched_feedback[area] = []
            for feedback_text in feedbacks:
                if feedback_text in text_to_item:
                    # Include original row data
                    original_item = text_to_item[feedback_text]
                    original_row = original_item.get("original_row", {"text": feedback_text})
                    original_row = clean_json_values(original_row)
                    
                    # Add sentiment score if available
                    if feedback_text in text_to_sentiment:
                        original_row["sentiment_score"] = text_to_sentiment[feedback_text]
                    else:
                        # Use a neutral sentiment as fallback
                        original_row["sentiment_score"] = 0.0
                        
                    enriched_feedback[area].append(original_row)
                else:
                    # Fallback if text not found in mapping
                    enriched_feedback[area].append({"text": feedback_text, "sentiment_score": 0.0})
        
        # Replace classified_feedback with enriched version
        classified_feedback = enriched_feedback
        
        return {
            "key_areas": final_areas,  # Now includes area_type
            "classified_feedback": classified_feedback,
            "insight_summary": insight_summary
        }
        
    except Exception as e:
        logger.error(f"Error in raw data analysis: {str(e)}\n{traceback.format_exc()}")
        # Create default areas if exception occurs
        basic_areas = [
            {"area": "General Issues", "problem": "Various general problems with the service or product", "area_type": "issue"},
            {"area": "Feature Requests", "problem": "Customer requests for new features", "area_type": "feature"},
            {"area": "Equipment Problems", "problem": "Issues with physical equipment or hardware", "area_type": "issue"},
            {"area": "Service Enhancements", "problem": "Requests for service improvements", "area_type": "feature"}
        ]
        # Create a simple fallback classification
        feedback_by_area = {}
        for i, area in enumerate(basic_areas):
            area_name = area.get('area')
            start_idx = (len(all_feedbacks) * i) // len(basic_areas)
            end_idx = (len(all_feedbacks) * (i+1)) // len(basic_areas)
            feedback_by_area[area_name] = [{"text": item["text"], "sentiment_score": 0.0} for item in all_feedbacks[start_idx:end_idx]]
        
        # Fallback insight summary
        fallback_insight_summary = {
            "user_loves": "Unable to determine what users love due to an error in data processing",
            "feature_request": "Unable to identify feature requests due to an error in data processing",
            "pain_point": "Unable to identify pain points due to an error in data processing",
            "overall_summary": "Unable to generate an overall summary due to an error"
        }
        
        return {
            "key_areas": basic_areas,
            "classified_feedback": feedback_by_area,
            "insight_summary": fallback_insight_summary
        }
async def process_csv_data(
    client: AzureOpenAI, 
    csv_data: str,
    source: str = None,
    slow_mode: bool = False,
    skip_insights: bool = False
) -> Dict[str, Any]:
    '''
    Process structured CSV data to extract and classify feedback.
    
    Args:
        client: Azure OpenAI client
        csv_data: CSV data as string
        source: Optional source value to override the Source column
        slow_mode: If True, use the original sequential processing
        
    Returns:
        Dictionary with key areas, classified feedback, and insight summaries
    '''
    try:
        # Parse CSV data
        df = pd.read_csv(StringIO(csv_data))
        
        if df.empty or len(df.columns) == 0:
            raise ValueError("Empty or invalid CSV data")
            
        logger.info(f"Parsed CSV with {len(df)} rows and {len(df.columns)} columns")
        logger.info(f"Original columns: {list(df.columns)}")
        
        # Use OpenAI to identify relevant columns
        columns = await identify_relevant_columns(AZURE_CLIENT, df)
        feedback_col = columns.get("feedback_col")
        rating_col = columns.get("rating_col")
        
        logger.info(f"Using columns: feedback={feedback_col}, rating={rating_col}")
        
        # Skip processing if no feedback column identified
        if not feedback_col or feedback_col not in df.columns:
            logger.warning(f"No feedback column identified, creating placeholder dataframe")
            # Create a placeholder dataframe with standardized columns
            df = pd.DataFrame({
                "Customer Feedback": df.iloc[:, 0] if len(df.columns) > 0 else ["No feedback data"],
                "Received": "N/A",
                "Name": "N/A",
                "Source": source if source else "N/A"
            })
            feedback_col = "Customer Feedback"
        else:
            # Apply standardization FIRST
            df = standardize_dataframe(df, columns, source)
            logger.info(f"Standardized columns: {list(df.columns)}")
            
            # Update column references to use standardized names
            feedback_col = "Customer Feedback"  # Always use standardized name
            
            # Check if rating column exists in standardized dataframe
            rating_col = None
            for col in df.columns:
                if col not in ["Received", "Name", "Customer Feedback", "Source"] and "rating" in col.lower():
                    rating_col = col
                    break
            
            # Skip 5-star reviews if rating column exists
            # if rating_col:
            #     try:
            #         df[rating_col] = pd.to_numeric(df[rating_col], errors='coerce')
            #         high_rating_mask = df[rating_col] >= 5
            #         five_star_count = high_rating_mask.sum()
            #         logger.info(f"Skipping {five_star_count} high-rated reviews (5-star)")
            #         df = df[~high_rating_mask]
            #     except Exception as e:
            #         logger.warning(f"Could not filter out 5-star reviews, processing all: {str(e)}")
        
        # Extract feedback with ratings - vectorized approach from standardized dataframe
        valid_rows = df[df[feedback_col].notna() & (df[feedback_col].astype(str).str.strip() != "")]
        all_feedback_items = []
        
        for idx, row in valid_rows.iterrows():
            feedback_text = str(row[feedback_col])
            
            # Add rating info if available
            if rating_col and rating_col in df.columns and pd.notna(row[rating_col]):
                feedback_text = f"[Rating: {row[rating_col]}] {feedback_text}"
            
            # Create row dictionary from standardized dataframe
            row_dict = row.replace([float('nan'), float('inf'), float('-inf')], -1).to_dict()
            
            all_feedback_items.append({
                "text": feedback_text,
                "original_index": idx,
                "original_row": row_dict  # Now contains standardized column names
            })
        
        if not all_feedback_items:
            raise ValueError("No valid feedback found in the CSV data")
            
        logger.info(f"Extracted {len(all_feedback_items)} feedback items")
            
        # Sample for key area identification
        # Sample for key area identification - ensure we use all reviews for small datasets
        if len(all_feedback_items) <= 10:
            sample_feedbacks = [item["text"] for item in all_feedback_items]
        else:
            sample_feedbacks = [item["text"] for item in all_feedback_items[:min(1000, len(all_feedback_items) // 2)]]
        #sample_feedbacks = [item["text"] for item in all_feedback_items[:min(1000, len(all_feedback_items) // 2)]]
        sample_feedback_text = '\n'.join(sample_feedbacks)
        
        # Identify key areas in parallel with other initialization
        key_area_task = asyncio.create_task(identify_key_areas(client, sample_feedback_text))
        
        # Process feedback in batches for embedding - start while key areas are being identified
        feedback_texts = [item["text"] for item in all_feedback_items]
        
        # Get key areas from the parallel task
        key_areas = await key_area_task
        logger.info(f"Identified {len(key_areas)} key areas")
        
        # Initialize result structure with empty lists - use the area name as key
        classified_feedback = {area.get('area', f"Area {i+1}"): [] for i, area in enumerate(key_areas)}
        
        # Use embeddings for classification
        try:
            logger.info("Using embeddings for feedback classification")
            
            # Create a list of just the area names for embeddings
            problem_texts = [area.get('problem', f"Problem {i+1}") for i, area in enumerate(key_areas)]
            key_area_embeddings_task = asyncio.create_task(get_embeddings(client, problem_texts))
            area_names = [area.get('area', f"Area {i+1}") for i, area in enumerate(key_areas)]
            
            # Process feedback in batches for embedding - parallelize with area embeddings
            batch_size = 300
            feedback_embeddings_tasks = []
            
            for i in range(0, len(feedback_texts), batch_size):
                batch = feedback_texts[i:i + batch_size]
                logger.info(f"Generating embeddings for feedback batch {i//batch_size + 1}/{(len(feedback_texts) + batch_size - 1)//batch_size}")
                task = asyncio.create_task(get_embeddings(client, batch))
                feedback_embeddings_tasks.append(task)
            
            # Wait for all embedding tasks to complete
            all_embedding_tasks = [key_area_embeddings_task] + feedback_embeddings_tasks
            embedding_results = await asyncio.gather(*all_embedding_tasks)
            
            # First result is key area embeddings, rest are feedback embedding batches
            key_area_embeddings = embedding_results[0]
            all_feedback_embeddings = []
            for result in embedding_results[1:]:
                all_feedback_embeddings.extend(result)
            
            # Classify feedback by type (feature vs issue)
            feedback_types = await classify_feedback_by_type(client, all_feedback_embeddings, feedback_texts)
            
            # Generate insight summaries using the embeddings - pass classified_feedback and key_areas
            logger.info("Generating insight summaries from feedback embeddings")
            # We need to pre-calculate classified_feedback for summary generation
            # First do a preliminary classification
            key_area_matrix = np.array(key_area_embeddings)
            feedback_matrix = np.array(all_feedback_embeddings)
            
            # Calculate similarity with optimized matrix operations
            A_norm = feedback_matrix / np.linalg.norm(feedback_matrix, axis=1, keepdims=True)
            B_norm = key_area_matrix / np.linalg.norm(key_area_matrix, axis=1, keepdims=True)
            similarity_matrix = np.dot(A_norm, B_norm.T)
            
            # Quick preliminary classification for summary generation
            preliminary_classified = {area.get('area', f"Area {i+1}"): [] for i, area in enumerate(key_areas)}
            for i, similarities in enumerate(similarity_matrix):
                best_match = np.argmax(similarities)
                if best_match < len(area_names):
                    area = area_names[best_match]
                    preliminary_classified[area].append(feedback_texts[i])
            
            insight_summary = await generate_summary(
                client, 
                all_feedback_embeddings, 
                feedback_texts,
                preliminary_classified,  # Pass preliminary classification
                key_areas,  # Pass key areas with area_type
                skip_insights=skip_insights
            )
            
            # Calculate sentiment scores using the same embeddings we already generated
            logger.info("Calculating sentiment scores for all feedback items")
            sentiment_scores = await calculate_sentiment_scores(all_feedback_embeddings, client)
            
            # Create a mapping from text to sentiment score
            text_to_sentiment = {feedback_texts[i]: sentiment_scores[i] for i in range(len(feedback_texts))}
            
            # Create area type mapping
            area_type_map = {area["area"]: area.get("area_type", "issue") for area in key_areas}
            
            # Classify feedback based on similarity with type preference

            # Adjust similarity threshold based on dataset size
            if len(all_feedback_items) <= 5:
                similarity_threshold = 0.65  # Lower threshold for small datasets
            else:
                similarity_threshold = 0.8
            
            # Apply type-based weighting
            for i in range(similarity_matrix.shape[0]):
                feedback_type = feedback_types[i]
                for j in range(similarity_matrix.shape[1]):
                    area_name = area_names[j]
                    area_type = area_type_map.get(area_name, "issue")
                    
                    # Apply weighting based on type match
                    if feedback_type == area_type:
                        similarity_matrix[i, j] *= 1.2  # Boost similarity for matching types
                    else:
                        similarity_matrix[i, j] *= 0.8  # Reduce similarity for mismatched types
            
            # Vectorized matching with weighted similarities
            matches = similarity_matrix > similarity_threshold
            
            # For rows with no matches, use the best match
            no_matches = ~np.any(matches, axis=1)
            if np.any(no_matches):
                best_matches = np.argmax(similarity_matrix[no_matches], axis=1)
                for i, row_idx in enumerate(np.where(no_matches)[0]):
                    matches[row_idx, best_matches[i]] = True
            
            # Add feedback to matching areas using the matches matrix
            for i in range(matches.shape[0]):
                for j in np.where(matches[i])[0]:
                    if j < len(area_names):
                        area = area_names[j]
                        classified_feedback[area].append(feedback_texts[i])
            
            logger.info(f"Successfully classified {len(all_feedback_embeddings)} feedback items using embeddings")
            
        except Exception as e:
            logger.error(f"Error with embedding classification: {str(e)}")
            logger.info("Falling back to direct OpenAI classification")
            
            # Fallback for insight summary if embeddings failed
            insight_summary = {
                "user_loves": "Insights will be generated via completions" if skip_insights else "Unable to determine what users love due to an error in embeddings processing",
                "feature_request": "Insights will be generated via completions" if skip_insights else "Unable to identify feature requests due to an error in embeddings processing",
                "pain_point": "Insights will be generated via completions" if skip_insights else "Unable to identify pain points due to an error in embeddings processing",
                "overall_summary": "Insights will be generated via completions" if skip_insights else "Unable to generate an overall summary due to an error"
            }
            
            # Fall back to chunked classification via OpenAI
            chunk_size = 10
            
            # Using semaphore to limit concurrent API calls
            semaphore = asyncio.Semaphore(3)  # Limit to 3 concurrent API calls
            
            async def process_chunk(chunk_idx, chunk):
                async with semaphore:
                    logger.info(f"Processing feedback chunk {chunk_idx + 1}/{(len(all_feedback_items) + chunk_size - 1)//chunk_size}")
                    
                    chunk_texts = [item["text"] for item in chunk]
                    
                    # Create prompt and classify
                    areas_text = "\n".join([f"{j+1}. {area.get('area', f'Area {j+1}')} ({area.get('area_type', 'issue')})" for j, area in enumerate(key_areas)])
                    feedback_text = "\n".join([f"Feedback {j+1}: {text}" for j, text in enumerate(chunk_texts)])
                    
                    prompt = f'''
                    I need to classify these customer feedback items into the most relevant categories.
                    
                    CATEGORIES (with their type):
                    {areas_text}
                    
                    FEEDBACK ITEMS:
                    {feedback_text}
                    
                    For each feedback item, tell me which categories it belongs to (multiple categories allowed).
                    Prefer assigning feedback to categories of matching type (feature requests to feature categories, issues to issue categories).
                    
                    Format as JSON:
                    {{
                        "classifications": [
                            {{
                                "feedback_index": 0,
                                "category_indices": [0, 2]
                            }},
                            ...
                        ]
                    }}
                    
                    Category indices are 0-based, corresponding to the numbered list above.
                    '''
                    
                    try:
                        response = client.chat.completions.create(
                            model="gpt-4.1-mini",
                            messages=[
                                {"role": "system", "content": "You are a customer feedback classification system."},
                                {"role": "user", "content": prompt}
                            ],
                            temperature=0.1,
                            response_format={"type": "json_object"}
                        )
                        
                        result = json.loads(response.choices[0].message.content)
                        classifications = []
                        
                        # Process classifications
                        if "classifications" in result:
                            for classification in result["classifications"]:
                                feedback_idx = classification.get("feedback_index", 0)
                                category_indices = classification.get("category_indices", [])
                                
                                if feedback_idx < len(chunk_texts):
                                    feedback_text = chunk_texts[feedback_idx]
                                    
                                    if not category_indices and key_areas:
                                        first_area = key_areas[0].get('area', 'Area 1')
                                        classifications.append((first_area, feedback_text))
                                    else:
                                        for cat_idx in category_indices:
                                            if 0 <= cat_idx < len(key_areas):
                                                area = key_areas[cat_idx].get('area', f'Area {cat_idx+1}')
                                                classifications.append((area, feedback_text))
                        return classifications
                    
                    except Exception as chunk_error:
                        logger.error(f"Error processing chunk {chunk_idx + 1}: {str(chunk_error)}")
                        # Assign all feedbacks in this chunk to the first category as fallback
                        if key_areas:
                            first_area = key_areas[0].get('area', 'Area 1')
                            return [(first_area, text) for text in chunk_texts]
                        return []
            
            # Prepare chunks
            chunks = []
            for i in range(0, len(all_feedback_items), chunk_size):
                chunk = all_feedback_items[i:i + min(chunk_size, len(all_feedback_items) - i)]
                chunks.append(chunk)
            
            # Process all chunks in parallel with rate limiting
            chunk_results = await asyncio.gather(
                *[process_chunk(i, chunk) for i, chunk in enumerate(chunks)]
            )
            
            # Add all classified feedback to the result
            for classifications in chunk_results:
                for area, feedback_text in classifications:
                    classified_feedback[area].append(feedback_text)
            
            # Generate fallback sentiment scores if we don't have embeddings
            if not text_to_sentiment:
                # Use a simple keyword-based approach for fallback
                positive_keywords = ["good", "great", "excellent", "amazing", "love", "best", "happy", "perfect", "satisfied"]
                negative_keywords = ["bad", "poor", "terrible", "awful", "hate", "worst", "unhappy", "disappointed", "failure"]
                
                for text in feedback_texts:
                    text_lower = text.lower()
                    pos_count = sum(1 for word in positive_keywords if word in text_lower)
                    neg_count = sum(1 for word in negative_keywords if word in text_lower)
                    # Simple formula: (pos - neg) / (pos + neg + 1) to get score between -1 and 1
                    sentiment = (pos_count - neg_count) / (pos_count + neg_count + 1) if (pos_count + neg_count) > 0 else 0
                    text_to_sentiment[text] = sentiment
            
            # If we couldn't generate insight summary from embeddings, try to generate it via direct OpenAI call
            if "Unable to determine what users love" in insight_summary["user_loves"]:
                try:
                    # Take a sample of reviews for each category 
                    sample_size = min(200, len(feedback_texts))
                    sample_reviews = feedback_texts[:sample_size]
                    combined_reviews = "\n".join([f"- {review}" for review in sample_reviews])
                    
                    # Generate a summary directly from OpenAI
                    prompt = f'''
                    I have a collection of {len(feedback_texts)} customer reviews. Here's a sample:
                    
                    {combined_reviews}
                    
                    Based on these reviews, please provide:
                    1. What users love most about the product (focus on concrete features)
                    2. What features users are most commonly requesting
                    3. What are the biggest pain points users are experiencing
                    4. An overall summary of the feedback
                    
                    Keep each response to a single concise sentence.
                    Format as a JSON object with these keys: user_loves, feature_request, pain_point, overall_summary
                    '''
                    
                    response = client.chat.completions.create(
                        model="gpt-4.1-mini",
                        messages=[
                            {"role": "system", "content": "You are a customer insight specialist who extracts clear, actionable insights from reviews."},
                            {"role": "user", "content": prompt}
                        ],
                        temperature=0.3,
                        response_format={"type": "json_object"}
                    )
                    
                    direct_summary = json.loads(response.choices[0].message.content)
                    insight_summary = {
                        "user_loves": direct_summary.get("user_loves", "No clear indication of what users love"),
                        "feature_request": direct_summary.get("feature_request", "No clear feature requests identified"),
                        "pain_point": direct_summary.get("pain_point", "No clear pain points identified"),
                        "overall_summary": direct_summary.get("overall_summary", "Insufficient feedback for overall summary")
                    }
                except Exception as summary_error:
                    logger.error(f"Error generating direct summary: {str(summary_error)}")
                    # Keep existing fallback summaries
        
        # Log classification stats
        logger.info("Classification results:")
        for area, feedbacks in classified_feedback.items():
            area_type = area_type_map.get(area, "unknown") if 'area_type_map' in locals() else "unknown"
            logger.info(f"  - {area} [{area_type}]: {len(feedbacks)} items")
        
        # Map original data back to each feedback
        text_to_item = {item["text"]: item for item in all_feedback_items}
        enriched_feedback = {}
        
        for area, feedbacks in classified_feedback.items():
            enriched_feedback[area] = []
            for feedback_text in feedbacks:
                if feedback_text in text_to_item:
                    # Include original row data (which now has standardized column names)
                    original_item = text_to_item[feedback_text]
                    original_row = original_item.get("original_row", {"text": feedback_text})
                    
                    # Add sentiment score if available
                    if feedback_text in text_to_sentiment:
                        original_row["sentiment_score"] = text_to_sentiment[feedback_text]
                    else:
                        # Use a neutral sentiment as fallback
                        original_row["sentiment_score"] = 0.0
                        
                    enriched_feedback[area].append(original_row)
                else:
                    # Fallback if text not found in mapping - ensure standardized columns
                    enriched_feedback[area].append({
                        "Customer Feedback": feedback_text, 
                        "Received": "N/A",
                        "Name": "N/A",
                        "Source": source if source else "N/A",
                        "sentiment_score": 0.0
                    })
        
        # Replace classified_feedback with enriched version
        classified_feedback = enriched_feedback

        return {
            "key_areas": key_areas,  # Now includes area_type
            "classified_feedback": classified_feedback,
            "insight_summary": insight_summary
        }
        
    except Exception as e:
        logger.error(f"Error processing CSV data: {str(e)}\n{traceback.format_exc()}")
        # Fall back to raw data processing
        logger.info("Falling back to raw data processing")
        return await analyze_raw_data_chunks(client, csv_data)
async def format_analysis_results(analysis_results: Dict[str, Any], return_raw_feedback: bool = False) -> Dict[str, Any]:
    '''
    Format the analysis results into the final structure.
    
    Args:
        analysis_results: Dictionary with key areas, classified feedback, and insight summary
        return_raw_feedback: Whether to include raw feedback text in the response
        
    Returns:
        Formatted results matching the required output structure
    '''
    key_areas = analysis_results.get("key_areas", [])
    classified_feedback = analysis_results.get("classified_feedback", {})
    insight_summary = analysis_results.get("insight_summary", {
        "user_loves": "No clear insights on what users love",
        "feature_request": "No clear feature requests identified",
        "pain_point": "No clear pain points identified"
    })
    
    # Format the results
    formatted_results = []
    
    # Process each key area
    for area_obj in key_areas:
        if isinstance(area_obj, dict):
            # Get area name and problem from the dictionary
            area_name = area_obj.get("area", "Unknown Area")
            area_problem = area_obj.get("problem", f"Issues related to {area_name.lower()}")
            area_type = area_obj.get("area_type", "issue")  # NEW LINE - Get area_type
        else:
            # Fallback for string-only areas
            area_name = str(area_obj)
            area_problem = f"Issues related to {area_name.lower()}"
            area_type = "issue"  # NEW LINE - Default to issue
        
        # Get the feedback for this area
        area_feedbacks = classified_feedback.get(area_name, [])
        
        result = {
            "key_area": area_name,
            "customer_problem": area_problem,
            "number_of_users": len(area_feedbacks),
            "type": area_type  # NEW LINE - Add type field
        }
        
        # Only include raw feedback if requested
        if return_raw_feedback:
            # Filter the raw feedback to include only desired columns
            filtered_feedbacks = []
            allowed_columns = {"Received", "Name", "Customer Feedback", "Source", "sentiment_score"}
            
            for feedback in area_feedbacks:
                # Create a new dictionary with only the allowed columns
                filtered_feedback = {}
                for key, value in feedback.items():
                    if key in allowed_columns:
                        filtered_feedback[key] = value
                
                # Add the filtered feedback
                filtered_feedbacks.append(filtered_feedback)
            
            result["raw_feedbacks"] = filtered_feedbacks
            
        formatted_results.append(result)
    # Remove key areas with no feedback
    original_count = len(formatted_results)
    formatted_results = [result for result in formatted_results if result["number_of_users"] > 0]
    removed_count = original_count - len(formatted_results)
    if removed_count > 0:
        logger.info(f"Removed {removed_count} key areas with no feedback")

    # Sort by number of users (descending)
    formatted_results.sort(key=lambda x: x["number_of_users"], reverse=True)
    
    return {
        "analysis_results": formatted_results,
        "summary": {
            "total_feedback_items": sum(len(feedbacks) for feedbacks in classified_feedback.values()),
            "total_key_areas": len(formatted_results)
        },
        "insight_summary": insight_summary  # Add the insight summary to the return value
    }
def integrate_with_main_app(app):
    '''
    Integrates the customer feedback analysis router with a FastAPI app.
    
    Args:
        app: FastAPI application instance
    '''
    app.include_router(router, prefix="/feedback", tags=["customer_feedback"])
    logger.info("Customer feedback analysis API integrated with main app")

# Standalone app routes
@app.post("/analyze-feedback")
@async_timeout(180)
async def analyze_feedback(
    file: UploadFile = File(...),
    return_raw_feedback: bool = Form(False),
    source: str = Form(None),
    extraction_prompt: str = Form(None),
    mode: str = Form("auto")  # NEW: Processing mode - 'auto', 'csv', 'completions'
) -> JSONResponse:
    '''
    Analyzes customer feedback from any uploaded file type.
    
    Modes:
    - 'auto' (default): Enhanced processing with insight generation
    - 'csv': Original CSV/Excel processing without enhancements
    - 'completions': Force all processing through completions fallback
    '''
    job_id = f"feedback_analysis_{int(time.time())}"
    
    try:
        start_time = time.time()
        logger.info(f"[JOB {job_id}] Starting analysis of file: {file.filename}")
        logger.info(f"[JOB {job_id}] Mode: {mode}")
        logger.info(f"[JOB {job_id}] Return raw feedback: {return_raw_feedback}")
        logger.info(f"[JOB {job_id}] Source: {source or 'None'}")
        
        # Validate mode
        if mode not in ["auto", "csv", "completions"]:
            logger.warning(f"[JOB {job_id}] Invalid mode '{mode}', defaulting to 'auto'")
            mode = "auto"
        
        # Read file
        file_content = await file.read()
        file_size = len(file_content)
        logger.info(f"[JOB {job_id}] File size: {file_size/1024:.1f} KB")
        
        # Validate size
        MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
        if file_size > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413,
                detail=f"File too large. Maximum size is {MAX_FILE_SIZE/1024/1024:.0f}MB"
            )
        
        # Extract raw content (always needed for completions mode or fallback)
        raw_content = None
        if mode in ["auto", "completions"]:
            raw_content = await extract_raw_content_from_file(file_content, file.filename)
            logger.info(f"[JOB {job_id}] Extracted {len(raw_content)} chars of raw content")
        
        # COMPLETIONS MODE: Skip directly to completions fallback
        if mode == "completions":
            logger.info(f"[JOB {job_id}] Mode is 'completions', using direct generation")
            
            complete_analysis = await generate_complete_analysis_via_completions(
                AZURE_CLIENT,
                raw_content,
                source,
                return_raw_feedback
            )
            
            # Return directly with metadata
            complete_analysis["metadata"] = {
                "job_id": job_id,
                "original_filename": file.filename,
                "file_size_kb": round(file_size / 1024, 2),
                "processing_time_seconds": round(time.time() - start_time, 2),
                "file_type": os.path.splitext(file.filename)[1].lower(),
                "source": source,
                "processing_method": "completions_mode",
                "mode": mode
            }
            
            total_items = complete_analysis.get("summary", {}).get("total_feedback_items", 0)
            total_areas = complete_analysis.get("summary", {}).get("total_key_areas", 0)
            logger.info(f"[JOB {job_id}] Complete (completions mode): {total_areas} areas, {total_items} items")
            
            return JSONResponse(content=complete_analysis, status_code=200)
        
        # Determine file type
        file_ext = os.path.splitext(file.filename)[1].lower()
        
        # Track processing success
        analysis_results = None
        processing_method = None
        use_fallback = False
        
        # STEP 1: Try direct processing for CSV/Excel
        if file_ext in ['.csv', '.xlsx', '.xls', '.xlsm']:
            logger.info(f"[JOB {job_id}] Processing {file_ext} directly")
            try:
                if file_ext in ['.xlsx', '.xls', '.xlsm']:
                    analysis_results = await process_excel_data(AZURE_CLIENT, file_content, file.filename, source, skip_insights=(mode == "auto"))
                    processing_method = "excel_direct"
                else:  # CSV
                    # For CSV mode, we need raw content for validation
                    if mode == "csv" and not raw_content:
                        raw_content = await extract_raw_content_from_file(file_content, file.filename)
                    
                    # Validate it's actually CSV
                    if not raw_content or (',' in raw_content and '\n' in raw_content):
                        csv_content = raw_content if raw_content else file_content.decode('utf-8', errors='ignore')
                        analysis_results = await process_csv_data(AZURE_CLIENT, csv_content, source)
                        processing_method = "csv_direct"
                    else:
                        logger.warning(f"[JOB {job_id}] File has .csv extension but doesn't look like CSV")
                        file_ext = '.txt'  # Treat as text
                        
                # Check if results are valid (only in auto mode)
                if mode == "auto" and analysis_results and not sanity_check_analysis_results(analysis_results):
                    logger.warning(f"[JOB {job_id}] Direct processing produced poor results, will use fallback")
                    use_fallback = True
                    
            except Exception as e:
                logger.error(f"[JOB {job_id}] Direct processing failed: {str(e)}")
                analysis_results = None
        
        # STEP 2: Convert other files to CSV (skip in csv mode for non-CSV files)
        if not analysis_results or file_ext not in ['.csv', '.xlsx', '.xls', '.xlsm']:
            if mode == "csv":
                # In CSV mode, only process actual CSV/Excel files
                logger.info(f"[JOB {job_id}] CSV mode: Skipping non-CSV/Excel file")
                raise ValueError(f"CSV mode only supports CSV and Excel files, got {file_ext}")
            
            logger.info(f"[JOB {job_id}] Converting {file_ext} to CSV via API")
            
            # Build extraction prompt
            if not extraction_prompt:
                if file_ext in ['.txt', '.text', '.log', '.md', '']:
                    extraction_prompt = f'''
Extract customer feedback from this text file into CSV format.

CRITICAL INSTRUCTIONS:
1. Each distinct feedback item should be a separate row
2. Look for natural separators (line breaks, dates, names, numbers)
3. Extract VERBATIM - do not summarize or modify
4. If no clear pattern, treat paragraphs as separate feedback

OUTPUT COLUMNS (use EXACTLY these names):
- user: Customer name or "Customer 1", "Customer 2" etc
- review_text: The complete, unmodified feedback text
- date: Date if found (YYYY-MM-DD format), else "N/A"
- source: "{source or 'Text File'}"
- star_rating: Rating if found (1-5), else "N/A"

IMPORTANT: Extract ALL text that could be customer feedback.'''
                else:
                    extraction_prompt = f'''
Extract ALL customer feedback/reviews from this {file_ext} file.
Structure as CSV with EXACTLY these columns: user, review_text, date, source, star_rating.
Extract all feedback VERBATIM - do not modify or summarize.
Use 'N/A' for any missing values.'''
            
            # Try conversion with retries
            conversion_attempts = 3 if file_ext in ['.txt', '.text'] else 1
            
            for attempt in range(conversion_attempts):
                try:
                    logger.info(f"[JOB {job_id}] Conversion attempt {attempt + 1}/{conversion_attempts}")
                    
                    # Adjust prompt for retries
                    if attempt > 0:
                        extraction_prompt = "Extract to CSV format. Each line/paragraph is a feedback row. Columns: user,review_text,date,source,star_rating"
                    
                    csv_content = await convert_file_to_csv(
                        file_content=file_content,
                        filename=file.filename,
                        prompt=extraction_prompt,
                        mode="extract"
                    )
                    
                    if csv_content and csv_content.strip():
                        # Map API columns to our format
                        df = pd.read_csv(StringIO(csv_content))
                        if 'user' in df.columns or 'review_text' in df.columns:
                            df = df.rename(columns={
                                'user': 'Name',
                                'review_text': 'Customer Feedback',
                                'date': 'Received',
                                'source': 'Source',
                                'star_rating': 'Rating'
                            })
                            csv_content = df.to_csv(index=False)
                        
                        analysis_results = await process_csv_data(AZURE_CLIENT, csv_content, source or file_ext.upper(), skip_insights=(mode == "auto"))
                        processing_method = "csv_conversion"
                        
                        # Check if results are valid (only in auto mode)
                        if mode == "auto" and not sanity_check_analysis_results(analysis_results):
                            logger.warning(f"[JOB {job_id}] Conversion produced poor results")
                            use_fallback = True
                        else:
                            break  # Good results, stop retrying
                        
                except Exception as e:
                    logger.error(f"[JOB {job_id}] Conversion attempt {attempt + 1} failed: {str(e)}")
                    if attempt < conversion_attempts - 1:
                        await asyncio.sleep(1)
        
        # STEP 3: Enhance insights (only in auto mode)
        if mode == "auto" and analysis_results and not use_fallback:
            logger.info(f"[JOB {job_id}] Auto mode: Enhancing insights with direct generation")
            try:
                # Ensure we have raw content
                if not raw_content:
                    raw_content = await extract_raw_content_from_file(file_content, file.filename)
                
                # Generate enhanced insights from raw content
                enhanced_insights = await generate_insight_summary_direct(AZURE_CLIENT, raw_content)
                
                # Override existing insights with enhanced version
                analysis_results["insight_summary"] = enhanced_insights
                logger.info(f"[JOB {job_id}] Successfully enhanced insights")
                
            except Exception as e:
                logger.error(f"[JOB {job_id}] Insight enhancement failed: {str(e)}")
                # Keep original insights if enhancement fails
        
        # STEP 4: Use complete fallback if needed (only in auto mode)
        if mode == "auto" and (not analysis_results or use_fallback):
            logger.warning(f"[JOB {job_id}] Auto mode: Using complete analysis generation fallback")
            
            # Ensure we have raw content
            if not raw_content:
                raw_content = await extract_raw_content_from_file(file_content, file.filename)
            
            complete_analysis = await generate_complete_analysis_via_completions(
                AZURE_CLIENT,
                raw_content,
                source,
                return_raw_feedback
            )
            
            # Return directly with metadata
            complete_analysis["metadata"] = {
                "job_id": job_id,
                "original_filename": file.filename,
                "file_size_kb": round(file_size / 1024, 2),
                "processing_time_seconds": round(time.time() - start_time, 2),
                "file_type": file_ext,
                "source": source,
                "processing_method": "complete_generation_fallback",
                "mode": mode
            }
            
            total_items = complete_analysis.get("summary", {}).get("total_feedback_items", 0)
            total_areas = complete_analysis.get("summary", {}).get("total_key_areas", 0)
            logger.info(f"[JOB {job_id}] Complete (auto fallback): {total_areas} areas, {total_items} items")
            
            return JSONResponse(content=complete_analysis, status_code=200)
        
        # STEP 5: Format results for normal flow
        if not analysis_results:
            raise ValueError(f"No analysis results generated in {mode} mode")
            
        logger.info(f"[JOB {job_id}] Formatting results (method: {processing_method})")
        formatted_results = await format_analysis_results(analysis_results, return_raw_feedback)
        
        # Add metadata
        formatted_results["metadata"] = {
            "job_id": job_id,
            "original_filename": file.filename,
            "file_size_kb": round(file_size / 1024, 2),
            "processing_time_seconds": round(time.time() - start_time, 2),
            "file_type": file_ext,
            "source": source,
            "processing_method": processing_method,
            "insights_enhanced": mode == "auto",  # Only true in auto mode
            "mode": mode
        }
        
        # Log summary
        total_items = formatted_results.get("summary", {}).get("total_feedback_items", 0)
        total_areas = formatted_results.get("summary", {}).get("total_key_areas", 0)
        logger.info(f"[JOB {job_id}] Complete: {total_areas} areas, {total_items} items, {processing_method}, mode={mode}")
        
        return JSONResponse(content=formatted_results, status_code=200)
        
    except HTTPException:
        raise
        
    except Exception as e:
        logger.error(f"[JOB {job_id}] Critical error: {str(e)}\n{traceback.format_exc()}")
        
        # Emergency fallback response
        try:
            # For completions and auto mode, try to generate insights
            if mode in ["auto", "completions"] and 'raw_content' in locals() and raw_content:
                insights = await generate_insight_summary_direct(AZURE_CLIENT, raw_content)
            else:
                insights = {
                    "user_loves": "Unable to analyze due to critical error",
                    "feature_request": "Unable to analyze due to critical error",
                    "pain_point": f"System error: {str(e)[:100]}",
                    "overall_summary": "Analysis failed. Please try again or contact support."
                }
            
            # Return minimal valid structure
            return JSONResponse(
                content={
                    "analysis_results": [{
                        "key_area": "System Error",
                        "customer_problem": "Analysis could not be completed",
                        "number_of_users": 0,
                        "type": "issue"
                    }],
                    "summary": {"total_feedback_items": 0, "total_key_areas": 1},
                    "insight_summary": insights,
                    "metadata": {
                        "job_id": job_id,
                        "error": str(e),
                        "processing_time_seconds": round(time.time() - start_time, 2) if 'start_time' in locals() else 0,
                        "mode": mode if 'mode' in locals() else "unknown"
                    }
                },
                status_code=200
            )
        except:
            raise HTTPException(status_code=500, detail=f"Critical error: {str(e)}")
@app.post("/classify-single-review")
async def classify_single_review(
    request: Request = None,
    review: str = Form(None),
    existing_json: str = Form(None),
    file: UploadFile = File(None)
) -> JSONResponse:
    '''
    Classifies a single piece of text (review, email, meeting notes, etc.) into multiple categories.
    
    Args:
        request: FastAPI request object containing JSON
        review: Text to classify (when using form data)
        existing_json: Existing categories to consider (when using form data)
        
    Returns:
        JSONResponse with array of classifications, each containing:
        - key_area: Category name (2-4 words)
        - customer_problem: Specific problem/request description
        - sentiment_score: Float between -1.0 and 1.0
        - area_type: Either "feature" or "issue"
    '''
    try:
        logger.info(f"[classify-single-review] Endpoint called")
        if file is not None:
            logger.info(f"[classify-single-review] File uploaded: {file.filename}, extracting text")
            try:
                file_content = await file.read()
                review = await extract_text_internal(
                    file_content=file_content,
                    filename=file.filename,
                    strategy="auto",
                    languages=None,
                    encoding=None,
                    logger=logger
                )
                logger.info(f"[classify-single-review] Extracted {len(review)} characters from file")
            except Exception as e:
                logger.error(f"[classify-single-review] Failed to extract text from file: {str(e)}")
                raise HTTPException(status_code=422, detail=f"Failed to process file: {str(e)}")
        # Handle both form data and JSON body
        if review is not None:
            # Form data was provided
            logger.info(f"[classify-single-review] Processing form data input")
            review_text = review
            try:
                existing_categories = json.loads(existing_json) if existing_json else []
                logger.info(f"[classify-single-review] Parsed existing categories from form data")
            except json.JSONDecodeError as e:
                logger.error(f"[classify-single-review] Failed to parse existing_json: {str(e)}")
                raise HTTPException(status_code=400, detail="Invalid JSON format for existing categories")
        else:
            # Process as JSON body
            logger.info(f"[classify-single-review] Processing JSON body input")
            try:
                body = await request.json()
                review_text = body.get("review", "")
                existing_categories = body.get("existing_json", body.get("existing_categories", []))
                logger.info(f"[classify-single-review] Extracted data from JSON body")
            except Exception as e:
                logger.error(f"[classify-single-review] Failed to parse request body: {str(e)}")
                raise HTTPException(status_code=400, detail=f"Invalid request body: {str(e)}")
        
        if not review_text:
            raise HTTPException(status_code=400, detail="Text input is required")
        
        logger.info(f"[classify-single-review] Starting classification for text: {review_text[:100]}...")
        logger.info(f"[classify-single-review] Text length: {len(review_text)} characters")
        
        # Format existing categories for the prompt
        existing_categories_text = ""
        if existing_categories:
            logger.info(f"[classify-single-review] Found {len(existing_categories)} existing categories")
            for cat in existing_categories:
                logger.debug(f"[classify-single-review] Existing category: {cat.get('key_area', 'Unknown')} ({cat.get('area_type', 'issue')})")
            
            existing_categories_text = "\nEXISTING CATEGORIES TO CONSIDER:\n" + "\n".join([
                f"- {cat.get('key_area', 'Unknown')} ({cat.get('area_type', 'issue')}): {cat.get('customer_problem', 'No description')}"
                for cat in existing_categories
            ]) + "\n\nTry to use these existing categories when the content matches, but create new specific categories for anything that doesn't fit."
        else:
            logger.info(f"[classify-single-review] No existing categories provided")
        
        # Create a comprehensive prompt for multi-classification
        prompt = f'''
You are an expert at analyzing ANY type of text - customer reviews, meeting notes, transcripts, conversations, emails, chat logs, support tickets, or even casual discussions between people - and extracting ALL distinct problems, issues, complaints, feature requests, or suggestions mentioned.

INPUT TEXT TO ANALYZE:
"{review_text}"
{existing_categories_text}

CRITICAL INSTRUCTIONS:
1. The input could be ANYTHING: formal feedback, casual conversation between colleagues, meeting transcript, customer support chat, internal discussion, etc.
2. Extract problems/requests even if mentioned indirectly (e.g., "John was saying the dashboard is too slow" → dashboard performance issue)
3. Extract EACH distinct issue or feature as a separate item - don't over-consolidate
4. If someone mentions 3 different API needs, that's 3 separate items, not 1 "API improvements" item
5. Look for distinct action items, deliverables, or problems even if they're related to the same general topic
6. Classify each as either "feature" (requests/suggestions/enhancements) or "issue" (problems/complaints/bugs)
7. Calculate sentiment based on how each topic is discussed

IMPORTANT: Extract the most significant distinct items, aiming for 2-6 key areas maximum. While each distinct issue should be captured, consolidate closely related items when reasonable. For example, if an email mentions multiple API improvements, consider whether they represent one cohesive need or truly separate concerns. Quality over quantity - focus on the most impactful categories.

CLASSIFICATION RULES:
- "feature" = requests, suggestions, enhancements, improvements, "wish", "want", "need", "should have", "would be nice"
- "issue" = problems, bugs, complaints, difficulties, frustrations, "broken", "doesn't work", "confusing", "slow", "failing"

CONSOLIDATION GUIDELINES:
- Maximum 6 categories total - be selective and focus on the most significant items
- Consolidate related items if they represent aspects of the same core issue/request
- Example: "API is slow" and "API timeouts" can be one "API Performance" category
- But keep truly distinct items separate: "API performance" and "API documentation" are different
- Prioritize categories by impact and frequency of mention
- If more than 6 distinct areas exist, select the 6 most important/impactful ones

SENTIMENT SCORING:
- -1.0 to -0.7: Extremely negative (terrible, awful, hate, worst, critical failure)
- -0.6 to -0.3: Negative (frustrated, disappointed, annoying, problematic)
- -0.2 to 0.2: Neutral (factual, mixed feelings, matter-of-fact)
- 0.3 to 0.6: Positive (good, helpful, useful, appreciate)
- 0.7 to 1.0: Very positive (love, excellent, amazing, fantastic)

CRITICAL OUTPUT INSTRUCTIONS:
You MUST respond with ONLY a JSON array. The response should start with [ and end with ]
Do NOT wrap the array in any object. Do NOT add any text before or after the JSON.
Your ENTIRE response must be a valid JSON array like this:

[
    {{
        "key_area": "string - 2-4 words describing the category",
        "customer_problem": "string - comprehensive sentence covering all related issues/requests in this category",
        "sentiment_score": float between -1.0 and 1.0,
        "area_type": "issue" or "feature" (lowercase only)
    }},
    {{
        "key_area": "another category",
        "customer_problem": "another distinct problem or request",
        "sentiment_score": float between -1.0 and 1.0,
        "area_type": "issue" or "feature"
    }}
]

IMPORTANT: 
- Start your response with [ 
- End your response with ]
- Include AT LEAST 2 items in the array for any meaningful text
- Each distinct topic/issue/request gets its own object in the array
- Do NOT return a single object - always return an array even for one item

EXAMPLE 1 - Product Strategy Discussion (2 key areas):
Input: "Had a long discussion about our challenges. Picking the right features to build is really hard - we have multiple prioritization frameworks but still struggle. Also, our PMs don't spend enough time on customer research. The quarterly planning process remains painful and inefficient. On top of that, we keep building duplicate features across teams which is wasteful."

Output:
[
    {{
        "key_area": "Feature Prioritization",
        "customer_problem": "Difficulty in selecting the right features to build despite having multiple prioritization frameworks",
        "sentiment_score": -0.5,
        "area_type": "issue"
    }},
    {{
        "key_area": "Planning Process",
        "customer_problem": "Quarterly planning process is painful and inefficient",
        "sentiment_score": -0.6,
        "area_type": "issue"
    }},
    {{
        "key_area": "Feature Duplication",
        "customer_problem": "Teams are building duplicate features leading to wasted resources",
        "sentiment_score": -0.5,
        "area_type": "issue"
    }}
]

EXAMPLE 2 - Customer Service Chat (2 key areas):
Input: "Customer: I can't believe how bad this is. First, I can't find where to change my billing address anywhere in the settings. Second, I've been trying to download my invoices for tax purposes but the download button does nothing. Oh and one more thing - could you add a way to save multiple payment methods? I manage several company cards."

Output:
[
    {{
        "key_area": "Settings Navigation",
        "customer_problem": "Users cannot find where to change billing address in settings",
        "sentiment_score": -0.7,
        "area_type": "issue"
    }},
    {{
        "key_area": "Payment Methods",
        "customer_problem": "Users need ability to save and manage multiple payment methods for different company cards",
        "sentiment_score": 0.2,
        "area_type": "feature"
    }}
]

EXAMPLE 3 - Team Standup Notes (3 key areas):
Input: "Morning standup highlights: Lisa mentioned the new hire onboarding is taking 2 weeks which is way too long. Tom said customers in Europe can't access the platform after 6PM due to some timezone bug. Sarah wants us to add automated testing to speed up releases. The sales team is asking for a Salesforce integration by Q2. Also, the search feature returns irrelevant results according to user feedback."

Output:
[
    {{
        "key_area": "Onboarding Duration",
        "customer_problem": "New employee onboarding process takes 2 weeks which is excessive",
        "sentiment_score": -0.5,
        "area_type": "issue"
    }},
    {{
        "key_area": "Testing Automation",
        "customer_problem": "Development team wants automated testing to accelerate release cycles",
        "sentiment_score": 0.3,
        "area_type": "feature"
    }},
    {{
        "key_area": "CRM Integration",
        "customer_problem": "Sales team requires Salesforce integration by Q2",
        "sentiment_score": 0.2,
        "area_type": "feature"
    }},
]

EXAMPLE 4 - Executive Email (4 key areas):
Input: "Team, following up on our board meeting. We need to address these items: 1) Our mobile app ratings dropped to 2.5 stars mainly due to crashes 2) Implement SOC2 compliance for enterprise deals 3) The competitor analysis shows we're behind on AI features 4) Customer churn increased 15% last quarter 5) We need better analytics dashboards for the exec team 6) Consider acquiring that startup for their recommendation engine tech"

Output:
[

    {{
        "key_area": "Security Compliance",
        "customer_problem": "SOC2 compliance implementation needed for enterprise deals",
        "sentiment_score": 0.1,
        "area_type": "feature"
    }},
    {{
        "key_area": "AI Capabilities",
        "customer_problem": "Product lacking AI features compared to competitors",
        "sentiment_score": -0.4,
        "area_type": "issue"
    }},
    {{
        "key_area": "Executive Analytics",
        "customer_problem": "Executive team needs enhanced analytics dashboards for better decision making",
        "sentiment_score": 0.2,
        "area_type": "feature"
    }},
    {{
        "key_area": "Strategic Acquisition",
        "customer_problem": "Company considering acquiring startup for their recommendation engine technology",
        "sentiment_score": 0.4,
        "area_type": "feature"
    }}
]

EXAMPLE 5 - Slack Thread (2 key areas):
Input: "@channel just got off a call with BigCorp. They're threatening to cancel because our platform doesn't support their SSO provider. Also mentioned they'd love real-time collaboration features like Google Docs has."

Output:
[
    {{
        "key_area": "SSO Compatibility",
        "customer_problem": "Major client threatening cancellation due to lack of support for their SSO provider",
        "sentiment_score": -0.8,
        "area_type": "issue"
    }},
    {{
        "key_area": "Collaboration Features",
        "customer_problem": "Enterprise clients want real-time collaboration features similar to Google Docs",
        "sentiment_score": 0.3,
        "area_type": "feature"
    }}
]


SPECIAL CASES:
1. If the text contains NO clear issues, feature requests, or actionable feedback:
   - Return an empty array: []
   - This is valid for texts like "Everything is fine", "No comments", "N/A", etc.

2. If the text is too vague or general to classify:
   - Return a single item with key_area: "General Feedback"
   - Example: "I have some thoughts" → general feedback category

3. Category count guidelines:
   - 0 items: Only for text with no actionable content
   - 1 item: For vague or single-topic text
   - 2-6 items: For text mentioning multiple distinct topics (maximum 6)
   - If more than 6 potential categories exist, prioritize the most significant ones

Remember:
- ALWAYS return a JSON array starting with [ and ending with ]
- NEVER return a single object {...} - always an array [...]
- Extract the most significant distinct items (maximum 6 categories)
- 2-6 items for meaningful text, focusing on quality over quantity
- Consolidate closely related items when they represent the same core concern
- Keep truly distinct issues/features separate
- Be specific in the customer_problem - don't be vague or overly broad
- Only output the JSON array, nothing else
- Prioritize impact and importance when selecting categories
'''

        # Try GPT classification with retries
        max_retries = 3
        for attempt in range(max_retries):
            try:
                logger.info(f"[classify-single-review] Attempt {attempt + 1}/{max_retries} - Calling GPT-4.1-mini")
                logger.info(f"[classify-single-review] Prompt length: {len(prompt)} characters")
                
                response = AZURE_CLIENT.chat.completions.create(
                    model="gpt-4.1-mini",
                    messages=[
                        {
                            "role": "system",
                            "content": "You are a text analysis expert. Extract ALL distinct issues and features mentioned in any text. Output ONLY a valid JSON array starting with [ and ending with ]. Never return a single object."
                        },
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.2,
                    max_tokens=3000,
                )
                
                logger.info(f"[classify-single-review] GPT call successful")
                
                result_text = response.choices[0].message.content
                logger.info(f"[classify-single-review] GPT response length: {len(result_text)} characters")
                logger.debug(f"[classify-single-review] Raw GPT response: {result_text[:500]}...")
                
                # Clean the response - remove markdown code blocks and extra whitespace
                cleaned_text = result_text.strip()
                if cleaned_text.startswith("```json"):
                    cleaned_text = cleaned_text[7:]  # Remove ```json
                elif cleaned_text.startswith("```"):
                    cleaned_text = cleaned_text[3:]   # Remove ```
                if cleaned_text.endswith("```"):
                    cleaned_text = cleaned_text[:-3]  # Remove trailing ```
                cleaned_text = cleaned_text.strip()
                
                logger.debug(f"[classify-single-review] Cleaned response: {cleaned_text[:200]}...")
                
                # Try to parse the response
                try:
                    # First try direct parsing
                    result = json.loads(cleaned_text)
                    logger.info(f"[classify-single-review] Successfully parsed JSON response")
                    
                    # Handle different response formats
                    if isinstance(result, list):
                        logger.info(f"[classify-single-review] Response is a list with {len(result)} items")
                        # Good - we got a list
                        pass
                    elif isinstance(result, dict):
                        logger.info(f"[classify-single-review] Response is a dict with keys: {list(result.keys())}")
                        
                        # Check for error response
                        if "error" in result:
                            logger.error(f"[classify-single-review] GPT returned error: {result.get('error')}")
                            if attempt < max_retries - 1:
                                logger.info(f"[classify-single-review] Retrying after error response...")
                                await asyncio.sleep(1)
                                continue
                            else:
                                raise ValueError(f"GPT error: {result.get('error')}")
                        
                        # Look for an array in the response
                        array_found = False
                        for key in ["classifications", "results", "items", "analysis_results", "categories"]:
                            if key in result and isinstance(result[key], list):
                                result = result[key]
                                logger.info(f"[classify-single-review] Using '{key}' field containing {len(result)} items")
                                array_found = True
                                break
                        
                        if not array_found:
                            # Try to find any key that contains a list
                            for key, value in result.items():
                                if isinstance(value, list) and len(value) > 0:
                                    # Check if it looks like our expected format
                                    if isinstance(value[0], dict) and any(k in value[0] for k in ["key_area", "customer_problem"]):
                                        result = value
                                        logger.info(f"[classify-single-review] Found list in '{key}' field with {len(result)} items")
                                        array_found = True
                                        break
                            
                            if not array_found:
                                # Single item response - this is what we want to avoid
                                logger.warning(f"[classify-single-review] Response is a single object, not an array")
                                if all(key in result for key in ["key_area", "customer_problem", "area_type"]):
                                    # It's a valid single classification - wrap it
                                    result = [result]
                                    logger.info(f"[classify-single-review] Wrapped single classification in array")
                                else:
                                    # Invalid structure
                                    if attempt < max_retries - 1:
                                        logger.info(f"[classify-single-review] Invalid structure, retrying...")
                                        await asyncio.sleep(1)
                                        continue
                                    else:
                                        raise ValueError("Response is not in expected format")
                    else:
                        logger.error(f"[classify-single-review] Response is neither list nor dict: {type(result)}")
                        if attempt < max_retries - 1:
                            await asyncio.sleep(1)
                            continue
                        else:
                            raise ValueError(f"Response is {type(result)}, expected list or dict")
                    
                    # Validate it's a list at this point
                    if not isinstance(result, list):
                        raise ValueError("Failed to extract list from response")
                    
                    # Check if we have enough items
                    if len(result) == 0 and len(review_text) > 500:
                        # Only retry if text is substantial but no classifications found
                        logger.warning(f"[classify-single-review] No classifications for {len(review_text)} char text")
                        if attempt < max_retries - 1:
                            logger.info(f"[classify-single-review] Retrying to get classifications...")
                            await asyncio.sleep(1)
                            continue
                    elif len(result) < 2 and len(review_text) > 250:
                        # For medium-length text, we expect at least 2 classifications
                        logger.info(f"[classify-single-review] Only {len(result)} classifications for {len(review_text)} char text - accepting result")
                        # Don't retry, accept the result as is
                    
                    logger.info(f"[classify-single-review] Extracted {len(result)} items from response")
                    
                    # Validate and fix each item
                    validated_results = []
                    for i, item in enumerate(result):
                        if not isinstance(item, dict):
                            logger.warning(f"[classify-single-review] Item {i} is not a dict, skipping")
                            continue
                        
                        # Log missing fields
                        missing_fields = []
                        for field in ["key_area", "customer_problem", "sentiment_score", "area_type"]:
                            if field not in item:
                                missing_fields.append(field)
                        
                        if missing_fields:
                            logger.warning(f"[classify-single-review] Item {i} missing fields: {missing_fields}")
                        
                        # Ensure all required fields
                        validated_item = {
                            "key_area": item.get("key_area", "Unknown Area"),
                            "customer_problem": item.get("customer_problem", "Unspecified problem"),
                            "sentiment_score": 0.0,
                            "area_type": "issue"
                        }
                        
                        # Validate sentiment score
                        try:
                            score = float(item.get("sentiment_score", 0.0))
                            validated_item["sentiment_score"] = max(-1.0, min(1.0, score))
                        except:
                            logger.warning(f"[classify-single-review] Item {i} has invalid sentiment score")
                            validated_item["sentiment_score"] = 0.0
                        
                        # Validate area_type
                        area_type = item.get("area_type", "").lower()
                        if area_type in ["feature", "issue"]:
                            validated_item["area_type"] = area_type
                        else:
                            # Infer from problem description
                            inferred_type = infer_area_type(review_text, validated_item["customer_problem"])
                            logger.info(f"[classify-single-review] Item {i} area_type inferred as '{inferred_type}'")
                            validated_item["area_type"] = inferred_type
                        
                        validated_results.append(validated_item)
                    
                    # If no valid results, try again
                    if not validated_results:
                        # If GPT explicitly returned an empty array, that's a valid response
                        if isinstance(result, list) and len(result) == 0:
                            logger.info(f"[classify-single-review] GPT returned empty array - no classifications found in text")
                            # Return a default/neutral classification
                            return JSONResponse(content=[{
                                "key_area": "General Feedback",
                                "customer_problem": "No specific issues or feature requests identified in the provided text",
                                "sentiment_score": 0.0,
                                "area_type": "issue"
                            }], status_code=200)
                        else:
                            # This is an actual error - missing data in response
                            logger.error(f"[classify-single-review] No valid classifications found after validation")
                            if attempt < max_retries - 1:
                                logger.info(f"[classify-single-review] Retrying...")
                                await asyncio.sleep(1)
                                continue
                            else:
                                # Return a default response instead of raising error
                                logger.warning(f"[classify-single-review] All attempts failed to extract valid classifications, returning default")
                                return JSONResponse(content=[{
                                    "key_area": "Processing Error",
                                    "customer_problem": "Unable to classify the provided text after multiple attempts",
                                    "sentiment_score": -0.2,
                                    "area_type": "issue"
                                }], status_code=200)
                    
                    logger.info(f"[classify-single-review] Successfully validated {len(validated_results)} classifications")
                    
                    # Log the final results summary
                    for i, item in enumerate(validated_results):
                        logger.info(f"[classify-single-review] Item {i+1}: {item['key_area']} ({item['area_type']}) - sentiment: {item['sentiment_score']:.2f}")
                    
                    return JSONResponse(content=validated_results, status_code=200)
                    
                except (json.JSONDecodeError, ValueError) as e:
                    logger.error(f"[classify-single-review] Attempt {attempt + 1} - Failed to parse/validate GPT response: {str(e)}")
                    logger.error(f"[classify-single-review] Raw response: {result_text[:500]}...")
                    
                    if attempt < max_retries - 1:
                        logger.info(f"[classify-single-review] Retrying after parse error...")
                        await asyncio.sleep(1 ** (attempt + 1))  # Exponential backoff
                        continue
                    else:
                        raise
                        
            except Exception as e:
                logger.error(f"[classify-single-review] Attempt {attempt + 1} - Error during GPT call: {type(e).__name__}: {str(e)}")
                
                if attempt < max_retries - 1:
                    logger.info(f"[classify-single-review] Retrying after API error...")
                    await asyncio.sleep(2 ** (attempt + 1))  # Exponential backoff
                    continue
                else:
                    logger.error(f"[classify-single-review] All {max_retries} attempts failed")
                    raise
        
        # If we get here, all retries failed
        raise Exception("All classification attempts failed")
            
    except HTTPException:
        logger.error(f"[classify-single-review] HTTPException raised")
        raise
    except Exception as e:
        logger.error(f"[classify-single-review] Unexpected error: {type(e).__name__}: {str(e)}")
        logger.error(f"[classify-single-review] Full traceback: {traceback.format_exc()}")
        
        # Return error classification
        return JSONResponse(
            content=[{
                "key_area": "System Error",
                "customer_problem": f"Classification failed: {str(e)[:100]}",
                "sentiment_score": -0.8,
                "area_type": "issue"
            }],
            status_code=200
        )
def infer_area_type(review_text: str, customer_problem: str) -> str:
    '''
    Infers whether a review is a feature request or an issue based on text analysis.
    
    Args:
        review_text: The review text
        customer_problem: The identified customer problem
        
    Returns:
        "feature" or "issue"
    '''
    feature_keywords = [
        "want", "need", "request", "add", "feature", "enhance", "improve",
        "would like", "wish", "suggest", "should have", "could have",
        "implement", "missing", "lacks", "needs", "requires", "please add",
        "it would be great", "consider adding", "functionality"
    ]
    
    issue_keywords = [
        "broken", "bug", "crash", "error", "fail", "doesn't work", "not working",
        "issue", "problem", "glitch", "freeze", "slow", "hang", "stuck",
        "unable", "cannot", "malfunction", "defect", "fault", "wrong"
    ]
    
    combined_text = f"{review_text} {customer_problem}".lower()
    
    # Count keyword matches
    feature_count = sum(1 for keyword in feature_keywords if keyword in combined_text)
    issue_count = sum(1 for keyword in issue_keywords if keyword in combined_text)
    
    # Default to issue if unclear
    if feature_count > issue_count:
        return "feature"
    else:
        return "issue"


def extract_classification_from_text(text: str, existing_categories: List[Dict[str, str]], original_json: str = "[]") -> Dict[str, str]:
    '''
    Fallback method to extract key_area, customer_problem, sentiment_score, and area_type from GPT response when JSON parsing fails.
    
    Args:
        text: Raw text response from GPT
        existing_categories: List of existing categories
        original_json: Original JSON input string
        
    Returns:
        Dictionary with key_area, customer_problem, sentiment_score, and area_type
    '''
    # If no text provided, try to extract from original JSON or use default
    if not text or not text.strip():
        logger.warning("Empty response from GPT, attempting to use original data")
        try:
            # Try to extract from original JSON first
            if original_json and original_json.strip() not in ('[]', '{}', ''):
                try:
                    json_data = json.loads(original_json)
                    if isinstance(json_data, list) and len(json_data) > 0:
                        if "key_area" in json_data[0] and "customer_problem" in json_data[0]:
                            logger.info("Using first category from original JSON as fallback")
                            return {
                                "key_area": json_data[0]["key_area"],
                                "customer_problem": json_data[0]["customer_problem"],
                                "sentiment_score": 0.0,  # Default neutral sentiment
                                "area_type": "issue"  # Default to issue
                            }
                except:
                    pass
            
            # If original JSON parsing fails or is empty, use the first existing category
            if existing_categories:
                return {
                    "key_area": existing_categories[0]["key_area"],
                    "customer_problem": existing_categories[0]["customer_problem"],
                    "sentiment_score": 0.0,  # Default neutral sentiment
                    "area_type": "issue"  # Default to issue
                }
        except:
            pass
        
        # Final fallback
        return {
            "key_area": "Miscellaneous Issue",
            "customer_problem": "Issues related to miscellaneous customer concerns",
            "sentiment_score": -0.1,  # Slightly negative default for issues
            "area_type": "issue"
        }
    
    try:
        # Extract sentiment score from text if present
        sentiment_pattern = r'"sentiment_score"\s*:\s*(-?\d+\.?\d*)'
        sentiment_match = re.search(sentiment_pattern, text)
        sentiment_score = 0.0  # Default neutral
        
        if sentiment_match:
            try:
                sentiment_score = float(sentiment_match.group(1))
                sentiment_score = max(-1.0, min(1.0, sentiment_score))  # Clamp to [-1.0, 1.0]
            except (ValueError, TypeError):
                # Invalid float, use default
                sentiment_score = 0.0
        
        # Extract area_type from text if present
        area_type_pattern = r'"area_type"\s*:\s*"(feature|issue)"'
        area_type_match = re.search(area_type_pattern, text)
        area_type = area_type_match.group(1) if area_type_match else None
        
        # Multi-tiered parsing approach
        
        # 1. Try with the raw text as is
        try:
            raw_json = json.loads(text)
            if "key_area" in raw_json and "customer_problem" in raw_json:
                result = {
                    "key_area": raw_json["key_area"],
                    "customer_problem": raw_json["customer_problem"],
                    "sentiment_score": sentiment_score,
                    "area_type": raw_json.get("area_type", area_type or "issue")
                }
                
                # Update sentiment if present in JSON
                if "sentiment_score" in raw_json:
                    try:
                        result["sentiment_score"] = float(raw_json["sentiment_score"])
                        result["sentiment_score"] = max(-1.0, min(1.0, result["sentiment_score"]))
                    except (ValueError, TypeError):
                        pass
                
                # Infer area_type if not valid
                if result["area_type"] not in ["feature", "issue"]:
                    result["area_type"] = infer_area_type("", result["customer_problem"])
                
                return result
        except:
            pass
        
        # 2. Clean and try again
        cleaned_text = text.strip()
        # Remove markdown code block markers and extra spaces
        cleaned_text = re.sub(r'```json|```|`', '', cleaned_text).strip()
        try:
            cleaned_json = json.loads(cleaned_text)
            if "key_area" in cleaned_json and "customer_problem" in cleaned_json:
                result = {
                    "key_area": cleaned_json["key_area"],
                    "customer_problem": cleaned_json["customer_problem"],
                    "sentiment_score": sentiment_score,
                    "area_type": cleaned_json.get("area_type", area_type or "issue")
                }
                
                # Update sentiment if present
                if "sentiment_score" in cleaned_json:
                    try:
                        result["sentiment_score"] = float(cleaned_json["sentiment_score"])
                        result["sentiment_score"] = max(-1.0, min(1.0, result["sentiment_score"]))
                    except (ValueError, TypeError):
                        pass
                
                # Infer area_type if not valid
                if result["area_type"] not in ["feature", "issue"]:
                    result["area_type"] = infer_area_type("", result["customer_problem"])
                
                return result
        except:
            pass
        
        # 3. Try to extract any JSON-like content
        json_pattern = r'\{[^{}]*\}'
        json_matches = re.findall(json_pattern, cleaned_text)
        
        if json_matches:
            for json_str in json_matches:
                try:
                    # Try to parse this JSON fragment
                    result = json.loads(json_str)
                    if "key_area" in result and "customer_problem" in result:
                        extracted = {
                            "key_area": result["key_area"],
                            "customer_problem": result["customer_problem"],
                            "sentiment_score": sentiment_score,
                            "area_type": result.get("area_type", area_type or "issue")
                        }
                        
                        # Update sentiment if present
                        if "sentiment_score" in result:
                            try:
                                extracted["sentiment_score"] = float(result["sentiment_score"])
                                extracted["sentiment_score"] = max(-1.0, min(1.0, extracted["sentiment_score"]))
                            except (ValueError, TypeError):
                                pass
                        
                        # Infer area_type if not valid
                        if extracted["area_type"] not in ["feature", "issue"]:
                            extracted["area_type"] = infer_area_type("", extracted["customer_problem"])
                        
                        # Check if this is an existing pair
                        for cat in existing_categories:
                            if (cat["key_area"] == result["key_area"] and 
                                cat["customer_problem"] == result["customer_problem"]):
                                return extracted
                        
                        # If not an exact match, still return the extracted pair
                        return extracted
                except:
                    continue
        
        # 4. Try to extract key_area and customer_problem from text using regex
        key_area_pattern = r'"key_area"\s*:\s*"([^"]+)"'
        problem_pattern = r'"customer_problem"\s*:\s*"([^"]+)"'
        
        key_area_match = re.search(key_area_pattern, cleaned_text)
        problem_match = re.search(problem_pattern, cleaned_text)
        
        if key_area_match and problem_match:
            key_area = key_area_match.group(1)
            customer_problem = problem_match.group(1)
            
            result = {
                "key_area": key_area,
                "customer_problem": customer_problem,
                "sentiment_score": sentiment_score,
                "area_type": area_type or infer_area_type("", customer_problem)
            }
            
            # Check if this is an existing pair
            for cat in existing_categories:
                if (cat["key_area"] == key_area and cat["customer_problem"] == customer_problem):
                    return result
            
            # If not an exact match, return the extracted values
            return result
        
        # 5. Look for exact matches with existing categories in the text
        for category in existing_categories:
            # If both key_area and customer_problem appear in the text, prioritize that match
            if category["key_area"] in cleaned_text and category["customer_problem"] in cleaned_text:
                logger.info(f"Found exact category match in text: {category['key_area']}")
                return {
                    "key_area": category["key_area"],
                    "customer_problem": category["customer_problem"],
                    "sentiment_score": sentiment_score,
                    "area_type": area_type or infer_area_type("", category["customer_problem"])
                }
        
        # 6. Look for partial matches with existing categories
        for category in existing_categories:
            # Look for substantial parts of the customer_problem in the text
            problem_parts = category["customer_problem"].split()
            if len(problem_parts) >= 3:
                # Check if at least 3 consecutive words from the problem appear in the text
                for i in range(len(problem_parts) - 2):
                    phrase = " ".join(problem_parts[i:i+3])
                    if phrase.lower() in cleaned_text.lower():
                        logger.info(f"Found partial match on customer problem: {category['key_area']}")
                        return {
                            "key_area": category["key_area"],
                            "customer_problem": category["customer_problem"],
                            "sentiment_score": sentiment_score,
                            "area_type": area_type or infer_area_type("", category["customer_problem"])
                        }
        
        # 7. Check for key_area matches
        for category in existing_categories:
            if category["key_area"].lower() in cleaned_text.lower():
                logger.info(f"Found key_area match: {category['key_area']}")
                return {
                    "key_area": category["key_area"],
                    "customer_problem": category["customer_problem"],
                    "sentiment_score": sentiment_score,
                    "area_type": area_type or infer_area_type("", category["customer_problem"])
                }
        
        # 8. If all else fails, try to return raw JSON as last resort
        try:
            # If text looks like JSON but doesn't have the expected fields
            if text.strip().startswith('{') and text.strip().endswith('}'):
                raw_json_attempt = json.loads(text.strip())
                # If it's valid JSON with any fields, add the required fields with defaults
                if isinstance(raw_json_attempt, dict):
                    result = {
                        "key_area": raw_json_attempt.get("key_area", "Miscellaneous Issue"),
                        "customer_problem": raw_json_attempt.get("customer_problem", "General customer feedback"),
                        "sentiment_score": sentiment_score,
                        "area_type": raw_json_attempt.get("area_type", "issue")
                    }
                    
                    # Infer area_type if not valid
                    if result["area_type"] not in ["feature", "issue"]:
                        result["area_type"] = infer_area_type("", result["customer_problem"])
                    
                    return result
        except:
            pass
        
        # 9. Final fallback - return first category
        if existing_categories:
            logger.info("Using first existing category as final fallback")
            return {
                "key_area": existing_categories[0]["key_area"],
                "customer_problem": existing_categories[0]["customer_problem"],
                "sentiment_score": sentiment_score,
                "area_type": area_type or infer_area_type("", existing_categories[0]["customer_problem"])
            }
        
        # 10. Absolute last resort
        return {
            "key_area": "Miscellaneous Issue",
            "customer_problem": "Issues related to miscellaneous customer concerns",
            "sentiment_score": sentiment_score,
            "area_type": "issue"
        }
    
    except Exception as e:
        logger.error(f"Error in fallback extraction: {str(e)}")
        # Return the original JSON if all else fails
        try:
            if original_json and original_json not in ('[]', '{}', ''):
                json_data = json.loads(original_json)
                if isinstance(json_data, list) and len(json_data) > 0:
                    return {
                        "key_area": json_data[0].get("key_area", "Miscellaneous Issue"),
                        "customer_problem": json_data[0].get("customer_problem", "General customer feedback"),
                        "sentiment_score": -0.1,  # Slightly negative default for issues
                        "area_type": "issue"
                    }
        except:
            pass
            
        # Absolute final fallback
        return {
            "key_area": "Miscellaneous Issue",
            "customer_problem": "Issues related to miscellaneous customer concerns",
            "sentiment_score": -0.1,  # Slightly negative default for issues
            "area_type": "issue"
        }

        
@app.get("/", response_class=HTMLResponse)
async def serve_webpage():
    '''Serve the AI Assistant Hub webpage'''
    try:
        with open("webpage.html", "r", encoding="utf-8") as f:
            html_content = f.read()
        return HTMLResponse(content=html_content)
    except FileNotFoundError:
        return HTMLResponse(
            content="<h1>Error: webpage.html not found</h1><p>Please ensure webpage.html is in the same directory as your app.py file.</p>",
            status_code=404
        )
    except Exception as e:
        logger.error(f"Error serving webpage: {e}")
        return HTMLResponse(
            content="<h1>Error loading webpage</h1>",
            status_code=500
        )

# Optional: Add a favicon endpoint
@app.get("/favicon.ico")
async def favicon():
    '''Return a simple favicon'''
    return Response(content="", media_type="image/x-icon")
# Main function for running as standalone server
if __name__ == "__main__":
    import uvicorn
    
    # Set up server information
    port = int(os.environ.get("PORT", 8081))
    host = os.environ.get("HOST", "0.0.0.0")
    
    print(f"┌───────────────────────────────────────────────────────┐")
    print(f"│                                                       │")
    print(f"│  Customer Feedback Analysis API Server                │")
    print(f"│                                                       │")
    print(f"│  Endpoint:     http://{host}:{port}/analyze-feedback  │")
    print(f"│  Documentation: http://{host}:{port}/docs             │")
    print(f"│                                                       │")
    print(f"│  Logs:         customer_feedback.log                  │")
    print(f"│                                                       │")
    print(f"└───────────────────────────────────────────────────────┘")
    
    # Configure server settings
    uvicorn.run(
        app, 
        host=host, 
        port=port,
        log_level="info"
    )
